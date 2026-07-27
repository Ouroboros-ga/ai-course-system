"""
文档处理API接口
流程：上传文件 -> Docling解析为Markdown -> 豆包AI生成智课脚本 -> 存储到数据库 -> 返回结果
Updated: 2026-03-28 - 集成Docling解析和豆包AI，添加用户认证
Updated: 2026-04-09 - 重构：解析逻辑迁移到document_service服务层
"""

import os
import uuid
import tempfile
import json
import logging
import time
from pathlib import Path
from typing import Optional
from app.core.time_utils import utcnow_aware

from fastapi import APIRouter, UploadFile, File, HTTPException, Depends, Request, Query
from fastapi.responses import JSONResponse
from sqlmodel import Session, select, text, func

from app.schemas.document_schema import (
    DocumentUploadResponse,
    DocumentAnalyzeRequest,
    DocumentAnalyzeResponse,
)
from app.schemas.common_schema import UnifiedResponse
from app.core.exceptions import unified_response
from app.core.security import get_current_user, _get_user_id, _get_user_identity
from app.core.config import settings
from app.models.database import get_session
from app.models.course_model import (
    Course,
    CourseScript,
    ScriptNode,
    DoclingDocument,
    DoclingGroup,
    DoclingTable,
    DoclingTableCell,
    DoclingText,
    DoclingPicture,
    StudentEnrollment,
    CourseStatus,
    ParseStatus,
    ScriptNodeType,
)
from app.models.document_artifact_model import DocumentArtifact
from app.models.access_control_model import CourseMembership, CourseRole, MembershipStatus
from app.models.user_model import ChatHistory
from app.services.document_service import document_service
from app.services.course_access_service import (
    CourseAccessContext,
    activate_student_membership,
    course_permission,
    establish_course_access_baseline,
    require_course_permission,
)
from app.services import smart_course_service
from app.platform.adapters.base import classify_exception
from app.platform.adapters.errors import AdapterErrorCode
from app.platform.adapters.tts import TTSAdapter
from app.platform.tasks import TaskContext, TaskResult, TaskRunner, TaskStatus, TaskType

router = APIRouter(tags=["文档处理"])
logger = logging.getLogger(__name__)

UPLOAD_DIR = Path(tempfile.gettempdir()) / "ai_course_uploads"
UPLOAD_DIR.mkdir(exist_ok=True)

document_cache = {}
tts_generation_status = {}


def _tts_batch_task_status(
    total_count: int,
    success_count: int,
    failed_count: int,
    task_level_status: TaskStatus | None = None,
) -> TaskStatus:
    if task_level_status in {TaskStatus.TIMEOUT, TaskStatus.FAILED}:
        return task_level_status
    if total_count == 0:
        return TaskStatus.SUCCEEDED
    if success_count == total_count and failed_count == 0:
        return TaskStatus.SUCCEEDED
    if success_count > 0 and failed_count > 0:
        return TaskStatus.PARTIAL_SUCCESS
    if failed_count == total_count:
        return TaskStatus.FAILED
    return TaskStatus.RUNNING


def _tts_public_status(task_status: TaskStatus) -> str:
    if task_status == TaskStatus.PARTIAL_SUCCESS:
        return "partial"
    if task_status in {TaskStatus.FAILED, TaskStatus.TIMEOUT}:
        return "failed"
    if task_status == TaskStatus.SUCCEEDED:
        return "completed"
    return "processing"


def _tts_node_failure_result(
    context: TaskContext,
    exc: Exception,
    prior_result: TaskResult | None,
    duration_ms: float,
) -> TaskResult:
    if prior_result is not None:
        return TaskResult.fail(
            prior_result.error_code or AdapterErrorCode.UNKNOWN_ERROR,
            prior_result.error_message or str(exc),
            status=prior_result.status,
            provider=prior_result.provider,
            raw=prior_result.raw,
            duration_ms=duration_ms,
            context=context,
        )

    error_code = classify_exception(exc)
    return TaskResult.fail(
        error_code,
        str(exc),
        status=TaskStatus.TIMEOUT if error_code == AdapterErrorCode.TIMEOUT else TaskStatus.FAILED,
        duration_ms=duration_ms,
        context=context,
    )


def _tts_batch_task_result(
    course_id: int,
    total_count: int,
    node_results: list[TaskResult],
    errors: list[dict],
    duration_ms: float,
    task_level_result: TaskResult | None = None,
) -> TaskResult:
    success_count = sum(1 for result in node_results if result.success)
    failed_results = [result for result in node_results if not result.success]
    failed_count = len(failed_results)
    task_status = _tts_batch_task_status(
        total_count,
        success_count,
        failed_count,
        task_level_result.status if task_level_result else None,
    )
    failed_nodes = [
        {
            "node_id": result.context.node_id if result.context else None,
            "status": result.status.value,
            "error_code": result.error_code,
            "error": result.error_message,
        }
        for result in failed_results
    ]
    context = TaskContext(
        task_id=str(course_id),
        task_type=TaskType.TTS_BATCH,
        course_id=course_id,
        provider="tts",
        metadata={
            "total_count": total_count,
            "success_count": success_count,
            "failed_count": failed_count,
            "completed_count": len(node_results),
        },
    )
    data = {
        "total_count": total_count,
        "success_count": success_count,
        "failed_count": failed_count,
        "completed_count": len(node_results),
        "failed_nodes": failed_nodes,
        "errors": errors,
        "duration_ms": duration_ms,
        "status": task_status.value,
    }
    if task_status in {TaskStatus.SUCCEEDED, TaskStatus.PARTIAL_SUCCESS}:
        return TaskResult.ok(
            data=data,
            status=task_status,
            duration_ms=duration_ms,
            context=context,
        )

    error_codes = {result.error_code for result in failed_results if result.error_code}
    error_code = (
        task_level_result.error_code
        if task_level_result
        else next(iter(error_codes)) if len(error_codes) == 1 else "tts_batch_failed"
    )
    error_message = (
        task_level_result.error_message
        if task_level_result
        else f"TTS batch failed for {failed_count} node(s)"
    )
    result = TaskResult.fail(
        error_code,
        error_message or "TTS batch failed",
        status=task_status,
        raw=data,
        duration_ms=duration_ms,
        context=context,
    )
    result.data = data
    return result


def _update_tts_generation_status(status_info: dict, task_result: TaskResult) -> None:
    data = task_result.data or {}
    status_info.update({
        "status": _tts_public_status(task_result.status),
        "success_count": data.get("success_count", 0),
        "failed_count": data.get("failed_count", 0),
        "completed_count": data.get("completed_count", 0),
        "failed_nodes": data.get("failed_nodes", []),
        "error_code": task_result.error_code,
        "duration_ms": data.get("duration_ms", task_result.duration_ms),
    })


def _raise_for_tts_task_failure(result: TaskResult) -> None:
    if not result.success:
        raise RuntimeError(result.error_message or "TTS synthesis failed")

async def _background_synthesize_audio(course_id: int, script_id: int):
    import asyncio
    from app.models.database import engine
    from sqlmodel import create_engine as _ce, Session as _Session

    status_key = str(course_id)
    batch_started = time.perf_counter()
    node_results: list[TaskResult] = []
    total = 0
    completed = 0
    tts_generation_status[status_key] = {
        "status": "processing",
        "total": 0,
        "completed": 0,
        "errors": [],
    }

    try:
        with _Session(engine) as session:
            nodes = session.exec(
                select(ScriptNode)
                .where(ScriptNode.script_id == script_id)
                .order_by(ScriptNode.node_index)
            ).all()

            course_audio_dir = get_course_audio_dir(course_id)
            total = len([n for n in nodes if n.content and len(n.content.strip()) >= 10])
            tts_generation_status[status_key]["total"] = total

            for node in nodes:
                if not node.content or len(node.content.strip()) < 10:
                    continue

                node_started = time.perf_counter()
                node_context = TaskContext(
                    task_id=f"{course_id}:{node.id}",
                    task_type=TaskType.TTS_BATCH,
                    course_id=course_id,
                    node_id=node.id,
                    provider="tts",
                    input_summary=node.content.strip()[:120],
                    metadata={"stage": "background_batch"},
                )
                failed_task_result = None
                try:
                    from app.common.tts_client import tts_client

                    cleanup_old_node_audio(node, course_audio_dir)

                    content = node.content.strip()
                    if len(content) > 2000:
                        segments = []
                        current = ""
                        for char in content:
                            current += char
                            if char in "。！？；" and len(current) >= 500:
                                segments.append(current)
                                current = ""
                        if current:
                            segments.append(current)
                    else:
                        segments = [content]

                    all_audio = b""
                    voice = None
                    if node.extra_data:
                        voice = node.extra_data.get("voice")

                    for seg in segments:
                        tts_result = await TaskRunner().run(
                            node_context,
                            lambda seg=seg: TTSAdapter(tts_client).synthesize(
                                text=seg, voice=voice, sample_rate=16000, output_format="mp3",
                            ),
                        )
                        if not tts_result.success:
                            failed_task_result = tts_result
                        _raise_for_tts_task_failure(tts_result)
                        response = tts_result.data
                        all_audio += response.audio_data

                    audio_filename = f"node_{node.id}_{uuid.uuid4().hex[:8]}.mp3"
                    audio_path = course_audio_dir / audio_filename
                    audio_path.write_bytes(all_audio)

                    audio_url = f"/api/v1/document/audio/{course_id}/{audio_filename}"
                    node.audio_url = audio_url
                    node.audio_duration = len(all_audio) / 16000 / 2

                    session.add(node)
                    completed += 1
                    tts_generation_status[status_key]["completed"] = completed
                    node_results.append(TaskResult.ok(
                        data={
                            "node_id": node.id,
                            "audio_url": node.audio_url,
                            "audio_duration": node.audio_duration,
                        },
                        provider="tts",
                        duration_ms=(time.perf_counter() - node_started) * 1000,
                        context=node_context,
                    ))

                except Exception as e:
                    logger.warning(f"Background TTS failed for node {node.id}: {e}")
                    node_results.append(_tts_node_failure_result(
                        node_context,
                        e,
                        failed_task_result,
                        (time.perf_counter() - node_started) * 1000,
                    ))
                    tts_generation_status[status_key]["errors"].append({
                        "node_id": node.id, "title": node.title, "error": str(e),
                    })

            session.commit()

        batch_task_result = _tts_batch_task_result(
            course_id,
            total,
            node_results,
            tts_generation_status[status_key]["errors"],
            (time.perf_counter() - batch_started) * 1000,
        )
        _update_tts_generation_status(tts_generation_status[status_key], batch_task_result)
        logger.info(f"Background TTS generation completed for course {course_id}: {completed}/{total}")

    except Exception as e:
        logger.error(f"Background TTS generation failed for course {course_id}: {e}")
        tts_generation_status[status_key]["errors"].append({"error": str(e)})
        error_code = classify_exception(e)
        task_level_result = TaskResult.fail(
            error_code,
            str(e),
            status=TaskStatus.TIMEOUT if error_code == AdapterErrorCode.TIMEOUT else TaskStatus.FAILED,
        )
        batch_task_result = _tts_batch_task_result(
            course_id,
            total,
            node_results,
            tts_generation_status[status_key]["errors"],
            (time.perf_counter() - batch_started) * 1000,
            task_level_result=task_level_result,
        )
        _update_tts_generation_status(tts_generation_status[status_key], batch_task_result)

@router.post("/upload", response_model=UnifiedResponse)
async def upload_document(
    request: Request,
    file: UploadFile = File(..., description="上传的文档文件 (PDF, DOCX, PPTX等)"),
    session: Session = Depends(get_session),
    current_user: dict = Depends(get_current_user),
):
    """
    上传文档并解析，存储到数据库，生成智课脚本
    
    需要用户登录认证
    
    流程：
    1. 验证用户身份
    2. 存储文件信息到 courses 表
    3. 创建 docling_documents 记录 (status = PENDING)
    4. 调用document_service处理文档（解析+RAG+脚本生成）
    5. 存储解析结果到数据库
    6. 创建聊天记录归档
    7. 返回结果
    """
    try:
        user_id = int(current_user["user_id"])
        username = current_user.get("username", "user")
        print(f"[认证] 用户 {username} (ID: {user_id}) 上传文件")

        document_id = str(uuid.uuid4())

        file_ext = Path(file.filename).suffix.lower()
        safe_filename = f"{document_id}{file_ext}"
        file_path = UPLOAD_DIR / safe_filename

        with open(file_path, "wb") as f:
            content = await file.read()
            f.write(content)

        file_size = file_path.stat().st_size
        if file_size > 50 * 1024 * 1024:
            file_path.unlink()
            return unified_response(
                code=400,
                message="文件大小超过限制（最大50MB）",
                data=None
            )

        print(f"[步骤1] 存储文件信息到 courses 表")
        print(f"  文件已保存: {file_path}")

        course = Course(
            fanya_course_id=f"local_{document_id[:8]}",
            fanya_course_name=file.filename,
            title=Path(file.filename).stem,
            description=f"从文件 {file.filename} 导入的课程",
            teacher_id=user_id,
            status=CourseStatus.DRAFT,
            is_ai_generated=False,
            source_file_name=file.filename,
            source_file_path=str(file_path),
            source_mimetype=file.content_type or "application/octet-stream",
            total_pages=0,
        )
        session.add(course)
        session.commit()
        session.refresh(course)
        establish_course_access_baseline(session, course.id, user_id)
        session.commit()
        print(f"  创建课程记录: ID={course.id}, 标题={course.title}")

        try:
            from app.common.slide_converter import is_office_file, is_pdf_file, get_or_create_pdf
            if is_office_file(str(file_path)):
                print(f"[步骤1.5] 检测到Office文件，自动转换为PDF...")
                pdf_path = get_or_create_pdf(str(file_path))
                if pdf_path:
                    course.pdf_file_path = pdf_path
                    session.add(course)
                    session.commit()
                    print(f"  PDF转换成功: {pdf_path}")
                else:
                    print(f"  PDF转换失败，将在展示时重试")
            elif is_pdf_file(str(file_path)):
                course.pdf_file_path = str(file_path)
                session.add(course)
                session.commit()
                print(f"  PDF文件，直接存储路径")
        except Exception as e:
            print(f"  PDF转换异常(不影响主流程): {e}")

        print(f"[步骤2] 创建 docling_documents 记录 (status = PENDING)")
        docling_doc = DoclingDocument(
            course_id=course.id,
            schema_name="DoclingDocument",
            version="1.10.0",
            doc_name=file.filename,
            origin_filename=file.filename,
            origin_mimetype=file.content_type,
            origin_binary_hash=document_id,
            source_file_path=str(file_path),
            status=ParseStatus.PENDING,
        )
        session.add(docling_doc)
        session.commit()
        session.refresh(docling_doc)
        print(f"  创建文档解析记录: ID={docling_doc.id}, status=PENDING")

        print(f"[步骤3] 调用document_service处理文档")
        docling_doc.status = ParseStatus.PROCESSING
        session.commit()
        print(f"  更新状态: PROCESSING")

        process_result = await document_service.process_document(
            file_path=file_path,
            filename=file.filename,
            enable_rag=True,
            enable_script=True,
            course_id=course.id,
        )
        
        parse_result = process_result.parse_result
        structure_result = process_result.structure_result
        script_result = process_result.script_result
        rag_result = process_result.rag_result
        mind_map = process_result.mind_map

        print(f"  文档解析完成: {parse_result.parse_method}, {len(parse_result.markdown_content)} 字符")
        print(f"  RAG预处理: 公式{rag_result.formula_count}个, 表格{rag_result.table_count}个, 知识点{len(rag_result.knowledge_points)}个")
        print(f"  脚本生成: {len(script_result.nodes)} 个节点")

        print(f"[步骤4] 存储解析结果到相关表")
        total_texts = 0
        total_tables = 0
        total_pictures = 0
        total_groups = 0

        for group_data in structure_result.groups:
            group = DoclingGroup(
                doc_id=docling_doc.id,
                self_ref=group_data.get("self_ref", f"#/groups/{total_groups}"),
                name=group_data.get("name", f"Section {total_groups + 1}"),
                label=group_data.get("label", "section"),
                content_layer=group_data.get("content_layer", "body"),
                sort_order=total_groups,
            )
            session.add(group)
            session.commit()
            session.refresh(group)
            total_groups += 1
        print(f"  存储 {total_groups} 个分组记录")

        for idx, text_data in enumerate(structure_result.texts):
            text_record = DoclingText(
                doc_id=docling_doc.id,
                group_id=None,
                self_ref=text_data.get("self_ref", f"#/texts/{idx}"),
                label=text_data.get("label", "text"),
                text=text_data.get("text", ""),
                page_no=text_data.get("page_no", 1),
                sort_order=idx,
            )
            session.add(text_record)
            total_texts += 1
        session.commit()
        print(f"  存储 {total_texts} 条文本记录")

        docling_doc.status = ParseStatus.COMPLETED
        docling_doc.total_groups = total_groups
        docling_doc.total_texts = total_texts
        docling_doc.total_tables = total_tables
        docling_doc.total_pictures = total_pictures
        docling_doc.raw_json = {
            "groups": structure_result.groups,
            "texts": structure_result.texts,
            "tables": structure_result.tables,
            "pictures": structure_result.pictures,
            "raw_content": structure_result.raw_content,
        }
        docling_doc.updated_at = utcnow_aware()
        session.commit()
        print(f"  更新状态: COMPLETED")

        print(f"[步骤5] 存储智课脚本到数据库")
        course_script = CourseScript(
            course_id=course.id,
            version=1,
            version_name="v1.0",
            script_content=script_result.script_content,
            summary_text=script_result.summary,
            keywords=json.dumps(script_result.keywords, ensure_ascii=False),
            is_active=True,
            audio_url=None,
            audio_duration=0,
            created_by=user_id,
        )
        session.add(course_script)
        session.commit()
        session.refresh(course_script)
        print(f"  创建课程脚本记录: ID={course_script.id}")

        print(f"[步骤6] 拆分脚本节点")
        for idx, node in enumerate(script_result.nodes):
            script_node = ScriptNode(
                script_id=course_script.id,
                chapter_id=node.chapter_id,
                node_index=idx,
                node_type=ScriptNodeType(node.node_type),
                title=node.title,
                content=node.content,
                page_start=node.page_start,
                page_end=node.page_end,
                duration=node.duration,
                is_key_point=node.is_key_point,
                timestamp_start=node.timestamp_start,
                timestamp_end=node.timestamp_end,
            )
            session.add(script_node)
        session.commit()
        print(f"  创建 {len(script_result.nodes)} 个 script_nodes 记录（含时间戳数据）")

        course.total_nodes = len(script_result.nodes)
        course.total_duration = script_result.total_duration
        course.is_ai_generated = True
        course.status = CourseStatus.PUBLISHED
        course.updated_at = utcnow_aware()
        session.commit()
        print(f"  更新课程统计: total_nodes={course.total_nodes}, total_duration={course.total_duration}, status=PUBLISHED")

        print(f"[步骤6.5] 启动后台TTS音频生成任务")
        import asyncio
        asyncio.create_task(_background_synthesize_audio(course.id, course_script.id))
        print(f"  后台TTS任务已启动: course_id={course.id}, script_id={course_script.id}")

        print(f"[步骤7] 创建聊天记录归档")
        chat_record = ChatHistory(
            user_id=user_id,
            content=f"{file.filename} 解析",
        )
        session.add(chat_record)
        session.commit()
        session.refresh(chat_record)
        chat_id = chat_record.id
        print(f"  创建聊天记录: ID={chat_id}")

        document_cache[document_id] = {
            "course_id": course.id,
            "script_id": course_script.id,
            "doc_id": docling_doc.id,
            "chat_id": chat_id,
            "filename": file.filename,
            "markdown_content": parse_result.markdown_content,
            "script_content": script_result.script_content,
            "file_path": str(file_path),
            "rag_knowledge_points": rag_result.knowledge_points,
        }

        # 持久化到数据库
        try:
            existing = session.exec(
                select(DocumentArtifact).where(DocumentArtifact.document_id == document_id)
            ).first()
            if not existing:
                artifact = DocumentArtifact(
                    document_id=document_id,
                    course_id=course.id if course else 0,
                    file_name=file.filename if file else "",
                    mime_type=file.content_type if file else "",
                    parse_info={"status": "uploaded"},
                )
                session.add(artifact)
                session.commit()
        except Exception:
            pass  # 持久化失败不影响主流程

        print(f"[步骤8] 返回结果给前端")
        
        return unified_response(
            code=200,
            message="上传并解析成功，TTS语音正在后台生成",
            data={
                "fullContent": script_result.beautiful_markdown,
                "rawContent": parse_result.markdown_content,
                "title": course.title,
                "audioUrl": None,
                "mindMapJson": mind_map,
                "chatId": chat_id,
                "courseId": course.id,
                "ttsStatus": "processing",
                "ragInfo": {
                    "formulaCount": rag_result.formula_count,
                    "tableCount": rag_result.table_count,
                    "domainTermCount": rag_result.domain_term_count,
                    "treeNodeCount": rag_result.tree_node_count,
                    "knowledgePointCount": len(rag_result.knowledge_points),
                } if not rag_result.error else None,
            }
        )

    except HTTPException:
        raise
    except Exception as e:
        import traceback
        traceback.print_exc()
        return unified_response(
            code=500,
            message=f"文档处理失败: {str(e)}",
            data={"error": str(e)}
        )


@router.post("/analyze", response_model=DocumentAnalyzeResponse)
async def analyze_document(
    request: DocumentAnalyzeRequest,
    session: Session = Depends(get_session),
    current_user: dict = Depends(get_current_user),
):
    """
    对已有文档进行AI分析
    """
    try:
        if request.document_id not in document_cache:
            # 尝试从数据库回填
            artifact = session.exec(
                select(DocumentArtifact).where(DocumentArtifact.document_id == request.document_id)
            ).first()
            if artifact:
                document_cache[request.document_id] = {
                    "course_id": artifact.course_id,
                    "file_name": artifact.file_name,
                    "parse_info": artifact.parse_info,
                }
            else:
                raise HTTPException(status_code=404, detail="文档不存在或已过期")

        doc_data = document_cache[request.document_id]
        course_id = doc_data.get("course_id")

        course = session.get(Course, course_id)
        if not course:
            raise HTTPException(status_code=404, detail="课程不存在")
        require_course_permission(session, current_user, course_id, "course.content.read")

        # 优先从内存缓存取 script_content；重启后从 CourseScript 表重建
        script_content = doc_data.get("script_content")
        if not script_content:
            script = session.exec(
                select(CourseScript).where(
                    CourseScript.course_id == course_id,
                    CourseScript.is_active == True,
                )
            ).first()
            if script:
                script_content = {"summary": script.summary_text or "无摘要"}

        return DocumentAnalyzeResponse(
            success=True,
            message="分析完成",
            analysis=(script_content or {}).get("summary", "无摘要")
        )
        
    except HTTPException:
        raise
    except Exception as e:
        return DocumentAnalyzeResponse(
            success=False,
            message="分析失败",
            error=str(e)
        )


# ==================== 课程列表接口 (必须在 /{document_id} 之前) ====================

@router.get("/courses")
async def get_courses_list(
    status: Optional[str] = Query(None, description="课程状态筛选 (published/draft/archived)"),
    session: Session = Depends(get_session),
    current_user: dict = Depends(get_current_user),
):
    """
    获取课程列表（学生可查看已发布的课程）

    需要用户登录认证
    学生只能看到已发布的课程
    老师可以看到自己创建的所有课程 + 其他老师的已发布课程
    """
    try:
        user_id = int(current_user["user_id"])
        visible_memberships = session.exec(select(CourseMembership.course_id).where(
            CourseMembership.user_id == user_id,
            CourseMembership.status == MembershipStatus.ACTIVE,
        )).all()
        from sqlmodel import or_
        statement = select(Course).where(or_(
            Course.status == CourseStatus.PUBLISHED,
            Course.id.in_(list(visible_memberships) or [-1]),
        ))

        if status:
            try:
                status_enum = CourseStatus(status)
                statement = statement.where(Course.status == status_enum)
            except ValueError:
                pass

        statement = statement.order_by(Course.created_at.desc())
        courses = session.exec(statement).all()

        courses_data = []
        for course in courses:
            owner = session.exec(select(CourseMembership).where(
                CourseMembership.course_id == course.id,
                CourseMembership.role == CourseRole.OWNER,
                CourseMembership.status == MembershipStatus.ACTIVE,
            )).first()
            teacher_name = "未知教师"
            teacher_record = session.execute(
                text("SELECT username FROM users WHERE id = :uid"),
                {"uid": owner.user_id}
            ).fetchone() if owner else None
            if teacher_record:
                teacher_name = teacher_record[0]

            student_count = session.exec(
                select(func.count()).select_from(StudentEnrollment).where(
                    StudentEnrollment.course_id == course.id,
                    StudentEnrollment.is_active == True
                )
            ).one()

            courses_data.append({
                "id": course.id,
                "title": course.title,
                "description": course.description,
                "status": course.status.value,
                "teacher_id": owner.user_id if owner else None,
                "teacher_name": teacher_name,
                "total_nodes": course.total_nodes,
                "total_duration": course.total_duration,
                "source_file_name": course.source_file_name,
                "is_ai_generated": course.is_ai_generated,
                "student_count": student_count,
                "created_at": course.created_at.isoformat() if course.created_at else None,
            })

        return unified_response(
            code=200,
            message="获取课程列表成功",
            data={
                "courses": courses_data,
                "total": len(courses_data),
            }
        )

    except Exception as e:
        import traceback
        traceback.print_exc()
        return unified_response(
            code=500,
            message=f"获取课程列表失败: {str(e)}",
            data=None
        )


@router.get("/my-courses")
async def get_my_courses(
    session: Session = Depends(get_session),
    current_user: dict = Depends(get_current_user),
):
    """
    获取当前学生的选课列表

    返回学生已选的课程及学习进度，用于学生端"我的课程"展示
    """
    try:
        student_id = int(current_user["user_id"])
        memberships = session.exec(
            select(CourseMembership).where(
                CourseMembership.user_id == student_id,
                CourseMembership.role == CourseRole.STUDENT,
                CourseMembership.status == MembershipStatus.ACTIVE,
            ).order_by(CourseMembership.joined_at.desc())
        ).all()

        my_courses = []
        for membership in memberships:
            course = session.get(Course, membership.course_id)
            if not course:
                continue
            enr = session.exec(select(StudentEnrollment).where(
                StudentEnrollment.student_id == student_id,
                StudentEnrollment.course_id == course.id,
                StudentEnrollment.is_active == True,
            )).first()
            # Do not turn a partial historical migration into a 500 response
            # or invent an enrolment record during a read operation.
            if enr is None:
                continue

            my_courses.append({
                "enrollment_id": enr.id,
                "course_id": course.id,
                "title": course.title,
                "description": course.description,
                "teacher_name": _get_teacher_name(session, next((m.user_id for m in session.exec(select(CourseMembership).where(CourseMembership.course_id == course.id, CourseMembership.role == CourseRole.OWNER, CourseMembership.status == MembershipStatus.ACTIVE)).all()), None)),
                "total_nodes": course.total_nodes,
                "total_duration": course.total_duration,
                "overall_progress": round(enr.overall_progress, 1),
                "avg_understanding_score": round(enr.avg_understanding_score * 100, 1) if enr.avg_understanding_score else 0,
                "total_study_minutes": enr.total_study_minutes,
                "enrolled_at": enr.enrolled_at.isoformat() if enr.enrolled_at else None,
                "last_study_time": enr.last_study_time.isoformat() if enr.last_study_time else None,
            })

        return unified_response(
            code=200,
            message="获取我的课程成功",
            data={
                "courses": my_courses,
                "total": len(my_courses),
            }
        )

    except Exception as e:
        import traceback
        traceback.print_exc()
        return unified_response(code=500, message=f"获取课程失败: {str(e)}", data=None)

@router.get("/{document_id}")
async def get_document(
    document_id: str,
    session: Session = Depends(get_session),
    current_user: dict = Depends(get_current_user),
):
    """
    获取文档信息（从数据库读取）
    """
    if document_id not in document_cache:
        # 尝试从数据库回填
        artifact = session.exec(
            select(DocumentArtifact).where(DocumentArtifact.document_id == document_id)
        ).first()
        if artifact:
            document_cache[document_id] = {
                "course_id": artifact.course_id,
                "file_name": artifact.file_name,
                "parse_info": artifact.parse_info,
            }
        else:
            raise HTTPException(status_code=404, detail="文档不存在或已过期")

    doc_data = document_cache[document_id]
    course_id = doc_data.get("course_id")

    course = session.get(Course, course_id)
    if not course:
        raise HTTPException(status_code=404, detail="课程不存在")

    require_course_permission(session, current_user, course_id, "course.content.read")

    script = session.exec(
        select(CourseScript).where(CourseScript.course_id == course_id).where(CourseScript.is_active == True)
    ).first()

    nodes = []
    if script:
        nodes = session.exec(
            select(ScriptNode).where(ScriptNode.script_id == script.id).order_by(ScriptNode.node_index)
        ).all()

    return {
        "document_id": document_id,
        "course": {
            "id": course.id,
            "title": course.title,
            "status": course.status.value,
            "total_nodes": course.total_nodes,
            "total_duration": course.total_duration,
        },
        "script": {
            "id": script.id if script else None,
            "version": script.version if script else None,
            "summary": script.summary_text if script else None,
        } if script else None,
        "nodes_count": len(nodes),
        "nodes": [
            {
                "index": n.node_index,
                "type": n.node_type.value,
                "title": n.title,
                "duration": n.duration,
            }
            for n in nodes
        ],
    }


@router.post("/course/{course_id}/save")
async def save_course_nodes(
    course_id: int,
    request: Request,
    session: Session = Depends(get_session),
    access: CourseAccessContext = Depends(course_permission("course.script.edit")),
):
    """
    保存老师修改后的课程节点内容到数据库

    需要老师权限
    """
    try:
        course = session.get(Course, course_id)
        if not course:
            return unified_response(code=404, message="课程不存在")

        body = await request.json()
        nodes_data = body.get("nodes", [])

        updated_count = 0
        for node_data in nodes_data:
            node_id = node_data.get("id")
            if not node_id:
                continue

            node = session.get(ScriptNode, node_id)
            if node and node.script_id:
                script = session.get(CourseScript, node.script_id)
                if script and script.course_id == course_id:
                    if "title" in node_data:
                        node.title = node_data["title"]
                    if "content" in node_data:
                        node.content = node_data["content"]
                    if "page_start" in node_data:
                        node.page_start = node_data["page_start"]
                    if "page_end" in node_data:
                        node.page_end = node_data["page_end"]
                    if "extra_data" in node_data:
                        node.extra_data = node_data["extra_data"]
                    session.add(node)
                    updated_count += 1

        session.commit()

        return unified_response(code=200, message=f"成功保存 {updated_count} 个节点的修改", data={"updated_count": updated_count})

    except Exception as e:
        import traceback
        traceback.print_exc()
        return unified_response(code=500, message=f"保存失败: {str(e)}", data=None)


@router.get("/course/{course_id}")
async def get_course_detail(
    course_id: int,
    session: Session = Depends(get_session),
    _access: CourseAccessContext = Depends(course_permission("course.view")),
):
    """
    获取课程完整详情（包括脚本和节点）
    """
    course = session.get(Course, course_id)
    if not course:
        raise HTTPException(status_code=404, detail="课程不存在")

    script = session.exec(
        select(CourseScript).where(CourseScript.course_id == course_id).where(CourseScript.is_active == True)
    ).first()
    
    nodes = []
    if script:
        nodes = session.exec(
            select(ScriptNode).where(ScriptNode.script_id == script.id).order_by(ScriptNode.node_index)
        ).all()
    
    docling_doc = session.exec(
        select(DoclingDocument).where(DoclingDocument.course_id == course_id)
    ).first()

    # 查询持久化的 document_id（上传 UUID），与 GET /{document_id} 端点契约一致
    artifact = session.exec(
        select(DocumentArtifact).where(DocumentArtifact.course_id == course_id)
    ).first()
    document_id = artifact.document_id if artifact else None
    
    return unified_response(
        code=200,
        message="获取课程详情成功",
        data={
            "document_id": document_id,
            "course": {
                "id": course.id,
                "title": course.title,
                "description": course.description,
                "status": course.status.value,
                "source_file_name": course.source_file_name,
                "total_pages": course.total_pages,
                "total_nodes": course.total_nodes,
                "total_duration": course.total_duration,
                "created_at": course.created_at.isoformat() if course.created_at else None,
            },
            "script": {
                "id": script.id,
                "version": script.version,
                "version_name": script.version_name,
                "summary_text": script.summary_text,
                "keywords": script.keywords,
                "script_content": script.script_content,
                "audio_url": script.audio_url,
                "audio_duration": script.audio_duration,
            } if script else None,
            "nodes": [
                {
                    "id": n.id,
                    "node_index": n.node_index,
                    "node_type": n.node_type.value,
                    "title": n.title,
                    "content": n.content,
                    "page_start": n.page_start,
                    "page_end": n.page_end,
                    "duration": n.duration,
                    "is_key_point": n.is_key_point,
                    "extra_data": n.extra_data,
                    "audio_url": n.audio_url,
                    "audio_duration": n.audio_duration,
                }
                for n in nodes
            ],
            "parse_info": {
                "status": docling_doc.status.value if docling_doc else None,
                "total_texts": docling_doc.total_texts if docling_doc else 0,
                "total_tables": docling_doc.total_tables if docling_doc else 0,
                "total_pictures": docling_doc.total_pictures if docling_doc else 0,
            } if docling_doc else None,
        }
    )


# ==================== TTS语音合成接口 ====================

from app.common.tts_client import tts_client
from fastapi.responses import Response


@router.post("/tts/synthesize")
async def synthesize_speech(
    request: Request,
):
    """
    语音合成接口（支持长文本，自动分段合成）

    将文本转换为语音，返回音频文件。
    长文本（>2000字）会自动分段合成后拼接。

    请求体(JSON):
    - text: 要合成的文本（支持长文本）
    - voice: 音色（可选，默认使用配置中的音色）
    - sample_rate: 采样率（可选，默认16000）
    - output_format: 输出格式（可选，默认mp3）

    返回:
    - 音频文件（二进制数据）
    """
    try:
        body = await request.json()
        text = body.get("text", "")
        voice = body.get("voice")
        sample_rate = body.get("sample_rate", 16000)
        output_format = body.get("output_format", "mp3")
        text_length = len(text)
        print(f"[TTS] 开始合成语音: 文本长度={text_length}字, 前50字: {text[:50]}...")
        
        # 长文本分段处理（每段不超过2000字，按句号分割）
        MAX_SEGMENT_LENGTH = 2000
        if text_length <= MAX_SEGMENT_LENGTH:
            response = await tts_client.synthesize(
                text=text,
                voice=voice,
                sample_rate=sample_rate,
                output_format=output_format
            )
            audio_data = response.audio_data
            total_latency = response.latency_ms
        else:
            # 按句号分段
            segments = []
            current_segment = ""
            for sentence in text.replace("。", "。\n").replace("！", "！\n").replace("？", "？\n").replace("；", "；\n").split("\n"):
                if len(current_segment) + len(sentence) > MAX_SEGMENT_LENGTH and current_segment:
                    segments.append(current_segment.strip())
                    current_segment = sentence
                else:
                    current_segment += sentence
            if current_segment.strip():
                segments.append(current_segment.strip())
            
            print(f"[TTS] 长文本分段: {len(segments)}段")
            
            # 逐段合成
            audio_parts = []
            total_latency = 0
            for i, seg in enumerate(segments):
                seg_response = await tts_client.synthesize(
                    text=seg,
                    voice=voice,
                    sample_rate=sample_rate,
                    output_format=output_format
                )
                audio_parts.append(seg_response.audio_data)
                total_latency += seg_response.latency_ms
                print(f"[TTS] 段{i+1}/{len(segments)}完成, {len(seg_response.audio_data)}字节")
            
            audio_data = b"".join(audio_parts)
        
        print(f"[TTS] 合成成功，音频大小: {len(audio_data)} 字节")
        
        content_type_map = {
            "mp3": "audio/mpeg",
            "wav": "audio/wav",
            "pcm": "audio/pcm",
        }
        
        content_type = content_type_map.get(output_format.lower(), "audio/mpeg")
        
        return Response(
            content=audio_data,
            media_type=content_type,
            headers={
                "Content-Disposition": f"attachment; filename=speech.{output_format}",
                "X-Latency-Ms": str(total_latency),
            }
        )
        
    except Exception as e:
        import traceback
        error_detail = traceback.format_exc()
        print(f"[TTS] 错误详情:\n{error_detail}")

        error_msg = str(e)
        if "401" in error_msg or "grant not found" in error_msg or "authentication" in error_msg.lower():
            raise HTTPException(
                status_code=503,
                detail="语音合成服务认证失败，TTS API凭证可能已过期，请联系管理员更新配置"
            )
        elif "timeout" in error_msg.lower() or "超时" in error_msg:
            raise HTTPException(
                status_code=504,
                detail="语音合成服务响应超时，请稍后重试"
            )
        else:
            raise HTTPException(status_code=500, detail=f"语音合成失败: {error_msg}")


# ==================== TTS健康检查接口 ====================


@router.get("/tts/health")
async def tts_health_check(
    current_user: dict = Depends(get_current_user),
):
    try:
        from app.common.tts_client import tts_client, TTSError

        test_response = await tts_client.synthesize(
            text="测试",
            output_format="mp3",
        )

        return unified_response(code=200, message="TTS服务正常", data={
            "provider": settings.TTS_PROVIDER,
            "status": "healthy",
            "audio_size": len(test_response.audio_data),
        })

    except TTSError as e:
        error_msg = str(e)
        is_auth_error = "401" in error_msg or "grant not found" in error_msg
        return unified_response(code=200, message="TTS服务异常", data={
            "provider": settings.TTS_PROVIDER,
            "status": "auth_error" if is_auth_error else "error",
            "error": error_msg,
        })
    except Exception as e:
        return unified_response(code=200, message="TTS服务异常", data={
            "provider": getattr(settings, 'TTS_PROVIDER', 'unknown'),
            "status": "error",
            "error": str(e),
        })


@router.get("/course/{course_id}/tts-status")
async def get_tts_generation_status(
    course_id: int,
    session: Session = Depends(get_session),
    _access: CourseAccessContext = Depends(course_permission("course.content.read")),
):
    status_key = str(course_id)
    status_info = tts_generation_status.get(status_key, None)

    if status_info is None:
        course = session.get(Course, course_id)
        if not course:
            return unified_response(code=404, message="课程不存在", data=None)

        script = session.exec(
            select(CourseScript).where(
                CourseScript.course_id == course_id,
                CourseScript.is_active == True,
            )
        ).first()

        if not script:
            return unified_response(code=200, message="无脚本", data={
                "status": "no_script", "total": 0, "completed": 0, "errors": [],
            })

        nodes = session.exec(
            select(ScriptNode).where(ScriptNode.script_id == script.id)
        ).all()

        nodes_with_audio = sum(1 for n in nodes if n.audio_url)
        total_nodes = len(nodes)

        if nodes_with_audio == total_nodes and total_nodes > 0:
            return unified_response(code=200, message="TTS已完成", data={
                "status": "completed", "total": total_nodes, "completed": nodes_with_audio, "errors": [],
            })
        elif nodes_with_audio > 0:
            return unified_response(code=200, message="TTS部分完成", data={
                "status": "partial", "total": total_nodes, "completed": nodes_with_audio, "errors": [],
            })
        else:
            return unified_response(code=200, message="TTS未开始", data={
                "status": "not_started", "total": total_nodes, "completed": 0, "errors": [],
            })

    return unified_response(code=200, message="获取TTS状态成功", data=status_info)


# ==================== 脚本版本管理接口 ====================

import copy


@router.post("/course/{course_id}/script/snapshot")
async def create_script_snapshot(
    course_id: int,
    request: Request,
    session: Session = Depends(get_session),
    access: CourseAccessContext = Depends(course_permission("course.script.edit")),
):
    """
    创建脚本版本快照
    
    将当前激活脚本的所有节点保存为新版本。
    请求体: { "version_name": "可选版本名称" }
    """
    try:
        course = session.get(Course, course_id)
        if not course:
            return unified_response(code=404, message="课程不存在", data=None)
        # 获取当前激活脚本
        active_script = session.exec(
            select(CourseScript).where(
                CourseScript.course_id == course_id,
                CourseScript.is_active == True,
            )
        ).first()
        if not active_script:
            return unified_response(code=404, message="未找到激活脚本", data=None)

        # 获取当前所有节点
        current_nodes = session.exec(
            select(ScriptNode).where(ScriptNode.script_id == active_script.id).order_by(ScriptNode.node_index)
        ).all()

        body = await request.json() if request else {}
        version_name = body.get("version_name", f"v{active_script.version + 1} 快照")

        # 创建新脚本版本（深拷贝script_content）
        new_script = CourseScript(
            course_id=course_id,
            version=active_script.version + 1,
            version_name=version_name,
            script_content=copy.deepcopy(active_script.script_content),
            summary_text=active_script.summary_text,
            keywords=active_script.keywords,
            is_active=True,
            created_by=access.user_id,
        )
        # 旧版本设为非激活
        active_script.is_active = False
        session.add(active_script)
        session.add(new_script)
        session.flush()  # 获取new_script.id

        # 拷贝所有节点到新脚本
        for node in current_nodes:
            new_node = ScriptNode(
                script_id=new_script.id,
                chapter_id=node.chapter_id,
                node_index=node.node_index,
                node_type=node.node_type,
                title=node.title,
                content=node.content,
                page_start=node.page_start,
                page_end=node.page_end,
                timestamp_start=node.timestamp_start,
                timestamp_end=node.timestamp_end,
                duration=node.duration,
                is_key_point=node.is_key_point,
                extra_data=copy.deepcopy(node.extra_data) if node.extra_data else None,
            )
            session.add(new_node)

        session.commit()
        session.refresh(new_script)

        return unified_response(
            code=200,
            message=f"已创建版本快照: {version_name}",
            data={
                "version": new_script.version,
                "version_name": new_script.version_name,
                "script_id": new_script.id,
                "node_count": len(current_nodes),
            },
        )

    except Exception as e:
        import traceback
        traceback.print_exc()
        return unified_response(code=500, message=f"创建快照失败: {str(e)}", data=None)


@router.get("/course/{course_id}/script/versions")
async def get_script_versions(
    course_id: int,
    session: Session = Depends(get_session),
    _access: CourseAccessContext = Depends(course_permission("course.script.edit")),
):
    """
    获取脚本版本列表
    """
    try:
        scripts = session.exec(
            select(CourseScript)
            .where(CourseScript.course_id == course_id)
            .order_by(CourseScript.version.desc())
        ).all()

        versions = []
        for s in scripts:
            node_count = len(session.exec(
                select(ScriptNode).where(ScriptNode.script_id == s.id)
            ).all())
            versions.append({
                "id": s.id,
                "version": s.version,
                "version_name": s.version_name,
                "is_active": s.is_active,
                "node_count": node_count,
                "created_at": s.created_at.isoformat() if s.created_at else None,
            })

        return unified_response(code=200, message="获取版本列表成功", data=versions)

    except Exception as e:
        return unified_response(code=500, message=f"获取版本列表失败: {str(e)}", data=None)


@router.post("/course/{course_id}/script/rollback/{script_id}")
async def rollback_script_version(
    course_id: int,
    script_id: int,
    session: Session = Depends(get_session),
    _access: CourseAccessContext = Depends(course_permission("course.rollback")),
):
    """
    回滚到指定脚本版本
    
    将指定版本设为激活，当前激活版本设为非激活。
    """
    try:
        course = session.get(Course, course_id)
        if not course:
            return unified_response(code=404, message="课程不存在", data=None)
        target_script = session.get(CourseScript, script_id)
        if not target_script or target_script.course_id != course_id:
            return unified_response(code=404, message="目标版本不存在", data=None)

        # 取消当前激活版本
        current_active = session.exec(
            select(CourseScript).where(
                CourseScript.course_id == course_id,
                CourseScript.is_active == True,
            )
        ).first()
        if current_active:
            current_active.is_active = False
            session.add(current_active)

        # 激活目标版本
        target_script.is_active = True
        session.add(target_script)
        session.commit()

        return unified_response(
            code=200,
            message=f"已回滚到版本 v{target_script.version}",
            data={"version": target_script.version, "version_name": target_script.version_name},
        )

    except Exception as e:
        return unified_response(code=500, message=f"回滚失败: {str(e)}", data=None)


# ==================== 课程发布与选课管理接口 ====================

@router.post("/course/{course_id}/publish")
async def publish_course(
    course_id: int,
    session: Session = Depends(get_session),
    _access: CourseAccessContext = Depends(course_permission("course.publish")),
):
    """
    发布课程（老师操作）
    
    将课程状态从 draft 改为 published，学生可以看到并选择该课程
    """
    try:
        course = session.get(Course, course_id)
        if not course:
            return unified_response(code=404, message="课程不存在")

        if course.status == CourseStatus.PUBLISHED:
            return unified_response(code=200, message="课程已是发布状态", data={"status": "published"})

        # 更新状态为已发布
        course.status = CourseStatus.PUBLISHED
        course.updated_at = utcnow_aware()
        session.add(course)
        session.commit()

        return unified_response(code=200, message="课程发布成功", data={"status": "published"})

    except Exception as e:
        import traceback
        traceback.print_exc()
        return unified_response(code=500, message=f"发布失败: {str(e)}", data=None)


@router.post("/course/{course_id}/unpublish")
async def unpublish_course(
    course_id: int,
    session: Session = Depends(get_session),
    _access: CourseAccessContext = Depends(course_permission("course.unpublish")),
):
    """
    取消发布课程（老师操作）
    
    将课程状态从 published 改为 draft，学生无法再看到该课程
    已选课的学生保留选课记录但标记为不活跃
    """
    try:
        course = session.get(Course, course_id)
        if not course:
            return unified_response(code=404, message="课程不存在")

        # 更新状态为草稿
        course.status = CourseStatus.DRAFT
        course.updated_at = utcnow_aware()
        session.add(course)
        session.commit()

        return unified_response(code=200, message="课程已取消发布", data={"status": "draft"})

    except Exception as e:
        import traceback
        traceback.print_exc()
        return unified_response(code=500, message=f"取消发布失败: {str(e)}", data=None)


@router.delete("/course/{course_id}")
async def delete_course(
    course_id: int,
    session: Session = Depends(get_session),
    _access: CourseAccessContext = Depends(course_permission("course.delete")),
):
    """
    删除课程（老师操作）

    完全删除课程及其所有关联数据（脚本、选课记录、学习进度等）
    此操作不可恢复，需要二次确认
    """
    try:
        course = session.get(Course, course_id)
        if not course:
            return unified_response(code=404, message="课程不存在")

        # 检查是否有学生已选课
        enrollments_count = session.exec(
            select(func.count()).select_from(StudentEnrollment).where(
                StudentEnrollment.course_id == course_id,
                StudentEnrollment.is_active == True
            )
        ).one() or 0

        # 删除所有相关数据（按依赖顺序）
        try:
            from app.models.progress_model import LearningProgress, NodeProgress, UnderstandingAnalysis
            from app.models.course_model import CourseScript, ScriptNode, DoclingDocument, DoclingGroup, DoclingTable, DoclingTableCell, DoclingText, DoclingPicture
            from app.models.video_generation_model import VideoGenerationTask
            from app.models.qa_model import QASession, QAMessage
            from app.models.user_model import ChatHistory, ChatMessage
            from app.models.access_control_model import CourseCapability, CourseMembership
            from app.models.question_bank_model import QuestionAttempt, QuestionBankItem, QuestionSourceMapping
            from app.models.cognitive_state_model import CognitiveState, LearningEvidenceRecord, RecommendationRecord
            from app.models.visualization_model import VisualizationPlanRecord
            from app.models.safety_policy_model import CourseSafetyPolicy, CourseSandboxPolicy, SafetyAuditLog
            from app.models.web_research_model import WebResearchConfig, WebResearchResult, ExternalReference
            from app.models.media_timeline_model import MediaAsset, MediaTimelineCue
            from app.models.graph_production_model import CourseEvidenceRecord, GraphSnapshotRecord, GraphNodeReview
            from app.models.mapping_model import KnowledgePageMap
            from app.models.note_model import Note
            from app.models.confirmation_model import CourseConfirmation
            from app.models.feedback_model import Feedback
            from app.models.progress_model import LearningJumpHistory
            from app.models.qa_model import QAContext
            from app.services.graph_production_service import mark_evidence_stale

            # Remove Phase B--E course-scoped records before their legacy
            # parents (scripts and nodes).  SQLite test databases do not
            # always enforce FKs; production databases do.
            artifacts = session.exec(
                select(DocumentArtifact).where(DocumentArtifact.course_id == course_id)
            ).all()
            for artifact in artifacts:
                # Keep the lifecycle operation atomic.  The rows are then
                # deleted with the course, so inaccessible citations cannot
                # resolve to a replacement document.
                mark_evidence_stale(
                    session,
                    course_id,
                    artifact.document_id,
                    reason="course_deleted",
                    commit=False,
                )

            for model in (GraphNodeReview, GraphSnapshotRecord, CourseEvidenceRecord):
                for record in session.exec(select(model).where(model.course_id == course_id)).all():
                    session.delete(record)

            for model in (ExternalReference, WebResearchResult, WebResearchConfig,
                          MediaTimelineCue, MediaAsset, VisualizationPlanRecord,
                          SafetyAuditLog, CourseSafetyPolicy, CourseSandboxPolicy,
                          RecommendationRecord, LearningEvidenceRecord, CognitiveState,
                          QuestionAttempt, QuestionSourceMapping, KnowledgePageMap,
                          Note, CourseConfirmation, Feedback, LearningJumpHistory,
                          QAContext):
                for record in session.exec(select(model).where(model.course_id == course_id)).all():
                    session.delete(record)

            for item in session.exec(
                select(QuestionBankItem).where(QuestionBankItem.course_id == course_id)
            ).all():
                session.delete(item)

            for artifact in artifacts:
                session.delete(artifact)

            for model in (CourseCapability, CourseMembership):
                for record in session.exec(select(model).where(model.course_id == course_id)).all():
                    session.delete(record)

            # 1. 删除理解度分析（依赖 learning_progress）
            learning_progresses = session.exec(
                select(LearningProgress).where(LearningProgress.course_id == course_id)
            ).all()
            for lp in learning_progresses:
                analyses = session.exec(
                    select(UnderstandingAnalysis).where(UnderstandingAnalysis.progress_id == lp.id)
                ).all()
                for analysis in analyses:
                    session.delete(analysis)

            # 2. 删除节点进度和学习进度
            for lp in learning_progresses:
                node_progresses = session.exec(
                    select(NodeProgress).where(NodeProgress.progress_id == lp.id)
                ).all()
                for np in node_progresses:
                    session.delete(np)
                session.delete(lp)

            # 3. 删除问答会话和消息
            qa_sessions = session.exec(
                select(QASession).where(QASession.course_id == course_id)
            ).all()
            for qs in qa_sessions:
                qa_messages = session.exec(
                    select(QAMessage).where(QAMessage.session_id == qs.id)
                ).all()
                for qm in qa_messages:
                    session.delete(qm)
                session.delete(qs)

            # 4. 删除视频生成任务
            video_tasks = session.exec(
                select(VideoGenerationTask).where(VideoGenerationTask.course_id == course_id)
            ).all()
            for vt in video_tasks:
                session.delete(vt)

            # 5. 删除选课记录
            all_enrollments = session.exec(
                select(StudentEnrollment).where(StudentEnrollment.course_id == course_id)
            ).all()
            for enrollment in all_enrollments:
                session.delete(enrollment)

            # 6. 删除课程脚本节点和脚本
            scripts = session.exec(
                select(CourseScript).where(CourseScript.course_id == course_id)
            ).all()
            for script in scripts:
                nodes = session.exec(
                    select(ScriptNode).where(ScriptNode.script_id == script.id)
                ).all()
                for node in nodes:
                    session.delete(node)
                session.delete(script)

            # 7. 删除 Docling 文档及其子表
            docling_docs = session.exec(
                select(DoclingDocument).where(DoclingDocument.course_id == course_id)
            ).all()
            for doc in docling_docs:
                for group in session.exec(
                    select(DoclingGroup).where(DoclingGroup.doc_id == doc.id)
                ).all():
                    session.delete(group)
                for tbl in session.exec(
                    select(DoclingTable).where(DoclingTable.doc_id == doc.id)
                ).all():
                    for cell in session.exec(
                        select(DoclingTableCell).where(DoclingTableCell.table_id == tbl.id)
                    ).all():
                        session.delete(cell)
                    session.delete(tbl)
                for txt in session.exec(
                    select(DoclingText).where(DoclingText.doc_id == doc.id)
                ).all():
                    session.delete(txt)
                for pic in session.exec(
                    select(DoclingPicture).where(DoclingPicture.doc_id == doc.id)
                ).all():
                    session.delete(pic)
                session.delete(doc)

            # 8. 删除关联的聊天历史（通过课程标题匹配）
            chat_histories = session.exec(
                select(ChatHistory).where(ChatHistory.content == course.title)
            ).all()
            for ch in chat_histories:
                chat_msgs = session.exec(
                    select(ChatMessage).where(ChatMessage.chat_id == ch.id)
                ).all()
                for cm in chat_msgs:
                    session.delete(cm)
                session.delete(ch)

            # 9. 最后删除课程本身
            session.delete(course)
            session.commit()

            # 10. 清理该课程的进程内 RAG 检索作用域（best-effort，不影响事务）
            try:
                from app.common.RAG import rag_pipeline
                rag_pipeline.clear_course_index(course_id)
            except Exception as clear_err:
                print(f"[删除课程] 清理课程 RAG 索引失败（可忽略）: {clear_err}")

            print(f"[删除课程] actor={_access.user_id} deleted course {course_id}; affected enrollments={enrollments_count}")

            return unified_response(code=200, message=f"课程《{course.title}》已成功删除", data={
                "deleted_course_id": course_id,
                "affected_students": enrollments_count,
            })

        except Exception as e:
            session.rollback()
            raise e

    except Exception as e:
        import traceback
        traceback.print_exc()
        return unified_response(code=500, message=f"删除失败: {str(e)}", data=None)


@router.post("/course/{course_id}/enroll")
async def enroll_course(
    course_id: int,
    session: Session = Depends(get_session),
    current_user: dict = Depends(get_current_user),
):
    """
    学生选择/加入课程

    学生可以加入老师发布的课程
    选课成功后会自动初始化学习进度记录
    """
    try:
        student_id = int(current_user["user_id"])

        # Enrollment is the lifecycle entrypoint for learner membership.
        # Allowing a teacher/owner token here would create a legacy
        # StudentEnrollment row that Course Access v1 correctly refuses to
        # treat as learner participation.
        from app.models.user_model import User, UserRole
        user = session.get(User, student_id)
        if user is None or not user.is_active:
            raise HTTPException(status_code=401, detail="Account is unavailable")
        if user.role != UserRole.STUDENT:
            raise HTTPException(status_code=403, detail="Only student accounts can enroll in a course")

        course = session.get(Course, course_id)
        if not course:
            return unified_response(code=404, message="课程不存在")

        if course.status != CourseStatus.PUBLISHED:
            return unified_response(code=400, message="该课程尚未发布，无法选择", data=None)

        # 检查是否已经选过
        existing = session.exec(
            select(StudentEnrollment).where(
                StudentEnrollment.student_id == student_id,
                StudentEnrollment.course_id == course_id
            )
        ).first()

        if existing and existing.is_active:
            activate_student_membership(session, course_id, student_id)
            session.commit()
            # 已选课，检查是否需要初始化进度数据
            _ensure_learning_progress(session, student_id, course_id, course.total_nodes or 0)
            return unified_response(code=200, message="您已选择过此课程", data={
                "enrollment_id": existing.id,
                "already_enrolled": True
            })

        # 如果有历史记录但不活跃，重新激活
        if existing and not existing.is_active:
            existing.is_active = True
            existing.enrolled_at = utcnow_aware()
            session.add(existing)
            activate_student_membership(session, course_id, student_id)
            session.commit()
            # 初始化学习进度
            _ensure_learning_progress(session, student_id, course_id, course.total_nodes or 0)
            return unified_response(code=200, message="重新加入课程成功", data={
                "enrollment_id": existing.id,
                "reactivated": True
            })

        # 创建新的选课记录
        enrollment = StudentEnrollment(
            student_id=student_id,
            course_id=course_id,
            total_nodes_count=course.total_nodes or 0,
        )
        session.add(enrollment)
        activate_student_membership(session, course_id, student_id)
        session.commit()
        session.refresh(enrollment)

        # 初始化学习进度记录（关键：确保学生学习数据能正确保存）
        _init_learning_progress_for_student(session, student_id, course_id, course.total_nodes or 0)

        print(f"[选课] 学生 {current_user.get('username')} 成功加入课程 {course.title} (ID:{course_id})")

        return unified_response(code=200, message="选课成功！您现在可以开始学习了。", data={
            "enrollment_id": enrollment.id,
            "enrolled_at": enrollment.enrolled_at.isoformat() if enrollment.enrolled_at else None,
            "total_nodes": course.total_nodes or 0,
        })

    except HTTPException:
        raise
    except Exception as e:
        import traceback
        traceback.print_exc()
        return unified_response(code=500, message=f"选课失败: {str(e)}", data=None)


def _init_learning_progress_for_student(session: Session, student_id: int, course_id: int, total_nodes: int):
    """
    为新选课的学生初始化学习进度记录

    创建LearningProgress和NodeProgress记录，确保后续学习数据能正确保存到数据库
    """
    try:
        from app.models.progress_model import LearningProgress, NodeProgress

        # 检查是否已有进度记录（使用正确的字段名 user_id）
        existing_progress = session.exec(
            select(LearningProgress).where(
                LearningProgress.user_id == student_id,
                LearningProgress.course_id == course_id
            )
        ).first()

        if existing_progress:
            print(f"[进度初始化] 学生{student_id} 课程{course_id} 进度记录已存在")
            return

        # 创建总的学习进度记录（使用正确的字段名）
        learning_progress = LearningProgress(
            user_id=student_id,
            course_id=course_id,
            current_node_index=0,
            completion_rate=0.0,
            total_learning_time=0,
            total_nodes=total_nodes,
            completed_nodes=0,
            status="not_started",
            last_accessed_at=utcnow_aware(),
        )
        session.add(learning_progress)
        session.commit()
        session.refresh(learning_progress)

        print(f"[进度初始化] 创建LearningProgress ID={learning_progress.id}")

        # 如果有节点信息，为每个节点创建初始进度记录
        if total_nodes > 0:
            from sqlmodel import select as sql_select
            script = session.exec(
                sql_select(CourseScript).where(CourseScript.course_id == course_id).where(CourseScript.is_active == True)
            ).first()

            if script:
                nodes = session.exec(
                    sql_select(ScriptNode).where(ScriptNode.script_id == script.id).order_by(ScriptNode.node_index)
                ).all()

                for node in nodes:
                    node_progress = NodeProgress(
                        progress_id=learning_progress.id,
                        node_id=node.id,
                        node_index=node.node_index,
                        is_completed=False,
                        understanding_score=0.0,
                        question_count=0,
                    )
                    session.add(node_progress)

                session.commit()
                print(f"[进度初始化] 为 {len(nodes)} 个节点创建初始进度记录")

    except Exception as e:
        print(f"[进度初始化] 失败: {e}")
        import traceback
        traceback.print_exc()


def _ensure_learning_progress(session: Session, student_id: int, course_id: int, total_nodes: int):
    """确保学生的学习进度记录存在"""
    try:
        from app.models.progress_model import LearningProgress

        existing = session.exec(
            select(LearningProgress).where(
                LearningProgress.user_id == student_id,
                LearningProgress.course_id == course_id
            )
        ).first()

        if not existing:
            _init_learning_progress_for_student(session, student_id, course_id, total_nodes)

    except Exception as e:
        print(f"[确保进度] 检查失败: {e}")


@router.post("/course/{course_id}/unenroll")
async def unenroll_course(
    course_id: int,
    session: Session = Depends(get_session),
    access: CourseAccessContext = Depends(course_permission("course.learn")),
):
    """
    学生退出课程
    """
    try:
        if access.role is None or access.role.value != "student":
            return unified_response(code=403, message="只有学生成员可以退出课程", data=None)
        student_id = access.user_id

        enrollment = session.exec(
            select(StudentEnrollment).where(
                StudentEnrollment.student_id == student_id,
                StudentEnrollment.course_id == course_id,
                StudentEnrollment.is_active == True
            )
        ).first()

        if not enrollment:
            return unified_response(code=400, message="未找到选课记录", data=None)

        enrollment.is_active = False
        session.add(enrollment)
        membership = session.exec(
            select(CourseMembership).where(
                CourseMembership.course_id == course_id,
                CourseMembership.user_id == student_id,
            )
        ).first()
        if membership is not None and membership.role.value == "student":
            membership.status = MembershipStatus.WITHDRAWN
            membership.left_at = utcnow_aware()
            membership.updated_at = utcnow_aware()
            session.add(membership)
        session.commit()

        return unified_response(code=200, message="已退出课程", data=None)

    except Exception as e:
        import traceback
        traceback.print_exc()
        return unified_response(code=500, message=f"退出课程失败: {str(e)}", data=None)


def _get_teacher_name(session: Session, teacher_id: int) -> str:
    """获取教师姓名"""
    from sqlmodel import text
    result = session.execute(
        text("SELECT username FROM users WHERE id = :uid"),
        {"uid": teacher_id}
    ).fetchone()
    return result[0] if result else "未知教师"


@router.get("/course/{course_id}/students")
async def get_course_students(
    course_id: int,
    session: Session = Depends(get_session),
    _access: CourseAccessContext = Depends(course_permission("analytics.view_member")),
):
    """
    获取课程的学生列表及学习进度（老师查看）
    
    返回所有选择了该课程的活跃学生及其学习进度统计
    """
    try:
        course = session.get(Course, course_id)
        if not course:
            return unified_response(code=404, message="课程不存在")

        # 查询所有活跃的选课记录
        enrollments = session.exec(
            select(StudentEnrollment).where(
                StudentEnrollment.course_id == course_id,
                StudentEnrollment.is_active == True
            ).order_by(StudentEnrollment.enrolled_at.desc())
        ).all()

        students_data = []
        for enr in enrollments:
            # 获取学生用户名
            from sqlmodel import text
            user_result = session.execute(
                text("SELECT username FROM users WHERE id = :uid"),
                {"uid": enr.student_id}
            ).fetchone()
            username = user_result[0] if user_result else f"学生{enr.student_id}"

            # 动态计算该学生的理解度（从NodeProgress表）
            student_avg_score = None
            student_level = "unknown"
            
            try:
                from app.models.progress_model import LearningProgress, NodeProgress
                
                # 获取学生的学习进度记录
                lp = session.exec(
                    select(LearningProgress).where(
                        LearningProgress.user_id == enr.student_id,
                        LearningProgress.course_id == course_id,
                    )
                ).first()
                
                if lp:
                    # 获取该学生所有节点的进度
                    node_progress_list = session.exec(
                        select(NodeProgress).where(
                            NodeProgress.progress_id == lp.id,
                            NodeProgress.understanding_score.isnot(None)
                        )
                    ).all()
                    
                    if node_progress_list:
                        # 计算平均理解度分数（0-1）
                        total_score = sum(np.understanding_score for np in node_progress_list)
                        avg_score = total_score / len(node_progress_list)
                        
                        # 转换为百分比（0-100）用于显示
                        student_avg_score = round(avg_score * 100, 1)
                        
                        # 确定理解度等级
                        if avg_score >= 0.9:
                            student_level = "excellent"
                        elif avg_score >= 0.75:
                            student_level = "high"
                        elif avg_score >= 0.5:
                            student_level = "medium"
                        elif avg_score > 0:
                            student_level = "low"
            except Exception as calc_error:
                print(f"[警告] 计算学生{enr.student_id}理解度失败: {calc_error}")

            students_data.append({
                "enrollment_id": enr.id,
                "student_id": enr.student_id,
                "username": username,
                "enrolled_at": enr.enrolled_at.isoformat() if enr.enrolled_at else None,
                "overall_progress": round(enr.overall_progress, 1),
                # 使用动态计算的理解度，而非可能过时的汇总字段
                "avg_understanding_score": student_avg_score if student_avg_score is not None else 0,
                "understanding_level": student_level if student_level != "unknown" else None,
                "total_study_minutes": enr.total_study_minutes or 0,
                "last_study_time": enr.last_study_time.isoformat() if enr.last_study_time else None,
                "nodes_completed": enr.total_nodes_completed or 0,
                "nodes_total": enr.total_nodes_count or 0,
            })

        return unified_response(code=200, message="获取成功", data={
            "course_id": course_id,
            "course_title": course.title,
            "total_students": len(students_data),
            "students": students_data
        })

    except Exception as e:
        import traceback
        traceback.print_exc()
        return unified_response(code=500, message=f"获取学生列表失败: {str(e)}", data=None)


@router.get("/course/{course_id}/stats")
async def get_course_stats(
    course_id: int,
    session: Session = Depends(get_session),
    _access: CourseAccessContext = Depends(course_permission("analytics.view_course")),
):
    """
    获取课程统计数据（老师查看）
    
    返回：总选课人数、平均进度、平均理解度等统计信息
    """
    try:
        course = session.get(Course, course_id)
        if not course:
            return unified_response(code=404, message="课程不存在")

        # 统计活跃选课数
        active_count = session.exec(
            select(func.count()).select_from(StudentEnrollment).where(
                StudentEnrollment.course_id == course_id,
                StudentEnrollment.is_active == True
            )
        ).one()

        # 计算平均进度和理解度
        all_enrollments = session.exec(
            select(StudentEnrollment).where(
                StudentEnrollment.course_id == course_id,
                StudentEnrollment.is_active == True
            )
        ).all()

        total_students = len(all_enrollments)
        avg_progress = 0.0
        avg_understanding = 0.0
        total_study_time = 0

        if total_students > 0:
            avg_progress = sum(e.overall_progress for e in all_enrollments) / total_students
            # avg_understanding_score 是0-1的值，转为百分比显示
            avg_understanding = sum(e.avg_understanding_score or 0 for e in all_enrollments) / total_students * 100
            # 计算总学习时长（分钟），后续会转为小时并计算平均值
            total_study_minutes_all = sum(e.total_study_minutes or 0 for e in all_enrollments)

        # 进度分布
        progress_distribution = {
            "not_started": sum(1 for e in all_enrollments if e.overall_progress == 0),
            "beginner": sum(1 for e in all_enrollments if 0 < e.overall_progress < 30),
            "intermediate": sum(1 for e in all_enrollments if 30 <= e.overall_progress < 70),
            "advanced": sum(1 for e in all_enrollments if 70 <= e.overall_progress < 100),
            "completed": sum(1 for e in all_enrollments if e.overall_progress >= 100),
        }

        # 按节点的学习进度统计
        node_progress_stats = []
        from app.models.course_model import ScriptNode, CourseScript
        script = session.exec(
            select(CourseScript).where(CourseScript.course_id == course_id)
        ).first()
        if script:
            nodes = session.exec(
                select(ScriptNode)
                .where(ScriptNode.script_id == script.id)
                .order_by(ScriptNode.node_index)
            ).all()

            for node in nodes:
                from app.models.progress_model import LearningProgress, NodeProgress
                completed_count = 0
                node_avg_understanding = 0.0
                total_accessed = 0

                for enr in all_enrollments:
                    lp = session.exec(
                        select(LearningProgress).where(
                            LearningProgress.user_id == enr.student_id,
                            LearningProgress.course_id == course_id,
                        )
                    ).first()
                    if lp:
                        np = session.exec(
                            select(NodeProgress).where(
                                NodeProgress.progress_id == lp.id,
                                NodeProgress.node_id == node.id,
                            )
                        ).first()
                        if np:
                            total_accessed += 1
                            if np.is_completed:
                                completed_count += 1
                            if np.understanding_score is not None:
                                node_avg_understanding += np.understanding_score

                # 节点理解度：understanding_score是0-1，转为百分比
                node_avg_understanding = (node_avg_understanding / total_accessed * 100) if total_accessed > 0 else 0
                completion_rate = (completed_count / total_students * 100) if total_students > 0 else 0

                node_progress_stats.append({
                    "node_id": node.id,
                    "node_index": node.node_index,
                    "title": node.title or f"节点 {node.node_index + 1}",
                    "node_type": node.node_type.value if node.node_type else "lecture",
                    "is_key_point": node.is_key_point,
                    "completed_count": completed_count,
                    "total_students": total_students,
                    "completion_rate": round(completion_rate, 1),
                    "avg_understanding": round(node_avg_understanding, 1),
                    "accessed_count": total_accessed,
                })

        return unified_response(code=200, message="获取成功", data={
            "course_id": course_id,
            "course_title": course.title,
            "status": course.status.value,
            "total_students": total_students,
            "total_nodes": course.total_nodes or 0,
            "avg_progress": round(avg_progress, 1),
            # avg_understanding 已经在计算时转为百分比（0-100）
            "avg_understanding": round(avg_understanding, 1),
            # 改为学生平均学习时长（小时）
            "avg_study_hours_per_student": round(total_study_minutes_all / total_students / 60, 1) if total_students > 0 else 0,
            "progress_distribution": progress_distribution,
            "node_progress": node_progress_stats,
        })

    except Exception as e:
        import traceback
        traceback.print_exc()
        return unified_response(code=500, message=f"获取统计失败: {str(e)}", data=None)


PPT_SLIDES_DIR = Path(tempfile.gettempdir()) / "ai_course_ppt_slides"
PPT_SLIDES_DIR.mkdir(exist_ok=True)

BASE_DIR = Path(__file__).resolve().parent.parent.parent.parent
AUDIO_STORAGE_DIR = BASE_DIR / "audio_storage"
AUDIO_STORAGE_DIR.mkdir(exist_ok=True)


def get_course_audio_dir(course_id: int) -> Path:
    course_dir = AUDIO_STORAGE_DIR / str(course_id)
    course_dir.mkdir(parents=True, exist_ok=True)
    return course_dir


def cleanup_old_node_audio(node: ScriptNode, course_dir: Path):
    if node.audio_url:
        old_filename = node.audio_url.split("/")[-1]
        if old_filename:
            old_path = course_dir / old_filename
            if old_path.exists():
                try:
                    old_path.unlink()
                    logger.info(f"Cleaned up old audio for node {node.id}: {old_filename}")
                except OSError as e:
                    logger.warning(f"Failed to delete old audio {old_path}: {e}")


@router.get("/course/{course_id}/slides")
async def get_course_slides(
    course_id: int,
    session: Session = Depends(get_session),
    _access: CourseAccessContext = Depends(course_permission("course.content.read")),
):
    course = session.get(Course, course_id)
    if not course:
        return unified_response(code=404, message="课程不存在")

    source_path = course.source_file_path
    if not source_path or not Path(source_path).exists():
        return unified_response(code=200, message="无PPT文件", data={
            "course_id": course_id, "total_pages": 0, "slides": [],
        })

    course_slide_dir = PPT_SLIDES_DIR / str(course_id)
    slide_list_file = course_slide_dir / "slides_meta.json"

    if slide_list_file.exists():
        meta = json.loads(slide_list_file.read_text(encoding="utf-8"))
        if meta.get("source_mtime"):
            try:
                current_mtime = Path(source_path).stat().st_mtime
                if current_mtime <= meta["source_mtime"]:
                    return unified_response(code=200, message="获取成功", data=meta)
            except OSError:
                pass

    course_slide_dir.mkdir(parents=True, exist_ok=True)

    try:
        from app.common.slide_converter import get_or_create_pdf, render_pdf_to_images, is_pdf_file

        pdf_path = course.pdf_file_path
        if pdf_path and Path(pdf_path).exists():
            effective_pdf = pdf_path
        elif is_pdf_file(source_path):
            effective_pdf = source_path
        else:
            effective_pdf = get_or_create_pdf(source_path)
            if effective_pdf:
                course.pdf_file_path = effective_pdf
                session.add(course)
                session.commit()

        if effective_pdf:
            image_paths = render_pdf_to_images(effective_pdf, str(course_slide_dir), dpi=150)

            if image_paths:
                total_pages = len(image_paths)
                slides_info = []
                for i in range(total_pages):
                    slides_info.append({
                        "page": i + 1,
                        "url": f"/api/v1/document/course/{course_id}/slide/{i + 1}",
                    })

                source_mtime = Path(source_path).stat().st_mtime
                meta = {
                    "course_id": course_id,
                    "total_pages": total_pages,
                    "slides": slides_info,
                    "source_mtime": source_mtime,
                }
                slide_list_file.write_text(json.dumps(meta, ensure_ascii=False), encoding="utf-8")

                if course.total_pages != total_pages:
                    course.total_pages = total_pages
                    session.add(course)
                    session.commit()

                return unified_response(code=200, message="获取成功", data=meta)

        logger.warning(f"PDF conversion failed, falling back to text rendering for course {course_id}")
        return _fallback_text_slides(source_path, course_id, course_slide_dir, slide_list_file, session, course)

    except Exception as e:
        import traceback
        traceback.print_exc()
        return unified_response(code=500, message=f"转换幻灯片失败: {str(e)}", data=None)


def _fallback_text_slides(source_path, course_id, course_slide_dir, slide_list_file, session, course):
    try:
        from pptx import Presentation
        from PIL import Image, ImageDraw

        prs = Presentation(source_path)
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        total_pages = len(prs.slides)

        slides_info = []
        for i, slide in enumerate(prs.slides):
            img_filename = f"slide_{i + 1}.png"
            img_path = course_slide_dir / img_filename

            if not img_path.exists():
                scale = 2
                img_w = int(slide_width / 914400 * 96 * scale)
                img_h = int(slide_height / 914400 * 96 * scale)
                img = Image.new("RGB", (img_w, img_h), (255, 255, 255))
                draw = ImageDraw.Draw(img)

                y_offset = 40
                title_text = ""
                if slide.shapes.title and slide.shapes.title.text.strip():
                    title_text = slide.shapes.title.text.strip()
                    draw.text((40, y_offset), title_text, fill=(0, 0, 0))
                    y_offset += 60

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        if text != title_text:
                            draw.text((40, y_offset), text, fill=(51, 51, 51))
                            y_offset += 30

                img.save(str(img_path), "PNG")

            slides_info.append({
                "page": i + 1,
                "url": f"/api/v1/document/course/{course_id}/slide/{i + 1}",
            })

        source_mtime = Path(source_path).stat().st_mtime
        meta = {
            "course_id": course_id,
            "total_pages": total_pages,
            "slides": slides_info,
            "source_mtime": source_mtime,
        }
        slide_list_file.write_text(json.dumps(meta, ensure_ascii=False), encoding="utf-8")

        return unified_response(code=200, message="获取成功(文本降级模式)", data=meta)

    except Exception as e:
        import traceback
        traceback.print_exc()
        return unified_response(code=500, message=f"转换幻灯片失败: {str(e)}", data=None)


@router.get("/course/{course_id}/slide/{page_num}")
async def get_slide_image(
    course_id: int,
    page_num: int,
    session: Session = Depends(get_session),
    _access: CourseAccessContext = Depends(course_permission("course.content.read")),
):
    course_slide_dir = PPT_SLIDES_DIR / str(course_id)
    img_path = course_slide_dir / f"slide_{page_num}.png"

    if not img_path.exists():
        raise HTTPException(status_code=404, detail="幻灯片图片不存在")

    from fastapi.responses import FileResponse
    return FileResponse(
        str(img_path), media_type="image/png",
        filename=f"slide_{page_num}.png",
        headers={"Cache-Control": "public, max-age=3600"},
    )


@router.post("/course/{course_id}/node/{node_id}/synthesize-audio")
async def synthesize_node_audio(
    course_id: int,
    node_id: int,
    session: Session = Depends(get_session),
    _access: CourseAccessContext = Depends(course_permission("course.media.generate")),
):
    node = session.get(ScriptNode, node_id)
    if not node:
        raise HTTPException(status_code=404, detail="节点不存在")

    course = session.get(Course, course_id)
    if not course:
        raise HTTPException(status_code=404, detail="课程不存在")

    script = session.get(CourseScript, node.script_id)
    if script is None or script.course_id != course_id:
        raise HTTPException(status_code=400, detail="Node does not belong to course")

    if not node.content or len(node.content.strip()) < 10:
        return unified_response(code=400, message="节点内容过短，无法合成音频", data=None)

    try:
        from app.common.tts_client import tts_client
        from app.platform.adapters.tts import TTSAdapter

        course_audio_dir = get_course_audio_dir(course_id)
        cleanup_old_node_audio(node, course_audio_dir)

        content = node.content.strip()
        if len(content) > 2000:
            segments = []
            current = ""
            for char in content:
                current += char
                if char in "。！？；" and len(current) >= 500:
                    segments.append(current)
                    current = ""
            if current:
                segments.append(current)
        else:
            segments = [content]

        all_audio = b""
        total_latency = 0.0
        voice = None
        if node.extra_data:
            voice = node.extra_data.get("voice")

        for seg in segments:
            tts_result = await TaskRunner().run(
                TaskContext(
                    task_id=f"{course_id}:{node.id}",
                    task_type=TaskType.TTS_NODE,
                    course_id=course_id,
                    node_id=node.id,
                    provider="tts",
                    input_summary=seg[:120],
                    metadata={"stage": "single_node"},
                ),
                lambda seg=seg: TTSAdapter(tts_client).synthesize(
                    text=seg, voice=voice, sample_rate=16000, output_format="mp3",
                ),
            )
            if not tts_result.success:
                raise RuntimeError(tts_result.error_message or "TTS synthesis failed")
            response = tts_result.data
            all_audio += response.audio_data
            total_latency += response.latency_ms

        audio_filename = f"node_{node_id}_{uuid.uuid4().hex[:8]}.mp3"
        audio_path = course_audio_dir / audio_filename
        audio_path.write_bytes(all_audio)

        audio_url = f"/api/v1/document/audio/{course_id}/{audio_filename}"
        node.audio_url = audio_url

        import subprocess
        try:
            result = subprocess.run(
                ["ffprobe", "-v", "quiet", "-show_entries", "format=duration",
                 "-of", "default=noprint_wrappers=1:nokey=1", str(audio_path)],
                capture_output=True, text=True, timeout=10,
            )
            if result.returncode == 0 and result.stdout.strip():
                node.audio_duration = float(result.stdout.strip())
            else:
                node.audio_duration = len(all_audio) / 16000 / 2
        except Exception:
            node.audio_duration = len(all_audio) / 16000 / 2

        session.add(node)
        session.commit()
        session.refresh(node)

        return unified_response(code=200, message="音频合成成功", data={
            "node_id": node_id,
            "audio_url": node.audio_url,
            "audio_duration": node.audio_duration,
            "latency_ms": total_latency,
        })

    except Exception as e:
        import traceback
        traceback.print_exc()
        return unified_response(code=500, message=f"音频合成失败: {str(e)}", data=None)


@router.get("/audio/{course_id}/{filename}")
async def stream_audio(
    course_id: int,
    filename: str,
    session: Session = Depends(get_session),
    _access: CourseAccessContext = Depends(course_permission("course.content.read")),
):
    course_audio_dir = get_course_audio_dir(course_id)
    audio_path = course_audio_dir / filename
    if not audio_path.exists():
        legacy_path = AUDIO_STORAGE_DIR / filename
        if legacy_path.exists():
            audio_path = legacy_path
        else:
            raise HTTPException(status_code=404, detail="音频文件不存在")

    file_size = audio_path.stat().st_size
    ext = audio_path.suffix.lower()
    media_types = {".mp3": "audio/mpeg", ".wav": "audio/wav", ".pcm": "audio/pcm"}
    media_type = media_types.get(ext, "audio/mpeg")

    async def iterfile():
        with open(audio_path, "rb") as f:
            while chunk := f.read(1024 * 64):
                yield chunk

    from fastapi.responses import StreamingResponse
    return StreamingResponse(
        iterfile(), media_type=media_type,
        headers={
            "Content-Length": str(file_size),
            "Accept-Ranges": "bytes",
            "Cache-Control": "public, max-age=86400",
        },
    )


@router.post("/course/{course_id}/synthesize-all-audio")
async def synthesize_all_node_audio(
    course_id: int,
    session: Session = Depends(get_session),
    _access: CourseAccessContext = Depends(course_permission("course.media.generate")),
):
    course = session.get(Course, course_id)
    if not course:
        raise HTTPException(status_code=404, detail="课程不存在")

    script = session.exec(
        select(CourseScript).where(
            CourseScript.course_id == course_id,
            CourseScript.is_active == True,
        )
    ).first()

    if not script:
        return unified_response(code=400, message="课程无激活脚本", data=None)

    nodes = session.exec(
        select(ScriptNode)
        .where(ScriptNode.script_id == script.id)
        .order_by(ScriptNode.node_index)
    ).all()

    course_audio_dir = get_course_audio_dir(course_id)
    batch_started = time.perf_counter()
    total_count = len([n for n in nodes if n.content and len(n.content.strip()) >= 10])
    node_results: list[TaskResult] = []
    results = []
    errors = []

    for node in nodes:
        if not node.content or len(node.content.strip()) < 10:
            continue

        node_started = time.perf_counter()
        node_context = TaskContext(
            task_id=f"{course_id}:{node.id}",
            task_type=TaskType.TTS_BATCH,
            course_id=course_id,
            node_id=node.id,
            provider="tts",
            input_summary=node.content.strip()[:120],
            metadata={"stage": "sync_batch"},
        )
        failed_task_result = None
        try:
            from app.common.tts_client import tts_client

            cleanup_old_node_audio(node, course_audio_dir)

            content = node.content.strip()
            if len(content) > 2000:
                segments = []
                current = ""
                for char in content:
                    current += char
                    if char in "。！？；" and len(current) >= 500:
                        segments.append(current)
                        current = ""
                if current:
                    segments.append(current)
            else:
                segments = [content]

            all_audio = b""
            voice = None
            if node.extra_data:
                voice = node.extra_data.get("voice")

            for seg in segments:
                tts_result = await TaskRunner().run(
                    node_context,
                    lambda seg=seg: TTSAdapter(tts_client).synthesize(
                        text=seg, voice=voice, sample_rate=16000, output_format="mp3",
                    ),
                )
                if not tts_result.success:
                    failed_task_result = tts_result
                _raise_for_tts_task_failure(tts_result)
                response = tts_result.data
                all_audio += response.audio_data

            audio_filename = f"node_{node.id}_{uuid.uuid4().hex[:8]}.mp3"
            audio_path = course_audio_dir / audio_filename
            audio_path.write_bytes(all_audio)

            audio_url = f"/api/v1/document/audio/{course_id}/{audio_filename}"
            node.audio_url = audio_url
            node.audio_duration = len(all_audio) / 16000 / 2

            session.add(node)
            results.append({"node_id": node.id, "title": node.title, "audio_url": audio_url})
            node_results.append(TaskResult.ok(
                data={
                    "node_id": node.id,
                    "audio_url": node.audio_url,
                    "audio_duration": node.audio_duration,
                },
                provider="tts",
                duration_ms=(time.perf_counter() - node_started) * 1000,
                context=node_context,
            ))

        except Exception as e:
            node_results.append(_tts_node_failure_result(
                node_context,
                e,
                failed_task_result,
                (time.perf_counter() - node_started) * 1000,
            ))
            errors.append({"node_id": node.id, "title": node.title, "error": str(e)})

    session.commit()

    batch_task_result = _tts_batch_task_result(
        course_id,
        total_count,
        node_results,
        errors,
        (time.perf_counter() - batch_started) * 1000,
    )
    logger.info(
        "Synchronous TTS batch completed for course %s with internal status %s",
        course_id,
        batch_task_result.status.value,
    )

    return unified_response(code=200, message=f"批量合成完成: {len(results)}成功, {len(errors)}失败", data={
        "course_id": course_id,
        "success_count": len(results),
        "error_count": len(errors),
        "results": results,
        "errors": errors,
    })
    script = session.get(CourseScript, node.script_id)
    if script is None or script.course_id != course_id:
        raise HTTPException(status_code=400, detail="Node does not belong to course")
