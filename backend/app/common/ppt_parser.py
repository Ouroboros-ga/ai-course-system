import io
import logging
import tempfile
import zipfile
from dataclasses import dataclass, asdict, field
from pathlib import Path
from typing import Dict, List, Optional, Any, Union, BinaryIO

# 延迟导入，避免在模块加载时失败
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    Presentation = None  # 类型占位

# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.enum.shapes import MSO_SHAPE_TYPE
# HAS_PPTX = True

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class PPTParseError(Exception):
    """PPT 解析过程中的自定义异常"""
    pass


@dataclass
class SlideContent:
    """单页幻灯片的结构化数据容器"""
    slide_number: int
    title: str
    text_content: List[str]
    text_raw: str = ""
    images: Optional[List[str]] = None
    tables: Optional[List[List[List[str]]]] = None
    notes: str = ""

    def to_dict(self) -> Dict:
        """转换为字典"""
        return asdict(self)


@dataclass
class PPTParseResult:
    """整个PPT的解析结果容器"""
    file_name: str
    slide_count: int
    slides: List[SlideContent]
    metadata: Optional[Dict[str, Any]] = None

    def to_dict(self) -> Dict:
        """转换为字典"""
        return {
            "file_name": self.file_name,
            "slide_count": self.slide_count,
            "slides": [s.to_dict() for s in self.slides],
            "metadata": self.metadata or {}
        }


class PPTParser:
    def __init__(self, extract_images: bool = False, image_output_dir: Optional[str] = None):
        """
        初始化PPT解析器

        :param extract_images: 是否提取图片到磁盘
        :param image_output_dir: 图片输出目录（若为None则使用临时目录）
        """
        if not HAS_PPTX:
            raise ImportError("未安装 python-pptx 库，请执行 pip install python-pptx")

        self.extract_images = extract_images
        self._temp_dir_obj: Optional[tempfile.TemporaryDirectory] = None
        self.image_output_path: Optional[Path] = None

        if self.extract_images:
            if image_output_dir:
                # 使用用户指定目录
                self.image_output_path = Path(image_output_dir)
                self.image_output_path.mkdir(parents=True, exist_ok=True)
            else:
                # 创建临时目录（推荐使用上下文管理器清理）
                self._temp_dir_obj = tempfile.TemporaryDirectory()
                self.image_output_path = Path(self._temp_dir_obj.name)

            logger.info(f"图片将提取到: {self.image_output_path}")

    def __enter__(self):
        """支持上下文管理器"""
        return self

    def __exit__(self, exc_type, exc_val, exc_tb):
        """退出上下文时自动清理临时目录"""
        self.cleanup()
        return False

    def cleanup(self):
        """手动清理临时资源"""
        if self._temp_dir_obj:
            try:
                self._temp_dir_obj.cleanup()
                self._temp_dir_obj = None
                logger.info("临时图片目录已清理")
            except Exception as e:
                logger.warning(f"清理临时目录失败: {str(e)}")

    def parse(self, file_source: Union[str, BinaryIO, Any], filename: str = "unknown.pptx") -> PPTParseResult:
        """
        通用解析接口（支持文件路径、文件对象、FastAPI UploadFile）

        :param file_source: 文件路径或文件类对象
        :param filename: 原始文件名（用于记录）
        :raises PPTParseError: 解析失败时抛出
        """
        # 适配 FastAPI/Starlette 的 UploadFile 对象
        if hasattr(file_source, 'file'):
            # 保存引用以便读取内容
            file_source = file_source.file
            # 尝试从 UploadFile 获取真实文件名
            if hasattr(file_source, 'name') and filename == "unknown.pptx":
                filename = getattr(file_source, 'name', filename)

        try:
            # python-pptx 原生支持文件路径和文件对象
            prs = Presentation(file_source)

            slides_content = [
                self._parse_single_slide(slide, i)
                for i, slide in enumerate(prs.slides, 1)
            ]

            metadata = self._extract_metadata(prs)

            return PPTParseResult(
                file_name=filename,
                slide_count=len(prs.slides),
                slides=slides_content,
                metadata=metadata
            )
        except zipfile.BadZipFile:
            logger.error("无效的PPTX文件（非ZIP格式）")
            raise PPTParseError("无效的PPTX文件，请确保文件格式正确且未损坏")
        except Exception as e:
            logger.error(f"PPT解析失败: {str(e)}", exc_info=True)
            raise PPTParseError(f"PPT解析失败: {str(e)}")

    async def parse_fastapi_uploadfile(self, upload_file: Any) -> PPTParseResult:
        """
        FastAPI专用异步解析接口

        :param upload_file: FastAPI的UploadFile实例
        """
        # 读取文件内容到内存（解决异步文件对象与python-pptx的兼容性）
        content = await upload_file.read()
        file_like = io.BytesIO(content)

        # 优先使用上传文件的原始文件名
        filename = getattr(upload_file, 'filename', 'unknown.pptx')

        return self.parse(file_like, filename=filename)

    def _parse_single_slide(self, slide, slide_number: int) -> SlideContent:
        """解析单页幻灯片"""
        title = self._extract_title(slide)
        texts = []
        images = []
        tables = []

        for shape in slide.shapes:
            # 1. 提取表格 (优先级高于普通文本，避免重复提取)
            if shape.has_table:
                tables.append(self._parse_table(shape.table))
                continue

            # 2. 提取图片
            if self.extract_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_path = self._save_image(shape, slide_number)
                if img_path:
                    images.append(str(img_path))
                continue

            # 3. 提取文本 (排除标题和已处理的形状)
            if hasattr(shape, "text") and shape.text.strip():
                # 跳过标题形状，因为标题已经单独提取
                if shape == slide.shapes.title:
                    continue
                texts.append(shape.text.strip())

        # 4. 提取备注
        notes = self._extract_notes(slide)

        return SlideContent(
            slide_number=slide_number,
            title=title,
            text_content=texts,
            text_raw="\n".join(texts),
            images=images or None,
            tables=tables or None,
            notes=notes
        )

    def _extract_title(self, slide) -> str:
        """提取标题"""
        if slide.shapes.title and slide.shapes.title.text:
            return slide.shapes.title.text.strip()
        return ""

    def _extract_notes(self, slide) -> str:
        """提取备注"""
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                return notes_slide.notes_text_frame.text.strip()
        return ""

    def _save_image(self, shape, slide_number: int) -> Optional[Path]:
        """保存图片到磁盘，返回Path对象"""
        if not self.image_output_path:
            return None

        try:
            image = shape.image
            # 规范化扩展名
            ext = image.ext.lower()
            if ext not in ['png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff']:
                ext = 'png'  # 默认兜底

            # 生成唯一文件名
            filename = f"slide_{slide_number}_shape_{shape.shape_id}.{ext}"
            filepath = self.image_output_path / filename

            # 写入二进制数据
            filepath.write_bytes(image.blob)

            logger.debug(f"图片已保存: {filepath}")
            return filepath
        except Exception as e:
            logger.warning(f"图片保存失败 (Slide {slide_number}, Shape {shape.shape_id}): {str(e)}")
            return None

    def _parse_table(self, table) -> List[List[str]]:
        """解析表格数据"""
        return [
            [cell.text.strip() for cell in row.cells]
            for row in table.rows
        ]

    def _extract_metadata(self, prs) -> Dict[str, Any]:
        """提取元数据"""
        props = prs.core_properties
        return {
            "author": props.author or "",
            "title": props.title or "",
            "subject": props.subject or "",
            "modified": str(props.modified) if props.modified else ""
        }

    # 兼容旧代码：保留 __del__ 作为最后的防线，但不再推荐依赖
    def __del__(self):
        self.cleanup()
