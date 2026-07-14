"""Document probing — low-cost pre-parse inspection.

DocumentProbe performs deterministic, low-cost checks on a source artifact
before any parser is invoked: format detection, encryption, corruption,
image-only pages, and basic structural hints.
"""

from __future__ import annotations

import hashlib
import struct
from dataclasses import dataclass, field
from enum import Enum
from typing import List, Optional, Tuple


# ---------------------------------------------------------------------------
# ProbeResult
# ---------------------------------------------------------------------------


class DetectedFormat(str, Enum):
    """Detected document format after probing."""

    PPTX = "pptx"
    PDF = "pdf"
    DOCX = "docx"
    IMAGE = "image"          # standalone image file
    UNSUPPORTED = "unsupported"
    CORRUPT = "corrupt"
    ENCRYPTED = "encrypted"


class EncryptionStatus(str, Enum):
    NONE = "none"
    PASSWORD_PROTECTED = "password_protected"
    CERTIFICATE_ENCRYPTED = "certificate_encrypted"
    UNKNOWN = "unknown"


@dataclass(frozen=True)
class ProbeResult:
    """Result of a low-cost document probe.

    This is the primary output of DocumentProbe.  All fields are determined
    through cheap, deterministic checks — no full parse, no OCR, no AI.
    """

    detected_format: DetectedFormat
    encryption: EncryptionStatus = EncryptionStatus.NONE
    page_or_slide_count: int = 0
    has_text_content: bool = False
    image_only_pages: Tuple[int, ...] = field(default_factory=tuple)
    estimated_text_coverage: float = 0.0
    estimated_image_ratio: float = 0.0
    has_tables: bool = False
    has_formulas: bool = False
    has_notes: bool = False
    file_size_bytes: int = 0
    probe_duration_ms: float = 0.0
    error: Optional[str] = None
    warnings: List[str] = field(default_factory=list)

    def is_parseable(self) -> bool:
        """Return True if the artifact can likely be parsed by some provider."""
        if self.detected_format in (DetectedFormat.UNSUPPORTED,
                                    DetectedFormat.CORRUPT,
                                    DetectedFormat.ENCRYPTED):
            return False
        return True

    def needs_ocr(self) -> bool:
        """Return True if OCR enrichment is likely needed."""
        if self.estimated_text_coverage < 0.3 and self.image_only_pages:
            return True
        return False


# ---------------------------------------------------------------------------
# DocumentProbe
# ---------------------------------------------------------------------------


class DocumentProbe:
    """Low-cost document probing engine.

    Uses file signatures (magic bytes), lightweight library calls, and basic
    structural analysis — never a full parse, never OCR, never AI.
    """

    # Magic bytes for supported formats
    _MAGIC: dict = {
        b"\x50\x4b\x03\x04": "zip",         # ZIP-based (PPTX, DOCX, XLSX)
        b"\x25\x50\x44\x46": "pdf",          # PDF
    }

    _PPTX_MIME_TYPES = (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    )
    _DOCX_MIME_TYPES = (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    )
    _IMAGE_MIME_TYPES = (
        "image/png", "image/jpeg", "image/tiff", "image/bmp", "image/gif",
        "image/webp",
    )

    def __init__(self) -> None:
        self._has_pptx = False
        self._has_pymupdf = False
        self._has_docx_lib = False
        self._check_dependencies()

    def _check_dependencies(self) -> None:
        """Check which probe libraries are available."""
        try:
            import pptx  # noqa: F401
            self._has_pptx = True
        except ImportError:
            pass
        try:
            import fitz  # noqa: F401  (PyMuPDF)
            self._has_pymupdf = True
        except ImportError:
            pass
        try:
            import docx  # noqa: F401
            self._has_docx_lib = True
        except ImportError:
            pass

    def probe(self, data: bytes, filename: str = "",
              mime: str = "") -> ProbeResult:
        """Run a low-cost probe on the given source bytes."""
        import time
        start = time.perf_counter()

        fmt, enc = self._detect_format(data, filename, mime)
        if fmt == DetectedFormat.CORRUPT:
            duration = (time.perf_counter() - start) * 1000
            return ProbeResult(
                detected_format=fmt,
                file_size_bytes=len(data),
                probe_duration_ms=duration,
                error="File appears corrupted or empty",
            )
        if fmt == DetectedFormat.ENCRYPTED:
            duration = (time.perf_counter() - start) * 1000
            return ProbeResult(
                detected_format=fmt,
                encryption=EncryptionStatus.PASSWORD_PROTECTED,
                file_size_bytes=len(data),
                probe_duration_ms=duration,
                error="File is password-protected or encrypted",
            )

        # Run format-specific probe
        result: Optional[ProbeResult] = None
        if fmt == DetectedFormat.PPTX and self._has_pptx:
            result = self._probe_pptx(data)
        elif fmt == DetectedFormat.PDF and self._has_pymupdf:
            result = self._probe_pdf(data)
        elif fmt == DetectedFormat.DOCX and self._has_docx_lib:
            result = self._probe_docx(data)
        elif fmt == DetectedFormat.IMAGE:
            result = self._probe_image(data)

        duration = (time.perf_counter() - start) * 1000
        if result is None:
            # Fallback: basic stats only
            return ProbeResult(
                detected_format=fmt,
                file_size_bytes=len(data),
                probe_duration_ms=duration,
                warnings=[f"No detailed probe available for {fmt.value}"],
            )

        object.__setattr__(result, "probe_duration_ms", duration)
        return result

    def _detect_format(
        self, data: bytes, filename: str, mime: str,
    ) -> Tuple[DetectedFormat, EncryptionStatus]:
        """Detect format from magic bytes, mime, and filename."""
        if len(data) < 4:
            return DetectedFormat.CORRUPT, EncryptionStatus.NONE

        # Check for ZIP-based formats (PPTX, DOCX)
        if data[:4] == b"\x50\x4b\x03\x04":
            # Need enough data for ZIP header fields
            if len(data) < 30:
                return DetectedFormat.CORRUPT, EncryptionStatus.NONE
            # Check encryption flag in ZIP header (bit 0 of general purpose flag)
            # byte 6-7: general purpose bit flag; bit 0 = encrypted
            if len(data) > 8:
                gp_flag = struct.unpack("<H", data[6:8])[0]
                if gp_flag & 0x01:
                    return DetectedFormat.ENCRYPTED, EncryptionStatus.PASSWORD_PROTECTED
            # Distinguish PPTX vs DOCX by filename/mime
            lower = filename.lower()
            if ".pptx" in lower or "presentation" in mime:
                return DetectedFormat.PPTX, EncryptionStatus.NONE
            if ".docx" in lower or "wordprocessing" in mime or "msword" in mime:
                return DetectedFormat.DOCX, EncryptionStatus.NONE
            # Check ZIP contents for [Content_Types].xml
            if self._check_zip_content_type(data, "presentation"):
                return DetectedFormat.PPTX, EncryptionStatus.NONE
            if self._check_zip_content_type(data, "word"):
                return DetectedFormat.DOCX, EncryptionStatus.NONE
            # Default to PPTX if no better match
            return DetectedFormat.PPTX, EncryptionStatus.NONE

        # PDF
        if data[:4] == b"\x25\x50\x44\x46":
            # Check for encryption in PDF trailer
            lower = data.lower()
            if b"/encrypt" in lower and b"/filter" in lower:
                return DetectedFormat.ENCRYPTED, EncryptionStatus.PASSWORD_PROTECTED
            return DetectedFormat.PDF, EncryptionStatus.NONE

        # Image formats by magic bytes
        if data[:8] == b"\x89\x50\x4e\x47\x0d\x0a\x1a\x0a":  # PNG
            return DetectedFormat.IMAGE, EncryptionStatus.NONE
        if data[:2] in (b"\xff\xd8",):
            return DetectedFormat.IMAGE, EncryptionStatus.NONE
        if data[:4] == b"\x49\x49\x2a\x00" or data[:4] == b"\x4d\x4d\x00\x2a":  # TIFF
            return DetectedFormat.IMAGE, EncryptionStatus.NONE
        if data[:2] == b"\x42\x4d":  # BMP
            return DetectedFormat.IMAGE, EncryptionStatus.NONE

        return DetectedFormat.UNSUPPORTED, EncryptionStatus.NONE

    @staticmethod
    def _check_zip_content_type(data: bytes, keyword: str) -> bool:
        """Check if ZIP contents include a [Content_Types].xml reference."""
        try:
            import zipfile
            import io
            with zipfile.ZipFile(io.BytesIO(data)) as zf:
                if "[Content_Types].xml" in zf.namelist():
                    ct = zf.read("[Content_Types].xml").lower()
                    return keyword.encode() in ct
        except Exception:
            pass
        return False

    def _probe_pptx(self, data: bytes) -> ProbeResult:
        """Probe a PPTX file for structural hints."""
        from pptx import Presentation
        import io

        try:
            prs = Presentation(io.BytesIO(data))
        except Exception:
            # Detailed probe failed, but format was correctly identified
            # Return a basic result with the detected format preserved
            return ProbeResult(
                detected_format=DetectedFormat.PPTX,
                page_or_slide_count=0,
                file_size_bytes=len(data),
                warnings=["python-pptx could not open file for detailed probe; "
                          "format detected by signature only"],
            )

        slide_count = len(prs.slides)
        image_only_pages: List[int] = []
        has_tables = False
        has_notes = False
        total_shapes = 0
        text_shapes = 0
        image_shapes = 0

        for i, slide in enumerate(prs.slides, 1):
            slide_has_text = False
            slide_has_image = False
            for shape in slide.shapes:
                total_shapes += 1
                if shape.has_table:
                    has_tables = True
                    slide_has_text = True
                if hasattr(shape, "text") and shape.text.strip():
                    text_shapes += 1
                    slide_has_text = True
                try:
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        image_shapes += 1
                        slide_has_image = True
                except Exception:
                    pass
            if not slide_has_text and slide_has_image:
                image_only_pages.append(i)
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame and notes_slide.notes_text_frame.text.strip():
                    has_notes = True

        text_ratio = text_shapes / max(total_shapes, 1)
        image_ratio = image_shapes / max(total_shapes, 1)
        has_text = text_shapes > 0

        return ProbeResult(
            detected_format=DetectedFormat.PPTX,
            page_or_slide_count=slide_count,
            has_text_content=has_text,
            image_only_pages=tuple(image_only_pages),
            estimated_text_coverage=text_ratio,
            estimated_image_ratio=image_ratio,
            has_tables=has_tables,
            has_notes=has_notes,
            file_size_bytes=len(data),
        )

    def _probe_pdf(self, data: bytes) -> ProbeResult:
        """Probe a PDF file for page count and text hints."""
        import fitz
        import io

        try:
            doc = fitz.open(stream=data, filetype="pdf")
        except Exception:
            return ProbeResult(
                detected_format=DetectedFormat.PDF,
                page_or_slide_count=0,
                file_size_bytes=len(data),
                warnings=["PyMuPDF could not open file for detailed probe; "
                          "format detected by signature only"],
            )

        page_count = doc.page_count
        text_pages = 0
        image_only_pages: List[int] = []
        total_chars = 0

        for i in range(page_count):
            page = doc[i]
            text = page.get_text().strip()
            chars = len(text)
            total_chars += chars
            if chars > 20:
                text_pages += 1
            else:
                # Check for images on this page
                images = page.get_images(full=True)
                if images:
                    image_only_pages.append(i + 1)

        doc.close()

        text_coverage = text_pages / max(page_count, 1)
        return ProbeResult(
            detected_format=DetectedFormat.PDF,
            page_or_slide_count=page_count,
            has_text_content=text_pages > 0,
            image_only_pages=tuple(image_only_pages),
            estimated_text_coverage=text_coverage,
            file_size_bytes=len(data),
        )

    def _probe_docx(self, data: bytes) -> ProbeResult:
        """Probe a DOCX file."""
        import docx
        import io

        try:
            doc = docx.Document(io.BytesIO(data))
        except Exception:
            return ProbeResult(
                detected_format=DetectedFormat.DOCX,
                page_or_slide_count=0,
                file_size_bytes=len(data),
                warnings=["python-docx could not open file for detailed probe; "
                          "format detected by signature only"],
            )

        paras = [p.text for p in doc.paragraphs if p.text.strip()]
        has_text = len(paras) > 0
        has_tables = len(doc.tables) > 0

        return ProbeResult(
            detected_format=DetectedFormat.DOCX,
            page_or_slide_count=len(paras) // 30 + 1,  # rough estimate
            has_text_content=has_text,
            estimated_text_coverage=1.0 if has_text else 0.0,
            has_tables=has_tables,
            file_size_bytes=len(data),
        )

    def _probe_image(self, data: bytes) -> ProbeResult:
        """Probe a standalone image file."""
        return ProbeResult(
            detected_format=DetectedFormat.IMAGE,
            page_or_slide_count=1,
            has_text_content=False,
            image_only_pages=(1,),
            estimated_text_coverage=0.0,
            estimated_image_ratio=1.0,
            file_size_bytes=len(data),
            warnings=["Image file detected; OCR provider required"],
        )
