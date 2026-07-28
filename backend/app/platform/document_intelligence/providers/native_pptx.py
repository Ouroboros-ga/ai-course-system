"""Native PPTX parser provider.

Wraps python-pptx to extract slide structure, text, tables, notes, and images,
producing a ParserOutput compatible with the DocumentIR mapper.

This provider uses the existing python-pptx library (same dependency as the
read-only ``common/ppt_parser.py::PPTParser``) but produces structured output
for the DocumentIR pipeline rather than the old V1 container types.
"""

from __future__ import annotations

import io
import posixpath
import re
import time
import zipfile
from dataclasses import dataclass, field
from typing import Any, Dict, List, Optional, Tuple
from xml.etree import ElementTree

from ..registry import (
    ParserCapabilities,
    ParserOutput,
    ParseError,
    ParseTimeoutError,
    ParseUnavailableError,
    ParseMalformedError,
)
from ..planner import ParsePlan
from ..source_artifact import SourceArtifact


# ---------------------------------------------------------------------------
# Check python-pptx availability
# ---------------------------------------------------------------------------

try:
    from pptx import Presentation as _PptxPresentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE as _MSO_SHAPE_TYPE

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    _PptxPresentation = None  # type: ignore
    _MSO_SHAPE_TYPE = None    # type: ignore


_PPT_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "dgm": "http://schemas.openxmlformats.org/drawingml/2006/diagram",
}


def _extract_ooxml_supplements(
    data: bytes, slide_count: int,
) -> tuple[dict[int, list[dict[str, str]]], list[str]]:
    """Conservatively expose OOXML text python-pptx does not model.

    The relationship graph is intentionally not reconstructed here: a source
    XML locator is still a precise citation target, while unlinked diagrams or
    charts remain explicitly marked for review rather than guessed onto a
    slide.  This parses only text-bearing XML and never manufactures geometry.
    """
    pages: dict[int, list[dict[str, str]]] = {index: [] for index in range(1, slide_count + 1)}
    warnings: list[str] = []
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            names = set(archive.namelist())
            for slide_index in range(1, slide_count + 1):
                name = f"ppt/slides/slide{slide_index}.xml"
                if name not in names:
                    continue
                root = ElementTree.fromstring(archive.read(name))
                for index, shape in enumerate(root.findall(".//p:pic", _PPT_NS), start=1):
                    alt = _text(shape.find("p:nvPicPr/p:cNvPr", _PPT_NS), "descr")
                    title = _text(shape.find("p:nvPicPr/p:cNvPr", _PPT_NS), "title")
                    if alt or title:
                        pages[slide_index].append({
                            "kind": "image_alt_text",
                            "text": alt or title,
                            "raw_locator": f"{name}#/p:sld/p:cSld/p:spTree/p:pic[{index}]/p:nvPicPr/p:cNvPr",
                        })
                if root.findall(".//a:graphicData[@uri='http://schemas.openxmlformats.org/drawingml/2006/diagram']", _PPT_NS):
                    warnings.append(f"PPTX_SMARTART_REVIEW_REQUIRED:{name}")
                if root.findall(".//a:graphicData[@uri='http://schemas.openxmlformats.org/drawingml/2006/chart']", _PPT_NS):
                    warnings.append(f"PPTX_CHART_REVIEW_REQUIRED:{name}")
                if root.findall(".//m:oMath", {**_PPT_NS, "m": "http://schemas.openxmlformats.org/officeDocument/2006/math"}):
                    warnings.append(f"PPTX_FORMULA_REVIEW_REQUIRED:{name}")

            # Diagram and chart XML contains text absent from python-pptx. If
            # a relationship cannot be proven cheaply, keep it unassigned so
            # the quality signal tells the reviewer exactly what is missing.
            for prefix, kind, xpath in (
                ("ppt/diagrams/data", "diagram_text", ".//dgm:t"),
                ("ppt/charts/chart", "chart_text", ".//a:t"),
            ):
                for name in sorted(item for item in names if item.startswith(prefix) and item.endswith(".xml")):
                    root = ElementTree.fromstring(archive.read(name))
                    text = " ".join(item.text.strip() for item in root.findall(xpath, _PPT_NS) if item.text and item.text.strip())
                    if text:
                        warnings.append(f"PPTX_UNASSIGNED_{kind.upper()}:{name}")
    except (OSError, ValueError, zipfile.BadZipFile, ElementTree.ParseError) as exc:
        warnings.append(f"PPTX_OOXML_SUPPLEMENT_UNAVAILABLE:{type(exc).__name__}")
    return pages, warnings


def _text(element: Any, attribute: str) -> str:
    return str(element.get(attribute, "")).strip() if element is not None else ""


# ---------------------------------------------------------------------------
# Slide data containers
# ---------------------------------------------------------------------------


@dataclass(frozen=True)
class SlideShape:
    """A shape extracted from a PPTX slide."""

    shape_id: int
    shape_name: str
    shape_type: int
    text: Optional[str] = None
    bbox_emu: Optional[Tuple[int, int, int, int]] = None  # x, y, w, h in EMU
    has_table: bool = False
    is_title: bool = False
    is_picture: bool = False
    is_group: bool = False
    raw_data: Optional[Dict[str, Any]] = None


@dataclass(frozen=True)
class SlideTable:
    """A table extracted from a PPTX slide."""

    rows: int
    columns: int
    cells: Tuple[Tuple[str, ...], ...] = field(default_factory=tuple)
    shape_id: int = 0
    bbox_emu: Optional[Tuple[int, int, int, int]] = None
    cell_bboxes_emu: Tuple[Tuple[Optional[Tuple[int, int, int, int]], ...], ...] = field(default_factory=tuple)


@dataclass(frozen=True)
class SlideData:
    """All data extracted from a single PPTX slide."""

    slide_index: int
    title: Optional[str] = None
    shapes: Tuple[SlideShape, ...] = field(default_factory=tuple)
    tables: Tuple[SlideTable, ...] = field(default_factory=tuple)
    notes: Optional[str] = None
    width_emu: int = 0
    height_emu: int = 0
    slide_layout: Optional[str] = None


# ---------------------------------------------------------------------------
# Parser provider
# ---------------------------------------------------------------------------


class NativePptxProvider:
    """Parser provider for native PPTX files using python-pptx.

    This provider extracts slide-level structure, text content, tables,
    notes, and shape metadata.  It does NOT perform OCR, formula recognition,
    or image analysis — those are enrichment steps.
    """

    name = "native-pptx"
    version = "1.0.0"
    capabilities = ParserCapabilities(
        supported_formats=("pptx",),
        supports_tables=True,
        supports_notes=True,
        supports_reading_order=True,
        supports_coordinates=True,
        supports_visual_assets=True,
        supports_provenance=True,
        max_file_size_bytes=200 * 1024 * 1024,
        max_pages=200,
    )

    def __init__(self) -> None:
        if not HAS_PPTX:
            self._available = False
        else:
            self._available = True

    async def parse(
        self,
        source: SourceArtifact,
        plan: ParsePlan,
    ) -> ParserOutput:
        """Parse a PPTX source artifact.

        Args:
            source: The source artifact (must be PPTX format).
            plan: The parse plan with configuration.

        Returns:
            ParserOutput with slide data.

        Raises:
            ParseUnavailableError: If python-pptx is not installed.
            ParseMalformedError: If the file cannot be opened as a PPTX.
            ParseTimeoutError: If parsing exceeds the configured timeout.
        """
        if not self._available:
            raise ParseUnavailableError(
                "python-pptx is not installed. "
                "Install with: pip install python-pptx"
            )

        timeout_ms = 120000
        if plan.steps:
            primary = plan.primary_step
            if primary:
                timeout_ms = primary.timeout_ms

        start = time.perf_counter()

        # Source data must be accessible
        data = self._get_source_data(source)
        if data is None:
            raise ParseMalformedError(
                "Source artifact has no accessible data"
            )

        try:
            prs = _PptxPresentation(io.BytesIO(data))
        except Exception as exc:
            raise ParseMalformedError(
                f"Failed to open PPTX: {exc}"
            ) from exc

        slides: List[SlideData] = []
        slide_width = prs.slide_width or 0
        slide_height = prs.slide_height or 0

        for i, pptx_slide in enumerate(prs.slides):
            if (time.perf_counter() - start) * 1000 > timeout_ms:
                raise ParseTimeoutError(
                    f"PPTX parsing timed out after {timeout_ms}ms "
                    f"at slide {i + 1}/{len(prs.slides)}"
                )

            slide_data = self._extract_slide(pptx_slide, i + 1,
                                             slide_width, slide_height)
            slides.append(slide_data)

        supplements, supplement_warnings = _extract_ooxml_supplements(data, len(slides))

        duration_ms = int((time.perf_counter() - start) * 1000)

        # Build pages tuple for ParserOutput
        pages: List[Dict[str, Any]] = []
        for slide in slides:
            page = self._slide_to_page_dict(slide)
            page["ooxml_supplements"] = supplements.get(slide.slide_index, [])
            pages.append(page)

        return ParserOutput(
            provider=self.name,
            provider_version=self.version,
            pages=tuple(pages),
            metadata={
                "slide_count": len(slides),
                "slide_width_emu": slide_width,
                "slide_height_emu": slide_height,
                "duration_ms": duration_ms,
            },
            warnings=supplement_warnings,
        )

    def _extract_slide(
        self,
        pptx_slide: Any,
        index: int,
        slide_width: int,
        slide_height: int,
    ) -> SlideData:
        """Extract all data from a single PPTX slide."""
        title: Optional[str] = None
        shapes: List[SlideShape] = []
        tables: List[SlideTable] = []
        notes: Optional[str] = None

        # Get slide title
        if pptx_slide.shapes.title:
            title_text = pptx_slide.shapes.title.text
            if title_text:
                title = title_text.strip()

        # Extract nested GroupShape members too; a group is not a leaf node.
        def extract_shape(shape: Any, parent_shape_id: Optional[int] = None) -> None:
            shape_type = int(shape.shape_type) if hasattr(shape, "shape_type") else 0
            is_title = (shape == pptx_slide.shapes.title) if pptx_slide.shapes.title else False

            # Extract text
            shape_text: Optional[str] = None
            if hasattr(shape, "text") and shape.text:
                shape_text = shape.text.strip()

            # Extract bounding box (EMU coordinates)
            bbox = None
            if hasattr(shape, "left") and hasattr(shape, "top"):
                left = int(getattr(shape, "left", 0))
                top = int(getattr(shape, "top", 0))
                width = int(getattr(shape, "width", 0))
                height = int(getattr(shape, "height", 0))
                bbox = (left, top, width, height)

            # Check for table
            has_table = False
            if hasattr(shape, "has_table") and shape.has_table:
                has_table = True
                table = shape.table
                rows_data: List[Tuple[str, ...]] = []
                cell_boxes: List[Tuple[Optional[Tuple[int, int, int, int]], ...]] = []
                top_offset = 0
                for row in table.rows:
                    cells = tuple(cell.text.strip() for cell in row.cells)
                    rows_data.append(cells)
                    left_offset = 0
                    row_boxes: List[Optional[Tuple[int, int, int, int]]] = []
                    for column in table.columns:
                        if bbox is None:
                            row_boxes.append(None)
                        else:
                            row_boxes.append((
                                bbox[0] + left_offset,
                                bbox[1] + top_offset,
                                int(column.width),
                                int(row.height),
                            ))
                        left_offset += int(column.width)
                    cell_boxes.append(tuple(row_boxes))
                    top_offset += int(row.height)
                tables.append(SlideTable(
                    rows=len(table.rows),
                    columns=len(table.columns),
                    cells=tuple(rows_data),
                    shape_id=int(getattr(shape, "shape_id", 0) or 0),
                    bbox_emu=bbox,
                    cell_bboxes_emu=tuple(cell_boxes),
                ))

            # Check for picture
            is_picture = False
            try:
                is_picture = (shape_type == 13)  # MSO_SHAPE_TYPE.PICTURE
            except Exception:
                pass

            is_group = (shape_type == 6)  # MSO_SHAPE_TYPE.GROUP

            # Group containers provide hierarchy, not independently citeable
            # content. Table shapes are represented by the structured table
            # block below, avoiding an additional empty/unknown block.
            if not is_group and not has_table:
                shapes.append(SlideShape(
                    shape_id=int(shape.shape_id) if hasattr(shape, "shape_id") else 0,
                    shape_name=str(getattr(shape, "name", "")),
                    shape_type=shape_type,
                    text=shape_text,
                    bbox_emu=bbox,
                    has_table=has_table,
                    is_title=is_title,
                    is_picture=is_picture,
                    is_group=is_group,
                    raw_data={"parent_shape_id": parent_shape_id} if parent_shape_id is not None else None,
                ))
            if is_group and hasattr(shape, "shapes"):
                for child in shape.shapes:
                    extract_shape(child, int(getattr(shape, "shape_id", 0) or 0))

        for shape in pptx_slide.shapes:
            extract_shape(shape)

        # Extract notes
        if pptx_slide.has_notes_slide:
            notes_slide = pptx_slide.notes_slide
            if notes_slide.notes_text_frame:
                notes_text = notes_slide.notes_text_frame.text
                if notes_text:
                    notes = notes_text.strip()

        return SlideData(
            slide_index=index,
            title=title,
            shapes=tuple(shapes),
            tables=tuple(tables),
            notes=notes,
            width_emu=slide_width,
            height_emu=slide_height,
            slide_layout=str(getattr(pptx_slide, "slide_layout", None)),
        )

    def _slide_to_page_dict(self, slide: SlideData) -> Dict[str, Any]:
        """Convert a SlideData to a dict for ParserOutput pages."""
        shape_dicts: List[Dict[str, Any]] = []
        for shape in slide.shapes:
            sd: Dict[str, Any] = {
                "shape_id": shape.shape_id,
                "shape_name": shape.shape_name,
                "shape_type": shape.shape_type,
                "is_title": shape.is_title,
                "is_picture": shape.is_picture,
                "is_group": shape.is_group,
            }
            if shape.text is not None:
                sd["text"] = shape.text
            if shape.bbox_emu is not None:
                sd["bbox_emu"] = list(shape.bbox_emu)
            if shape.has_table:
                sd["has_table"] = True
            if shape.raw_data:
                sd.update(shape.raw_data)
            shape_dicts.append(sd)

        table_dicts: List[Dict[str, Any]] = []
        for table in slide.tables:
            td: Dict[str, Any] = {
                "rows": table.rows,
                "columns": table.columns,
                "cells": [list(row) for row in table.cells],
                "shape_id": table.shape_id,
            }
            if table.bbox_emu is not None:
                td["bbox_emu"] = list(table.bbox_emu)
            if table.cell_bboxes_emu:
                td["cell_bboxes_emu"] = [
                    [list(cell) if cell is not None else None for cell in row]
                    for row in table.cell_bboxes_emu
                ]
            table_dicts.append(td)

        result: Dict[str, Any] = {
            "slide_index": slide.slide_index,
            "title": slide.title or "",
            "shapes": shape_dicts,
            "tables": table_dicts,
            "width_emu": slide.width_emu,
            "height_emu": slide.height_emu,
        }
        if slide.notes:
            result["notes"] = slide.notes
        if slide.slide_layout:
            result["slide_layout"] = slide.slide_layout

        return result

    @staticmethod
    def _get_source_data(source: SourceArtifact) -> Optional[bytes]:
        """从 SourceArtifact 读取源数据。

        P0-4：通过 source.uri（即 object_key）从对象存储读取文件内容。
        若 uri 为空，回退到 None（测试场景可由调用方直接注入数据）。
        """
        if source.data is not None:
            return source.data
        if not source.uri:
            return None
        try:
            from app.services.object_storage import get_object_storage
            storage = get_object_storage()
            return storage.get(source.uri)
        except FileNotFoundError:
            return None
        except Exception:
            return None


# ---------------------------------------------------------------------------
# IR mapper
# ---------------------------------------------------------------------------


def _reading_order_shapes(shapes: List[Dict[str, Any]], slide_height: int) -> List[Dict[str, Any]]:
    """Resolve visual reading order instead of leaking PPT z-order.

    Placeholder titles lead. Remaining text is read by approximate columns,
    top-to-bottom within a column; small elements at the bottom are deferred
    as footer-like material.  Group children retain their own coordinates.
    """
    def geometry(shape: Dict[str, Any]) -> tuple[float, float, float, float]:
        left, top, width, height = shape.get("bbox_emu") or (0, 0, 0, 0)
        return float(left), float(top), float(width), float(height)

    title = [shape for shape in shapes if shape.get("is_title")]
    body = [shape for shape in shapes if not shape.get("is_title")]
    threshold = float(slide_height or 1) * 0.86
    footer = [shape for shape in body if geometry(shape)[1] >= threshold and geometry(shape)[3] < float(slide_height or 1) * 0.12]
    content = [shape for shape in body if shape not in footer]
    content.sort(key=lambda shape: geometry(shape)[0] + geometry(shape)[2] / 2.0)
    columns: List[List[Dict[str, Any]]] = []
    for shape in content:
        x, _y, width, _height = geometry(shape)
        center = x + width / 2.0
        target = next((column for column in columns if abs((geometry(column[0])[0] + geometry(column[0])[2] / 2.0) - center) < max(width, geometry(column[0])[2], 1.0) * 0.75), None)
        if target is None:
            columns.append([shape])
        else:
            target.append(shape)
    columns.sort(key=lambda column: geometry(column[0])[0])
    ordered = list(title)
    for column in columns:
        ordered.extend(sorted(column, key=lambda shape: geometry(shape)[1]))
    ordered.extend(sorted(footer, key=lambda shape: geometry(shape)[0]))
    return ordered


def map_pptx_output_to_ir(
    output: ParserOutput,
    source: SourceArtifact,
    run_id: str,
    parser_run_id: str,
    normalization_version: str = "1",
) -> Tuple[List[Dict[str, Any]], List[Dict[str, Any]], List[Dict[str, Any]]]:
    """Map native-pptx ParserOutput to DocumentIR block/unit/asset dicts.

    Args:
        output: The ParserOutput from NativePptxProvider.
        source: The source artifact.
        run_id: The pipeline run ID.
        parser_run_id: The parser run ID.
        normalization_version: Normalization version for stable IDs.

    Returns:
        Tuple of (blocks, units, assets) as dicts ready for DocumentIR construction.
    """
    from ..contracts import CURRENT_SCHEMA_VERSION
    from ..document_ir.models import (
        compute_document_id, compute_unit_id,
        ContentBlock, DocumentUnit, Provenance, UnitType,
        AssetKind, VisualAsset,
    )

    doc_id = compute_document_id(
        artifact_id=source.artifact_id,
        schema_version=CURRENT_SCHEMA_VERSION.serialize(),
        normalization_version=normalization_version,
    )

    blocks: List[Dict[str, Any]] = []
    units: List[Dict[str, Any]] = []
    assets: List[Dict[str, Any]] = []
    page_blocks: Dict[int, List[str]] = {}

    slide_width = output.metadata.get("slide_width_emu", 0)
    slide_height = output.metadata.get("slide_height_emu", 0)

    for page_dict in output.pages:
        slide_index = page_dict["slide_index"]
        unit_id = compute_unit_id(
            document_id=doc_id,
            unit_type="slide",
            index=slide_index,
            normalization_version=normalization_version,
        )
        block_ids: List[str] = []

        shapes = _reading_order_shapes(page_dict.get("shapes", []), slide_height)
        for shape in shapes:
            block_id = f"blk_pptx_s{slide_index}_sh{shape['shape_id']}"

            # Build bounding box
            bbox = None
            bbox_emu = shape.get("bbox_emu")
            if bbox_emu and slide_width > 0 and slide_height > 0:
                from ..contracts import BoundingBox, CoordinateSpace
                left, top, w, h = bbox_emu
                x0 = left / slide_width
                y0 = top / slide_height
                x1 = (left + w) / slide_width
                y1 = (top + h) / slide_height
                try:
                    bbox = BoundingBox(
                        x0=round(x0, 6),
                        y0=round(y0, 6),
                        x1=round(x1, 6),
                        y1=round(y1, 6),
                        coordinate_space=CoordinateSpace.NORMALIZED,
                    )
                except ValueError:
                    bbox = None

            # Determine block type
            text = shape.get("text", "")
            block_type = "paragraph"
            heading_level = None
            if shape.get("is_title") or _is_numbered_section_heading(text):
                block_type = "heading"
                heading_level = 1
            elif not text:
                block_type = "image" if shape.get("is_picture") else "unknown"

            provenance = Provenance(
                artifact_id=source.artifact_id,
                run_id=run_id,
                parser_run_id=parser_run_id,
                provider="native-pptx",
                raw_locator=f"slides/{slide_index}/shapes/{shape['shape_id']}",
                page_or_slide=slide_index,
                bbox=bbox,
                confidence=1.0,
            )

            block = ContentBlock(
                block_id=block_id,
                page_or_slide=slide_index,
                bbox=bbox,
                reading_order=len(block_ids) + 1,
                block_type=block_type,
                heading_level=heading_level,
                text=text or None,
                provider="native-pptx",
                provenance=(provenance,),
                raw_result_ref=f"artifact://{parser_run_id}/raw.json"
                              f"#/slides/{slide_index}/shapes/{shape['shape_id']}",
            )
            blocks.append(block.to_dict())
            block_ids.append(block_id)

        # OOXML supplements intentionally lack fabricated geometry.  Their
        # XML locator is the authoritative native citation point until a
        # rendition can align the visual element to a slide bounding box.
        for supplement_index, supplement in enumerate(page_dict.get("ooxml_supplements", [])):
            text = supplement.get("text", "").strip()
            if not text:
                continue
            block_id = f"blk_pptx_s{slide_index}_xml{supplement_index}"
            kind = supplement.get("kind", "unknown")
            block_type = "diagram" if kind == "diagram_text" else (
                "chart" if kind == "chart_text" else "image"
            )
            provenance = Provenance(
                artifact_id=source.artifact_id,
                run_id=run_id,
                parser_run_id=parser_run_id,
                provider="native-pptx-ooxml",
                raw_locator=supplement["raw_locator"],
                page_or_slide=slide_index,
                confidence=1.0,
            )
            blocks.append(ContentBlock(
                block_id=block_id,
                page_or_slide=slide_index,
                reading_order=len(block_ids) + 1,
                block_type=block_type,
                text=text,
                provider="native-pptx-ooxml",
                provenance=(provenance,),
                raw_result_ref=f"artifact://{parser_run_id}/raw.json#{supplement['raw_locator']}",
            ).to_dict())
            block_ids.append(block_id)

        # Handle tables
        for t_idx, table in enumerate(page_dict.get("tables", [])):
            block_id = f"blk_pptx_s{slide_index}_sh{table.get('shape_id', t_idx)}"
            from ..document_ir.models import TableBlock, TableCell

            def normalized_bbox(raw_bbox: Any):
                if not raw_bbox or not slide_width or not slide_height:
                    return None
                from ..contracts import BoundingBox, CoordinateSpace
                left, top, width, height = raw_bbox
                try:
                    return BoundingBox(
                        x0=round(left / slide_width, 6),
                        y0=round(top / slide_height, 6),
                        x1=round((left + width) / slide_width, 6),
                        y1=round((top + height) / slide_height, 6),
                        coordinate_space=CoordinateSpace.NORMALIZED,
                    )
                except ValueError:
                    return None

            table_bbox = normalized_bbox(table.get("bbox_emu"))
            cells: List[TableCell] = []
            for r_idx, row in enumerate(table.get("cells", [])):
                for c_idx, cell_text in enumerate(row):
                    cells.append(TableCell(
                        row=r_idx,
                        col=c_idx,
                        text=cell_text,
                        bbox=normalized_bbox(
                            (table.get("cell_bboxes_emu") or [[]])[r_idx][c_idx]
                            if r_idx < len(table.get("cell_bboxes_emu") or [])
                            and c_idx < len((table.get("cell_bboxes_emu") or [[]])[r_idx])
                            else None
                        ),
                        header=r_idx == 0,
                    ))

            table_provenance = Provenance(
                artifact_id=source.artifact_id,
                run_id=run_id,
                parser_run_id=parser_run_id,
                provider="native-pptx",
                raw_locator=f"slides/{slide_index}/shapes/{table.get('shape_id', t_idx)}",
                page_or_slide=slide_index,
                bbox=table_bbox,
            )

            table_block = TableBlock(
                block_id=block_id,
                page_or_slide=slide_index,
                bbox=table_bbox,
                reading_order=len(block_ids) + 1,
                rows=table.get("rows", 0),
                columns=table.get("columns", 0),
                cells=tuple(cells),
                provider="native-pptx",
                provenance=(table_provenance,),
            )
            blocks.append(table_block.to_dict())
            block_ids.append(block_id)

        # Create unit
        unit = DocumentUnit(
            unit_id=unit_id,
            unit_type=UnitType.SLIDE,
            index=slide_index,
            label=page_dict.get("title") or None,
            width=slide_width / 914400 if slide_width else None,  # EMU to inches
            height=slide_height / 914400 if slide_height else None,
            coordinate_unit="inch",
            block_ids=tuple(block_ids),
            provenance=(
                Provenance(
                    artifact_id=source.artifact_id,
                    run_id=run_id,
                    parser_run_id=parser_run_id,
                    provider="native-pptx",
                    raw_locator=f"slides/{slide_index}",
                    page_or_slide=slide_index,
                ),
            ),
        )
        units.append(unit.to_dict())
        page_blocks[slide_index] = block_ids

        # Add slide render asset placeholder
        if slide_width and slide_height:
            asset = VisualAsset(
                asset_id=f"asset_pptx_slide_{slide_index}",
                kind=AssetKind.SLIDE_RENDER,
                page_or_slide=slide_index,
                width=slide_width / 914400,
                height=slide_height / 914400,
                linked_block_ids=tuple(block_ids),
                provider="native-pptx",
            )
            assets.append(asset.to_dict())

    return blocks, units, assets


def _is_numbered_section_heading(text: str) -> bool:
    """Recover common PPT section headings when the layout has no title placeholder.

    Some imported decks use plain text boxes for slide titles.  Restrict the
    fallback to short Chinese top-level numbered headings so body text and
    captions do not become graph concepts merely because of their position.
    """
    normalized = " ".join((text or "").split())
    return bool(
        normalized
        and len(normalized) <= 80
        and re.match(r"^[一二三四五六七八九十]+[、．.]", normalized)
    )
