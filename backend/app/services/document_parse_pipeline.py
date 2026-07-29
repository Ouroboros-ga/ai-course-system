"""P0-4 真实 DocumentIR/OCR/图谱构建流水线。

由 document_parse_handler 调用，串联：
1. 从 SourceMaterialVersion 读取 object_key
2. 从对象存储读取文件内容
3. DocumentProbe 探测格式
4. ParsePlanner 规划解析步骤
5. ParserRegistry 查找 Provider
6. Provider.parse 解析
7. map_*_output_to_ir 生成 DocumentIR blocks
8. 写入 DocumentBlock / EvidenceSpan / GraphCandidateBatch

设计要点：
- 真实 Provider 不可用时降级到 Fake Provider（保证端到端流程可用）
- 解析失败必须返回结构化错误，不伪装成功
- EvidenceSpan 候选基于规则抽取（每个非空 text block 生成一个候选）
- GraphCandidateBatch 自动创建并标记 succeeded
"""
from __future__ import annotations

import logging
import hashlib
import json
import io
import os
import tempfile
import unicodedata
from dataclasses import replace
from difflib import SequenceMatcher
from typing import Any, Optional

from sqlmodel import Session, select

from app.models.course_build_model import SourceMaterialVersion
from app.models.document_parse_model import (
    DocumentBlock,
    DocumentParseRun,
    EvidenceRenderAsset,
    EvidenceSpan,
    GraphCandidateBatch,
    ParseRunStatus,
    DocumentIRVersion,
)
from app.services.document_parse_service import (
    document_parse_service,
    graph_candidate_service,
)

logger = logging.getLogger(__name__)

_CACHE_PARSER_VERSIONS = {
    "canonical-projector": "1.0.0",
    "native-pptx": "1.0.0",
    "pdf-plumber": "1.0.0",
    "python-docx": "1.0.0",
    "tesseract-ocr": "1.0.0",
}


def _canonical_cache_key(*, source_sha256: str, parse_profile: str, pipeline: str) -> str:
    """Stable pre-parse key; changing parser or IR contract invalidates reuse."""
    from app.platform.document_intelligence.contracts import CURRENT_SCHEMA_VERSION
    identity = {
        "source_sha256": source_sha256,
        "parser_profile": parse_profile,
        "pipeline": pipeline,
        "parser_versions": _CACHE_PARSER_VERSIONS,
        "ir_schema_version": CURRENT_SCHEMA_VERSION.serialize(),
    }
    return hashlib.sha256(json.dumps(identity, sort_keys=True).encode("utf-8")).hexdigest()


def _find_cached_ir(session: Session, cache_key: str) -> Optional[DocumentIRVersion]:
    from app.services.object_storage import get_object_storage
    candidates = session.exec(select(DocumentIRVersion).where(
        DocumentIRVersion.cache_key == cache_key,
    ).order_by(DocumentIRVersion.created_at.desc())).all()
    storage = get_object_storage()
    return next((item for item in candidates if item.object_key and storage.exists(item.object_key)), None)


# ---------------------------------------------------------------------------
# Provider 注册表（懒加载，避免导入时副作用）
# ---------------------------------------------------------------------------


_PARSER_REGISTRY = None


def _get_parser_registry():
    """获取 ParserRegistry 单例，注册内置 Provider。

    P1-3: 用真实 Provider 替换 fake 实现：
    - NativePptxProvider: 真实 PPTX 解析（python-pptx）
    - PdfPlumberProvider: 真实 PDF 解析（pdfplumber，CPU-only）
    - TesseractOcrProvider: 真实 OCR（pytesseract，CPU-only）
    - OcrProvider: PaddleOCR 占位（GPU，未启用）

    所有 Provider 在依赖不可用时抛 ParseUnavailableError，绝不返回伪造结果。
    """
    global _PARSER_REGISTRY
    if _PARSER_REGISTRY is not None:
        return _PARSER_REGISTRY

    from app.platform.document_intelligence.registry import ParserRegistry
    from app.platform.document_intelligence.providers.native_pptx import NativePptxProvider
    from app.platform.document_intelligence.providers.pdf_plumber import PdfPlumberProvider
    from app.platform.document_intelligence.providers.ocr_provider import (
        TesseractOcrProvider,
        OcrProvider,
    )
    from app.platform.document_intelligence.providers.python_docx import (
        PythonDocxProvider,
    )

    registry = ParserRegistry()
    registry.register(NativePptxProvider())
    registry.register(PdfPlumberProvider())
    registry.register(TesseractOcrProvider())
    registry.register(OcrProvider())
    registry.register(PythonDocxProvider())  # Step 3: DOCX native parsing
    _PARSER_REGISTRY = registry
    return registry


# ---------------------------------------------------------------------------
# 主入口
# ---------------------------------------------------------------------------


async def run_parse_pipeline(
    session: Session,
    *,
    course_id: int,
    run_id: str,
    material_id: str,
    material_version_id: Optional[str],
    pipeline: str = "standard",
    stale_strategy: str = "mark_stale",
) -> tuple[int, int, int]:
    """执行真实解析流水线，返回 (block_count, evidence_span_count, graph_candidate_count)。

    失败时抛 ParsePipelineError，由 handler 捕获并转为 mark_failed。
    """
    # 0. 读取 parse_run，获取 initiated_by（用于 GraphCandidateBatch）
    parse_run = session.exec(
        select(DocumentParseRun).where(
            DocumentParseRun.run_id == run_id,
            DocumentParseRun.course_id == course_id,
        )
    ).first()
    if parse_run is None:
        raise ParsePipelineError(
            "VALIDATION_FAILED",
            f"DocumentParseRun not found: run_id={run_id}, course_id={course_id}",
        )
    initiated_by = parse_run.initiated_by or 0

    # 1. 读取 SourceMaterialVersion
    version = _resolve_material_version(
        session, course_id=course_id, material_id=material_id, version_id=material_version_id,
    )
    object_key = version.file_path
    if not object_key:
        raise ParsePipelineError(
            "SOURCE_UNAVAILABLE",
            f"SourceMaterialVersion {version.version_id} has no file_path (object_key)",
        )

    # 2. 从对象存储读取文件内容
    from app.services.object_storage import get_object_storage
    storage = get_object_storage()
    try:
        content = storage.get(object_key)
    except FileNotFoundError as exc:
        raise ParsePipelineError(
            "SOURCE_UNAVAILABLE",
            f"Object not found in storage: {object_key}: {exc}",
        ) from exc
    except Exception as exc:
        raise ParsePipelineError(
            "SOURCE_UNAVAILABLE",
            f"Failed to read object {object_key}: {exc}",
        ) from exc

    # 3. Build the stable identity from the original SourceMaterialVersion.
    # Legacy Office conversion below is only a transient parser input and must
    # never replace the original source artifact or stored version.
    from app.platform.document_intelligence.source_artifact import SourceArtifact

    source = SourceArtifact.from_bytes(
        content,
        filename=object_key,
        mime=version.mime_type or "application/octet-stream",
        uri=object_key,
    )

    # Page renders are a first-class provenance projection, independent of
    # whether the parser needed OCR.  Native PPTX parsing can therefore still
    # provide a clickable source page when LibreOffice is unavailable.
    source_name = object_key.lower()
    if source_name.endswith(".pptx") and not session.exec(select(EvidenceRenderAsset).where(
        EvidenceRenderAsset.course_id == course_id,
        EvidenceRenderAsset.run_id == run_id,
    )).first():
        _persist_python_pptx_renders(
            content_bytes=content,
            source=source,
            target_pages=[],
            run_id=run_id,
            course_id=course_id,
            document_id=parse_run.document_id,
            session=session,
            warnings=[],
        )

    cache_key = _canonical_cache_key(
        source_sha256=version.file_hash or source.sha256,
        parse_profile=parse_run.parse_profile,
        pipeline=pipeline,
    )
    cached = _find_cached_ir(session, cache_key)
    if cached is not None:
        logger.info(
            "document_parse cache hit: run=%s source_sha256=%s canonical_ir=%s",
            run_id, (version.file_hash or source.sha256)[:16], cached.ir_version_id,
        )
        from app.platform.document_intelligence.canonical import DocumentIRProjector
        from app.platform.document_intelligence.document_ir.serialization import deserialize_document_ir
        from app.services.object_storage import get_object_storage

        cached_ir = deserialize_document_ir(get_object_storage().get(cached.object_key).decode("utf-8"))
        stored_ir, block_count, evidence_span_count = DocumentIRProjector().persist_and_project(
            session,
            course_id=course_id,
            material_version_id=material_version_id,
            run_id=run_id,
            previous_run_id=parse_run.prev_run_id,
            document_ir=cached_ir,
            quality_decision=None,
            provider_versions=cached.parser_versions,
            parse_outcome=cached.parse_outcome,
            parser_profile=parse_run.parse_profile,
            cache_key=cache_key,
            cache_source=cached,
        )
        parse_run.document_id = cached_ir.document_id
        parse_run.document_ir_version_id = stored_ir.ir_version_id
        _bind_render_assets_to_document(
            session, course_id=course_id, run_id=run_id, document_id=cached_ir.document_id,
        )
        session.add(parse_run)
        batch = graph_candidate_service.create_batch(session, course_id=course_id, parse_run_id=run_id, initiated_by=initiated_by)
        nodes, relations = _build_graph_candidates(session, course_id=course_id, run_id=run_id, ir_version_id=stored_ir.ir_version_id)
        graph_candidate_service.mark_succeeded(session, batch_id=batch.batch_id, course_id=course_id,
                                              node_candidate_count=len(nodes), relation_candidate_count=len(relations),
                                              node_candidates=nodes, relation_candidates=relations)
        return block_count, evidence_span_count, 1

    # Legacy binary Office is accepted by the import API but has no native
    # parser. Convert it inside this worker to a transient PDF and use that
    # only for Probe/Provider/OCR execution; canonical identity stays source.
    parse_source = source
    source_name = object_key.lower()
    if source_name.endswith((".doc", ".ppt")):
        from app.platform.document_intelligence.libreoffice_converter import ConversionError, libreoffice_converter
        suffix = os.path.splitext(object_key)[1]
        try:
            with tempfile.TemporaryDirectory(prefix="legacy_office_parse_") as temp_dir:
                source_path = os.path.join(temp_dir, f"source{suffix}")
                with open(source_path, "wb") as handle:
                    handle.write(content)
                conversion = libreoffice_converter.convert_to_pdf(source_path, output_dir=temp_dir)
                with open(conversion.pdf_path, "rb") as handle:
                    converted_pdf = handle.read()
        except ConversionError as exc:
            raise ParsePipelineError(exc.error_code, exc.message) from exc
        except Exception as exc:
            raise ParsePipelineError(
                "CONVERSION_FAILED", f"Legacy Office conversion failed: {type(exc).__name__}: {exc}",
            ) from exc
        parse_source = replace(
            source,
            filename=f"{os.path.splitext(object_key)[0]}.pdf",
            mime="application/pdf",
            size_bytes=len(converted_pdf),
            uri=None,
            data=converted_pdf,
        )

    # 4. Probe and plan the actual parser input while retaining the original
    # source artifact ID for plan auditability and canonical projection.
    from app.platform.document_intelligence.probe import DocumentProbe
    probe = DocumentProbe()
    probe_result = probe.probe(
        parse_source.data or b"", filename=parse_source.filename, mime=parse_source.mime,
    )
    if not probe_result.is_parseable():
        raise ParsePipelineError("UNSUPPORTED_FORMAT", probe_result.error or "Source format is not parseable")

    from app.platform.document_intelligence.planner import ParsePlanner
    planner = ParsePlanner()
    registry = _get_parser_registry()
    planner.set_available_providers(list(registry.list_providers()))
    plan = planner.plan(probe_result, source.artifact_id)
    required_ocr_pages: set[int] = set()
    for step in plan.steps:
        if step.priority.value != "enrichment" or step.provider_name not in {"tesseract-ocr", "paddleocr"}:
            continue
        pages = [int(page) for page in (step.config.get("pages") or []) if str(page).isdigit()]
        # An OCR-required plan without an explicit page list applies to every
        # known page.  In particular, PDF enrichment is intentionally all-page.
        if not pages and probe_result.page_or_slide_count:
            pages = list(range(1, int(probe_result.page_or_slide_count) + 1))
        required_ocr_pages.update(pages)
    scope_pages = _scoped_pages(parse_run.reparse_scope, probe_result.page_or_slide_count)
    high_quality_ocr = parse_run.parse_profile == "high_quality_ocr"

    parser_run_id = f"prun_{run_id}"
    primary_blocks: list[dict[str, Any]] = []
    enrichment_blocks: list[dict[str, Any]] = []
    units_data: list[dict[str, Any]] = []
    assets_data: list[dict[str, Any]] = []
    warnings: list[str] = []
    used_providers: list[str] = []
    provider_versions: dict[str, str] = {}

    from app.platform.document_intelligence.planner import ParsePriority
    from app.platform.document_intelligence.providers.python_docx import (
        PythonDocxProvider,
        map_docx_output_to_ir,
    )

    # ------------------------------------------------------------------
    # Step 3 组合式解析：消费完整 ParsePlan（PRIMARY + ENRICHMENT）。
    # 旧实现：第一个成功 provider 就 break -> PDF 不 OCR、PPTX 低文本页无补充。
    # 新实现：PRIMARY 收集原生文本块；ENRICHMENT 经 DocumentOcrPort 补充 OCR 块；
    #         若两者都有，用 BlockReconciler 合并（原生优先，OCR 补图片/扫描）。
    # ------------------------------------------------------------------

    def _map_step_output(step_name: str, output):
        """把某个 provider 的 ParserOutput 映射为 IR block dicts。"""
        if step_name == "native-pptx":
            from app.platform.document_intelligence.providers.native_pptx import (
                map_pptx_output_to_ir,
            )
            b, u, a = map_pptx_output_to_ir(output, source, run_id, parser_run_id)
        elif step_name == "pdf-plumber":
            from app.platform.document_intelligence.providers.pdf_plumber import (
                map_pdf_plumber_output_to_ir,
            )
            b, u, a = map_pdf_plumber_output_to_ir(output, source, run_id, parser_run_id)
        elif step_name in ("tesseract-ocr", "paddleocr"):
            from app.platform.document_intelligence.providers.ocr_provider import (
                map_ocr_output_to_ir,
            )
            b, u, a = map_ocr_output_to_ir(output, source, run_id, parser_run_id)
        elif step_name == "python-docx":
            b, u, a = map_docx_output_to_ir(output, source, run_id, parser_run_id)
        else:
            b, u, a = [], [], []
        return b, u, a

    # 执行 PRIMARY 步骤（收集原生文本，不 break）
    for step in plan.steps:
        if step.priority != ParsePriority.PRIMARY:
            continue
        provider = registry.get(step.provider_name)
        if provider is None:
            warnings.append(f"primary provider {step.provider_name} not registered; skipped")
            continue
        try:
            output = await provider.parse(parse_source, plan)
            used_providers.append(step.provider_name)
            provider_versions[step.provider_name] = output.provider_version
            warnings.extend(str(warning) for warning in output.warnings)
            b, u, a = _map_step_output(step.provider_name, output)
            if u:
                units_data.extend(u)
            if a:
                assets_data.extend(a)
            primary_blocks.extend(b)
        except Exception as exc:
            warnings.append(
                f"primary provider {step.provider_name} failed: {type(exc).__name__}: {exc}"
            )
            continue

    # 执行 ENRICHMENT 步骤（OCR 补充）。优先用独立 PaddleOCR 服务（DocumentOcrPort），
    # 其次回退到进程内 Tesseract provider。
    for step in plan.steps:
        if step.priority != ParsePriority.ENRICHMENT:
            continue
        if high_quality_ocr:
            step = type(step)(
                provider_name=step.provider_name,
                priority=step.priority,
                timeout_ms=step.timeout_ms,
                config={**step.config, "pages": scope_pages or step.config.get("pages", []), "quality": "high"},
                enrichment_for=step.enrichment_for,
            )
        # 先尝试经 DocumentOcrPort 调独立 PaddleOCR 服务
        if step.provider_name in ("tesseract-ocr", "paddleocr"):
            ocr_blocks = await _ocr_enrichment_via_port(
                parse_source, plan, step, warnings, run_id, parser_run_id,
                session=session, course_id=course_id, document_id=parse_run.document_id,
            )
            if ocr_blocks:
                enrichment_blocks.extend(ocr_blocks)
                provider_versions["paddleocr-service"] = ocr_blocks[0].get(
                    "provider_version", "paddleocr-service"
                )
                continue
        # Office sources must be rendered before OCR.  Their ZIP bytes are not
        # image input and must never be passed to the in-process OCR fallback.
        if _is_office_source(source):
            warnings.append("Office OCR port produced no output; skipped unsafe raw-byte OCR fallback")
            continue
        if high_quality_ocr:
            warnings.append("High-quality OCR requires the isolated OCR service; skipped in-process fallback")
            continue
        # 回退：进程内 provider（如 TesseractOcrProvider）
        provider = registry.get(step.provider_name)
        if provider is None:
            continue
        try:
            output = await provider.parse(parse_source, plan)
            used_providers.append(step.provider_name)
            provider_versions[step.provider_name] = output.provider_version
            b, _, _ = _map_step_output(step.provider_name, output)
            enrichment_blocks.extend(b)
        except Exception as exc:
            warnings.append(
                f"enrichment provider {step.provider_name} failed: {type(exc).__name__}: {exc}"
            )
            continue

    # A high-quality OCR request is intentionally scoped and OCR-first.  It
    # must not silently widen into a full-document or in-process fallback run.
    if high_quality_ocr and not enrichment_blocks:
        from app.platform.document_intelligence.planner import ParseStep
        forced_step = ParseStep(
            provider_name="paddleocr", priority=ParsePriority.ENRICHMENT,
            timeout_ms=300000, config={"pages": scope_pages, "quality": "high"},
        )
        enrichment_blocks = await _ocr_enrichment_via_port(
            parse_source, plan, forced_step, warnings, run_id, parser_run_id,
            session=session, course_id=course_id, document_id=parse_run.document_id,
        )
        if not enrichment_blocks:
            raise ParsePipelineError(
                "HIGH_QUALITY_OCR_UNAVAILABLE",
                "Scoped high-quality OCR produced no output; the prior IR remains unchanged",
            )

    # A scoped OCR reparse is its own candidate IR.  Combining it with every
    # native block would make an adopted page-range proposal look like a full
    # document replacement.
    source_name = (getattr(parse_source, "filename", "") or getattr(parse_source, "uri", "")).lower()
    docx_aligned_ocr: set[str] = set()
    if source_name.endswith(".docx") and primary_blocks and enrichment_blocks:
        docx_aligned_ocr = _align_docx_native_to_rendition(
            primary_blocks, enrichment_blocks, warnings,
        )
    if high_quality_ocr:
        blocks_data = enrichment_blocks
    # 合并：原生文本优先；有 enrichment 时经 BlockReconciler 去重/补充
    elif source_name.endswith(".docx") and enrichment_blocks and primary_blocks:
        blocks_data = _merge_docx_native_and_rendition(
            primary_blocks, enrichment_blocks, docx_aligned_ocr,
        )
    elif enrichment_blocks and primary_blocks:
        blocks_data = _reconcile_blocks(primary_blocks, enrichment_blocks, run_id, parser_run_id, warnings)
    elif primary_blocks:
        blocks_data = primary_blocks
    else:
        # 无原生文本（如纯图片）：enrichment 即主结果
        blocks_data = enrichment_blocks

    if not blocks_data:
        raise ParsePipelineError(
            "PARSE_FAILED",
            f"All providers failed for {object_key}. Warnings: {warnings}",
        )

    # Required OCR is a quality boundary, not a transaction rollback.  Keep
    # the native candidate IR for review while making its partial state durable.
    ocr_pages_done = {int(block.get("page_or_slide") or 0) for block in enrichment_blocks}
    missing_required_ocr = sorted(required_ocr_pages - ocr_pages_done)
    if missing_required_ocr:
        warnings.append(f"required OCR unavailable for pages: {missing_required_ocr}")

    # 6. Assemble, validate, store and then project exactly one canonical IR.
    # Database rows are consumers of the immutable object, not a competing
    # source of IDs or truncated evidence snippets.
    from app.platform.document_intelligence.canonical import (
        CanonicalDocumentIRAssembler,
        DocumentIRProjector,
    )
    from app.platform.document_intelligence.document_ir.models import (
        validate_no_duplicate_ids,
        validate_reference_integrity,
    )
    from app.platform.document_intelligence.quality import QualityScorer, QualityVerdict

    document_ir = CanonicalDocumentIRAssembler().assemble(
        source=source,
        run_id=run_id,
        parser_run_id=parser_run_id,
        blocks=blocks_data,
        units=units_data,
        assets=assets_data,
        warnings=warnings,
        provider_versions=provider_versions,
    )
    integrity_errors = validate_reference_integrity(document_ir) + validate_no_duplicate_ids(document_ir)
    if integrity_errors:
        raise ParsePipelineError("IR_VALIDATION_FAILED", "; ".join(integrity_errors[:5]))
    # Planner thresholds describe desired capabilities, whereas QualityScorer
    # consumes its own explicit ``*_min`` / ``*_max`` gate schema.
    quality_decision = QualityScorer().evaluate(document_ir)
    if quality_decision.verdict == QualityVerdict.FAIL:
        raise ParsePipelineError("QUALITY_GATE_FAILED", ",".join(reason.value for reason in quality_decision.reasons))
    if missing_required_ocr:
        quality_decision = replace(
            quality_decision,
            verdict=QualityVerdict.NEEDS_REVIEW,
            needs_review=True,
        )
    stored_ir, block_count, evidence_span_count = DocumentIRProjector().persist_and_project(
        session,
        course_id=course_id,
        material_version_id=material_version_id,
        run_id=run_id,
        previous_run_id=parse_run.prev_run_id,
        document_ir=document_ir,
        quality_decision=quality_decision,
        provider_versions=provider_versions,
        parse_outcome=_parse_outcome(quality_decision, provider_versions, warnings),
        parser_profile=parse_run.parse_profile,
        cache_key=cache_key,
    )
    parse_run.document_id = document_ir.document_id
    parse_run.document_ir_version_id = stored_ir.ir_version_id
    _bind_render_assets_to_document(
        session, course_id=course_id, run_id=run_id, document_id=document_ir.document_id,
    )
    session.add(parse_run)

    # 7. 创建 GraphCandidateBatch
    batch = graph_candidate_service.create_batch(
        session,
        course_id=course_id,
        parse_run_id=run_id,
        initiated_by=initiated_by,
    )
    node_candidates, relation_candidates = _build_graph_candidates(
        session,
        course_id=course_id,
        run_id=run_id,
        ir_version_id=stored_ir.ir_version_id,
    )
    graph_candidate_service.mark_succeeded(
        session,
        batch_id=batch.batch_id,
        course_id=course_id,
        node_candidate_count=len(node_candidates),
        relation_candidate_count=len(relation_candidates),
        node_candidates=node_candidates,
        relation_candidates=relation_candidates,
    )

    return block_count, evidence_span_count, 1  # graph_candidate_count=1（一个 batch）


# ---------------------------------------------------------------------------
# 辅助
# ---------------------------------------------------------------------------


def _bind_render_assets_to_document(
    session: Session,
    *,
    course_id: int,
    run_id: str,
    document_id: str,
) -> None:
    """Replace the provisional source-artifact scope with canonical IR ID."""
    assets = session.exec(select(EvidenceRenderAsset).where(
        EvidenceRenderAsset.course_id == course_id,
        EvidenceRenderAsset.run_id == run_id,
        EvidenceRenderAsset.document_id.is_(None),
    )).all()
    for asset in assets:
        asset.document_id = document_id
        session.add(asset)


def _build_graph_candidates(
    session: Session,
    *,
    course_id: int,
    run_id: str,
    ir_version_id: str,
) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    """Create reviewable, evidence-linked graph candidates without publishing them.

    A title-like source block becomes one concept candidate.  The immediately
    following explanatory blocks are evidence for that concept.  Consecutive
    concepts receive an explicit ``next_topic`` relation, which teachers can
    relabel or reject before a graph snapshot is published.
    """
    from app.models.document_parse_model import DocumentBlock, EvidenceAnchor

    blocks = list(session.exec(select(DocumentBlock).where(
        DocumentBlock.course_id == course_id,
        DocumentBlock.run_id == run_id,
        DocumentBlock.document_ir_version_id == ir_version_id,
    ).order_by(DocumentBlock.order_index)).all())
    anchors = list(session.exec(select(EvidenceAnchor).where(
        EvidenceAnchor.course_id == course_id,
        EvidenceAnchor.run_id == run_id,
        EvidenceAnchor.ir_version_id == ir_version_id,
    )).all())
    anchor_by_block = {anchor.block_id: anchor.anchor_id for anchor in anchors}

    candidates: list[dict[str, Any]] = []
    current: dict[str, Any] | None = None
    for block in blocks:
        text = (block.text or "").strip()
        if not text:
            continue
        title_like = block.semantic_role in {"section_title", "knowledge_title"}
        if not title_like:
            title_like = block.block_type == "title" and len(text) <= 120
        if title_like:
            candidate_id = "gcn_" + hashlib.sha256(
                f"{ir_version_id}:{block.block_id}".encode("utf-8")
            ).hexdigest()[:24]
            current = {
                "candidate_id": candidate_id,
                "label": text[:300],
                "kind": "concept",
                "status": "proposed",
                "confidence": round(float(block.confidence or 0.0), 3),
                "source_block_ids": [block.block_id],
                "anchor_ids": [anchor_by_block[block.block_id]] if block.block_id in anchor_by_block else [],
                "page_or_slide": block.page_or_slide,
            }
            candidates.append(current)
        elif current is not None:
            current["source_block_ids"].append(block.block_id)
            if block.block_id in anchor_by_block:
                current["anchor_ids"].append(anchor_by_block[block.block_id])

    relations = [{
        "candidate_id": "gcr_" + hashlib.sha256(
            f"{ir_version_id}:{left['candidate_id']}:{right['candidate_id']}".encode("utf-8")
        ).hexdigest()[:24],
        "source_candidate_id": left["candidate_id"],
        "target_candidate_id": right["candidate_id"],
        "relation_type": "next_topic",
        "status": "proposed",
        "confidence": min(left["confidence"], right["confidence"]),
        "anchor_ids": list(dict.fromkeys(left["anchor_ids"] + right["anchor_ids"])),
    } for left, right in zip(candidates, candidates[1:])]
    return candidates, relations


async def _ocr_enrichment_via_port(
    source,
    plan,
    step,
    warnings: list[str],
    run_id: str,
    parser_run_id: str,
    session: Optional[Session] = None,
    course_id: Optional[int] = None,
    document_id: Optional[str] = None,
) -> list[dict[str, Any]]:
    """经 DocumentOcrPort 调独立 PaddleOCR 服务，做 OCR enrichment。

    渲染源文档的指定页为图片，逐页调 OCR 服务，返回 IR block dicts。
    若 OCR 服务不可用（UnavailableOcrPort / 服务 503 / 超时），返回空列表，
    由调用方回退到进程内 Tesseract provider；不伪造输出。

    ``source`` 需有 ``uri``（对象存储 object_key）以便读取字节。
    """
    blocks: list[dict[str, Any]] = []
    try:
        from app.platform.document_intelligence.ocr_port import (
            get_ocr_port, OcrUnavailable,
        )
        port = get_ocr_port()
        if not port.is_available:
            warnings.append("OCR service unavailable via DocumentOcrPort; will try in-process fallback")
            return []

        content_bytes = _source_bytes(source)

        # 确定要 OCR 的页码（来自 plan step config）
        target_pages = step.config.get("pages", [1]) if step.config else [1]

        # Office documents are rendered page-by-page before OCR.  Passing the
        # raw DOCX/PPTX ZIP to image OCR would silently create unusable output.
        if _is_office_source(source):
            return _ocr_office_pages_via_port(
                port, content_bytes, source, target_pages, warnings, run_id, parser_run_id,
                session=session, course_id=course_id, document_id=document_id,
            )

        # 若源是 PDF，用 ocr_pdf；否则按图片逐页 ocr_image
        mime = (source.mime or "").lower()
        if "pdf" in mime:
            try:
                result = None
                from app.platform.document_intelligence.page_renderer import PdfiumPageRenderer
                rendered_pages = PdfiumPageRenderer().render(content_bytes, target_pages, dpi=180)
                blocks = _ocr_rendered_pages(
                    port, rendered_pages, source, run_id, parser_run_id,
                    session=session, course_id=course_id, document_id=document_id,
                    warnings=warnings,
                )
                if blocks:
                    return blocks
            except OcrUnavailable as exc:
                warnings.append(f"OCR service rendered PDF OCR failed: {exc.error_code}: {exc.message}")
                return []
            except Exception as exc:
                warnings.append(f"Selective PDF rendering failed: {type(exc).__name__}: {exc}")
                return []
        else:
            # 图片型 PPT 页或其他图片：逐页 OCR。此处 content_bytes 是整个文件；
            # 图片型 PPT 的逐页渲染由 planner/preprobe 决定页范围，简化为整文件一次 OCR。
            try:
                result = port.ocr_image(content_bytes, page=int(target_pages[0]) if target_pages else 1)
            except OcrUnavailable as exc:
                warnings.append(f"OCR service ocr_image failed: {exc.error_code}: {exc.message}")
                return []

        # 把 OcrResult 转成 IR block dicts（与 map_ocr_output_to_ir 同形）
        from app.platform.document_intelligence.contracts import (
            BoundingBox, CoordinateSpace,
        )
        from app.platform.document_intelligence.document_ir.models import Provenance
        idx = 0
        for page in result.pages:
            for blk in page.blocks:
                bbox = None
                if blk.bbox and len(list(blk.bbox)) == 4:
                    try:
                        bbox = BoundingBox(
                            x0=round(float(blk.bbox[0]), 6),
                            y0=round(float(blk.bbox[1]), 6),
                            x1=round(float(blk.bbox[2]), 6),
                            y1=round(float(blk.bbox[3]), 6),
                            coordinate_space=CoordinateSpace.NORMALIZED,
                        )
                    except (ValueError, TypeError):
                        bbox = None
                text = blk.text or ""
                provenance = Provenance(
                    artifact_id=source.artifact_id,
                    run_id=run_id,
                    parser_run_id=parser_run_id,
                    provider="paddleocr-service",
                    raw_locator=f"pages/{page.page}/blocks/{idx}",
                    page_or_slide=page.page,
                    bbox=bbox,
                    confidence=blk.confidence,
                )
                blocks.append({
                    "block_id": f"blk_ocr_p{page.page}_b{idx}",
                    "block_type": "paragraph",
                    "text": text,
                    "page_or_slide": page.page,
                    "bbox": bbox,
                    "char_start": 0,
                    "char_end": len(text),
                    "order_index": idx,
                    "provenance": provenance,
                    "provider_version": result.provider_version,
                })
                idx += 1
        return blocks
    except Exception as exc:
        warnings.append(f"OCR enrichment via port failed: {type(exc).__name__}: {exc}")
        return []


def _reconcile_blocks(
    primary_blocks: list[dict[str, Any]],
    enrichment_blocks: list[dict[str, Any]],
    run_id: str,
    parser_run_id: str,
    warnings: list[str],
) -> list[dict[str, Any]]:
    """用 BlockReconciler 合并原生文本块与 OCR 块。

    原生文本优先保留；OCR 补充图片/扫描文本；通过内容坐标去重避免双重文本。
    Reconciler 输入是 DocumentIR，需把 block dicts 转成 Block 对象再合并。
    """
    if not enrichment_blocks:
        return primary_blocks
    try:
        from app.platform.document_intelligence.reconciliation import BlockReconciler
        from app.platform.document_intelligence.document_ir.models import (
            DocumentIR, ContentBlock, Provenance,
        )
        from app.platform.document_intelligence.contracts import (
            BoundingBox, CoordinateSpace,
        )

        def _to_ir(blocks: list[dict[str, Any]]) -> DocumentIR:
            ir_blocks = []
            for b in blocks:
                bbox = b.get("bbox")
                ir_bbox = None
                if bbox is not None and not isinstance(bbox, dict):
                    try:
                        ir_bbox = BoundingBox(
                            x0=float(getattr(bbox, "x0", 0)),
                            y0=float(getattr(bbox, "y0", 0)),
                            x1=float(getattr(bbox, "x1", 0)),
                            y1=float(getattr(bbox, "y1", 0)),
                            coordinate_space=CoordinateSpace.NORMALIZED,
                        )
                    except Exception:
                        ir_bbox = None
                prov = b.get("provenance") or Provenance(
                    artifact_id="",
                    run_id=run_id,
                    parser_run_id=parser_run_id,
                    provider="unknown",
                    raw_locator="",
                    page_or_slide=int(b.get("page_or_slide", 1) or 1),
                    bbox=ir_bbox,
                    confidence=None,
                )
                ir_blocks.append(ContentBlock(
                    block_id=b.get("block_id", f"blk_{id(b)}"),
                    block_type=b.get("block_type", "paragraph"),
                    text=b.get("text", ""),
                    page_or_slide=int(b.get("page_or_slide", 1) or 1),
                    bbox=ir_bbox,
                    reading_order=b.get("reading_order", b.get("order_index", 0)),
                    heading_level=b.get("heading_level"),
                    style_hints=b.get("style_hints") or {},
                    parent_id=b.get("parent_id"),
                    visual_description=b.get("visual_description"),
                    provenance=(prov,) if not isinstance(prov, (tuple, list)) else tuple(prov),
                ))
            return DocumentIR(
                document_id="reconcile",
                artifact_id="reconcile",
                blocks=tuple(ir_blocks),
                units=(),
                assets=(),
                warnings=(),
            )

        primary_ir = _to_ir(primary_blocks)
        enrichment_ir = _to_ir(enrichment_blocks)
        reconciler = BlockReconciler()
        result = reconciler.reconcile(primary_ir, enrichment_ir, primary_priority=True)

        # 转回 block dicts（保留原生块的原结构，补充 OCR 新块）
        out: list[dict[str, Any]] = []
        for b in result.blocks:
            # 若该 block 来自 primary（按 block_id 匹配），保留原 dict
            match = next((pb for pb in primary_blocks if pb.get("block_id") == b.block_id), None)
            if match is not None:
                out.append(match)
            else:
                # 来自 enrichment 的新块
                prov = getattr(b, "provenance", None)
                out.append({
                    "block_id": b.block_id,
                    "block_type": getattr(b, "block_type", "paragraph"),
                    "text": getattr(b, "text", "") or "",
                    "page_or_slide": getattr(b, "page_or_slide", 1) or 1,
                    "bbox": getattr(b, "bbox", None),
                    "char_start": 0,
                    "char_end": len(getattr(b, "text", "") or ""),
                    "order_index": getattr(b, "reading_order", len(out)),
                    "heading_level": getattr(b, "heading_level", None),
                    "style_hints": getattr(b, "style_hints", {}) or {},
                    "parent_id": getattr(b, "parent_id", None),
                    "visual_description": getattr(b, "visual_description", None),
                    "provenance": prov,
                })
        return out
    except Exception as exc:
        warnings.append(
            f"BlockReconciler failed; falling back to primary+enrichment concat: {exc}"
        )
        # 降级：原生优先，enrichment 全量追加（不去重）
        return primary_blocks + enrichment_blocks


def _align_docx_native_to_rendition(
    native_blocks: list[dict[str, Any]],
    rendition_blocks: list[dict[str, Any]],
    warnings: list[str],
    *,
    minimum_similarity: float = 0.82,
) -> set[str]:
    """Attach rendition coordinates to DOCX blocks only after strong matching.

    DOCX semantic order and LibreOffice pagination are separate facts. The
    native block keeps its OOXML locator; this adds page/bbox metadata only
    when a SequenceMatcher comparison with OCR text is strong enough. Failed
    matches remain explicitly page-unknown instead of silently pointing at an
    arbitrary rendition page.
    """
    available = [block for block in rendition_blocks if _normalise_alignment_text(block.get("text", ""))]
    unmatched = 0
    matched_ids: set[str] = set()
    for native in native_blocks:
        native_text = _normalise_alignment_text(native.get("text", ""))
        if not native_text:
            continue
        best = None
        best_score = 0.0
        for candidate in available:
            score = SequenceMatcher(
                None, native_text, _normalise_alignment_text(candidate.get("text", "")),
                autojunk=False,
            ).ratio()
            if score > best_score:
                best, best_score = candidate, score
        if best is None or best_score < minimum_similarity:
            unmatched += 1
            continue
        style_hints = dict(native.get("style_hints") or {})
        style_hints.update({
            "rendition_page": best.get("page_or_slide"),
            "rendition_bbox": _serialise_bbox(best.get("bbox")),
            "alignment_confidence": round(best_score, 4),
            "rendition_locator": _provenance_locator(best.get("provenance")),
        })
        native["style_hints"] = style_hints
        available.remove(best)
        if best.get("block_id"):
            matched_ids.add(str(best["block_id"]))
    if unmatched:
        warnings.append(f"DOCX_RENDITION_ALIGNMENT_UNRESOLVED:{unmatched}")
    return matched_ids


def _merge_docx_native_and_rendition(
    native_blocks: list[dict[str, Any]],
    rendition_blocks: list[dict[str, Any]],
    aligned_rendition_ids: set[str],
) -> list[dict[str, Any]]:
    """Keep DOCX source semantics and append only visually new OCR blocks."""
    return [
        *native_blocks,
        *[
            block for block in rendition_blocks
            if str(block.get("block_id", "")) not in aligned_rendition_ids
        ],
    ]


def _normalise_alignment_text(value: Any) -> str:
    return "".join(unicodedata.normalize("NFKC", str(value or "")).split()).lower()


def _serialise_bbox(value: Any) -> Any:
    return value.to_dict() if hasattr(value, "to_dict") else value


def _provenance_locator(value: Any) -> str | None:
    if isinstance(value, (tuple, list)):
        value = value[0] if value else None
    return getattr(value, "raw_locator", None) if value is not None else None


def _resolve_material_version(
    session: Session,
    *,
    course_id: int,
    material_id: str,
    version_id: Optional[str],
) -> SourceMaterialVersion:
    """根据 version_id 解析 SourceMaterialVersion；未指定时取最新的 current。"""
    if version_id:
        stmt = select(SourceMaterialVersion).where(
            SourceMaterialVersion.version_id == version_id,
            SourceMaterialVersion.course_id == course_id,
            SourceMaterialVersion.material_id == material_id,
        )
    else:
        stmt = (
            select(SourceMaterialVersion)
            .where(SourceMaterialVersion.course_id == course_id)
            .where(SourceMaterialVersion.material_id == material_id)
            .order_by(SourceMaterialVersion.created_at.desc())
            .limit(1)
        )
    version = session.exec(stmt).first()
    if version is None:
        raise ParsePipelineError(
            "SOURCE_UNAVAILABLE",
            f"SourceMaterialVersion not found: version_id={version_id}, material_id={material_id}, course_id={course_id}",
        )
    return version


def _is_office_source(source: Any) -> bool:
    name = (getattr(source, "filename", "") or getattr(source, "uri", "")).lower()
    return name.endswith((".doc", ".docx", ".ppt", ".pptx"))


def _parse_outcome(quality_decision: Any, provider_versions: dict[str, str], warnings: list[str]) -> str:
    """Keep non-binary parse quality visible to teachers and downstream gates."""
    warning_text = " ".join(warnings).lower()
    if (
        "smartart" in warning_text
        or "table_structure_unresolved" in warning_text
        or "chart_review_required" in warning_text
        or "formula_review_required" in warning_text
        or "floating_drawing_review_required" in warning_text
    ):
        return "unsupported_visual_structure"
    if getattr(quality_decision, "needs_review", False):
        return "manual_review_required"
    verdict = getattr(getattr(quality_decision, "verdict", None), "value", "")
    if verdict == "borderline":
        return "partial_success"
    if any("ocr" in provider.lower() for provider in provider_versions):
        return "native_with_ocr"
    return "native_complete"


def _scoped_pages(scope: dict | None, page_count: int) -> list[int]:
    """Validate an explicit OCR scope; no scope means all known pages."""
    pages = list((scope or {}).get("pages") or [])
    if not pages and page_count:
        return list(range(1, page_count + 1))
    cleaned: list[int] = []
    for page in pages:
        try:
            value = int(page)
        except (TypeError, ValueError):
            continue
        if value >= 1:
            cleaned.append(value)
    cleaned = sorted(set(cleaned))
    if page_count:
        cleaned = [page for page in cleaned if page <= page_count]
    return cleaned


def _source_bytes(source: Any) -> bytes:
    data = getattr(source, "data", None)
    if data is not None:
        return data
    from app.services.object_storage import get_object_storage
    if not getattr(source, "uri", None):
        raise ValueError("source has neither in-memory bytes nor an object-storage URI")
    return get_object_storage().get(source.uri)


def _ocr_office_pages_via_port(
    port: Any,
    content_bytes: bytes,
    source: Any,
    target_pages: list[int],
    warnings: list[str],
    run_id: str,
    parser_run_id: str,
    session: Optional[Session] = None,
    course_id: Optional[int] = None,
    document_id: Optional[str] = None,
) -> list[dict[str, Any]]:
    """Render Office pages before recognition and preserve their page numbers."""
    suffix = os.path.splitext(getattr(source, "filename", "") or getattr(source, "uri", ""))[1] or ".bin"
    try:
        from app.platform.document_intelligence.libreoffice_converter import libreoffice_converter
        with tempfile.TemporaryDirectory(prefix="document_ocr_") as temp_dir:
            source_path = os.path.join(temp_dir, f"source{suffix}")
            with open(source_path, "wb") as handle:
                handle.write(content_bytes)
            conversion = libreoffice_converter.convert_to_pdf(source_path, output_dir=temp_dir)
            images = libreoffice_converter.render_pages(conversion.pdf_path, output_dir=temp_dir)
            selected = set(target_pages) if target_pages else set(range(1, len(images) + 1))
            rows: list[tuple[int, Any]] = []
            for page_number, image_path in enumerate(images, start=1):
                if page_number not in selected:
                    continue
                with open(image_path, "rb") as handle:
                    image_bytes = handle.read()
                _persist_evidence_render_asset(
                    session=session, course_id=course_id, document_id=document_id,
                    run_id=run_id, page_number=page_number, image_bytes=image_bytes,
                    warnings=warnings,
                )
                result = port.ocr_image(image_bytes, page=page_number)
                rows.extend((page_number, block) for page in result.pages for block in page.blocks)
                provider_version = result.provider_version
    except Exception as exc:
        warnings.append(f"Office OCR rendering failed: {type(exc).__name__}: {exc}")
        _persist_python_pptx_renders(
            content_bytes=content_bytes,
            source=source,
            target_pages=target_pages,
            run_id=run_id,
            course_id=course_id,
            document_id=document_id,
            session=session,
            warnings=warnings,
        )
        return []

    from app.platform.document_intelligence.contracts import BoundingBox, CoordinateSpace
    from app.platform.document_intelligence.document_ir.models import Provenance
    blocks: list[dict[str, Any]] = []
    for index, (page_number, raw_block) in enumerate(rows):
        bbox = None
        if raw_block.bbox and len(list(raw_block.bbox)) == 4:
            try:
                bbox = BoundingBox(
                    x0=round(float(raw_block.bbox[0]), 6), y0=round(float(raw_block.bbox[1]), 6),
                    x1=round(float(raw_block.bbox[2]), 6), y1=round(float(raw_block.bbox[3]), 6),
                    coordinate_space=CoordinateSpace.NORMALIZED,
                )
            except (ValueError, TypeError):
                pass
        text = raw_block.text or ""
        blocks.append({
            "block_id": f"blk_ocr_p{page_number}_b{index}", "block_type": "paragraph", "text": text,
            "page_or_slide": page_number, "bbox": bbox, "char_start": 0, "char_end": len(text),
            "order_index": index, "provider_version": provider_version,
            "provenance": Provenance(
                artifact_id=source.artifact_id, run_id=run_id, parser_run_id=parser_run_id,
                provider="paddleocr-service", raw_locator=f"pages/{page_number}/blocks/{index}",
                page_or_slide=page_number, bbox=bbox, confidence=raw_block.confidence,
            ),
        })
    return blocks


def _persist_python_pptx_renders(
    *,
    content_bytes: bytes,
    source: Any,
    target_pages: list[int],
    run_id: str,
    course_id: Optional[int],
    document_id: Optional[str],
    session: Optional[Session],
    warnings: list[str],
) -> None:
    """Persist a deterministic local PPTX rendition when LibreOffice is absent.

    This is a source-derived fallback, not a fake page: text boxes and embedded
    images are read from the actual PPTX package and placed on a normalized
    slide canvas.  Native LibreOffice rendering remains preferred when present.
    """
    if session is None or course_id is None:
        return
    try:
        from pptx import Presentation
        from PIL import Image, ImageDraw, ImageFont

        presentation = Presentation(io.BytesIO(content_bytes))
        canvas_width, canvas_height = 1600, 900
        scale_x = canvas_width / float(presentation.slide_width)
        scale_y = canvas_height / float(presentation.slide_height)
        font_candidates = [
            os.path.join(os.environ.get("WINDIR", "C:/Windows"), "Fonts", "msyh.ttc"),
            os.path.join(os.environ.get("WINDIR", "C:/Windows"), "Fonts", "simhei.ttf"),
        ]
        font_path = next((path for path in font_candidates if os.path.isfile(path)), None)

        for page_number, slide in enumerate(presentation.slides, start=1):
            if target_pages and page_number not in set(target_pages):
                continue
            canvas = Image.new("RGB", (canvas_width, canvas_height), "white")
            draw = ImageDraw.Draw(canvas)
            for shape in slide.shapes:
                left = max(0, int(shape.left * scale_x))
                top = max(0, int(shape.top * scale_y))
                width = max(1, int(shape.width * scale_x))
                height = max(1, int(shape.height * scale_y))
                if getattr(shape, "shape_type", None) == 13 and getattr(shape, "image", None):
                    try:
                        with Image.open(io.BytesIO(shape.image.blob)) as embedded:
                            embedded = embedded.convert("RGB")
                            embedded.thumbnail((width, height))
                            canvas.paste(embedded, (left, top))
                    except Exception as image_error:
                        warnings.append(f"PPTX image rendition skipped on page {page_number}: {image_error}")
                if not getattr(shape, "has_text_frame", False):
                    continue
                text = "\n".join(
                    paragraph.text for paragraph in shape.text_frame.paragraphs if paragraph.text
                ).strip()
                if not text:
                    continue
                font_size = 22
                try:
                    first_run = next(run for paragraph in shape.text_frame.paragraphs for run in paragraph.runs if run.text)
                    if first_run.font.size:
                        font_size = max(12, min(72, int(first_run.font.size.pt * scale_y)))
                except (StopIteration, AttributeError, TypeError, ValueError):
                    pass
                font = ImageFont.truetype(font_path, font_size) if font_path else ImageFont.load_default()
                draw.multiline_text((left + 8, top + 8), text, fill=(25, 35, 50), font=font, spacing=6)

            output = io.BytesIO()
            canvas.save(output, format="PNG")
            _persist_evidence_render_asset(
                session=session,
                course_id=course_id,
                document_id=document_id,
                run_id=run_id,
                page_number=page_number,
                image_bytes=output.getvalue(),
                warnings=warnings,
            )
    except Exception as exc:
        warnings.append(f"Python PPTX fallback rendering failed: {type(exc).__name__}: {exc}")


def _persist_evidence_render_asset(
    *,
    session: Optional[Session],
    course_id: Optional[int],
    document_id: Optional[str],
    run_id: str,
    page_number: int,
    image_bytes: bytes,
    warnings: list[str],
) -> None:
    """Store a rendered Office page as a course-scoped Evidence Viewer asset."""
    if session is None or course_id is None:
        return
    import hashlib

    content_hash = hashlib.sha256(image_bytes).hexdigest()
    existing = session.exec(select(EvidenceRenderAsset).where(
        EvidenceRenderAsset.course_id == course_id,
        EvidenceRenderAsset.run_id == run_id,
        EvidenceRenderAsset.document_id == document_id,
        EvidenceRenderAsset.page_number == page_number,
        EvidenceRenderAsset.content_hash == content_hash,
    )).first()
    if existing:
        return
    object_key = f"evidence-render/course{course_id}/{document_id or run_id}/{run_id}/page-{page_number}.png"
    try:
        from app.services.object_storage import get_object_storage
        storage = get_object_storage()
        if not storage.exists(object_key):
            storage.put(object_key, image_bytes, mime_type="image/png")
        from PIL import Image
        image = Image.open(io.BytesIO(image_bytes))
        session.add(EvidenceRenderAsset(
            course_id=course_id,
            run_id=run_id,
            document_id=document_id,
            page_number=page_number,
            object_key=object_key,
            width=image.width,
            height=image.height,
            content_hash=content_hash,
        ))
    except Exception as exc:
        warnings.append(f"Evidence render asset persistence failed: {type(exc).__name__}: {exc}")


def _ocr_rendered_pages(
    port,
    rendered_pages,
    source,
    run_id: str,
    parser_run_id: str,
    *,
    session: Optional[Session] = None,
    course_id: Optional[int] = None,
    document_id: Optional[str] = None,
    warnings: Optional[list[str]] = None,
) -> list[dict[str, Any]]:
    """Convert page images to OCR blocks while preserving rendition provenance."""
    from app.platform.document_intelligence.contracts import BoundingBox, CoordinateSpace
    from app.platform.document_intelligence.document_ir.models import Provenance
    blocks: list[dict[str, Any]] = []
    for page in rendered_pages:
        _persist_evidence_render_asset(
            session=session,
            course_id=course_id,
            document_id=document_id,
            run_id=run_id,
            page_number=page.unit_index,
            image_bytes=page.image_bytes,
            warnings=warnings if warnings is not None else [],
        )
        result = port.ocr_image(page.image_bytes, page=page.unit_index)
        for raw_index, raw_block in enumerate(result.pages[0].blocks if result.pages else []):
            bbox = BoundingBox(
                x0=float(raw_block.bbox[0]), y0=float(raw_block.bbox[1]),
                x1=float(raw_block.bbox[2]), y1=float(raw_block.bbox[3]),
                coordinate_space=CoordinateSpace.NORMALIZED,
            ) if len(raw_block.bbox) == 4 else None
            text = raw_block.text or ""
            index = len(blocks)
            blocks.append({
                "block_id": f"blk_ocr_p{page.unit_index}_b{raw_index}",
                "block_type": "paragraph", "text": text,
                "page_or_slide": page.unit_index, "bbox": bbox,
                "char_start": 0, "char_end": len(text), "order_index": index,
                "provider_version": result.provider_version,
                "provenance": Provenance(
                    artifact_id=source.artifact_id, run_id=run_id, parser_run_id=parser_run_id,
                    provider="paddleocr-service",
                    raw_locator=f"rendition/{page.rendition_artifact_id}/blocks/{raw_index}",
                    page_or_slide=page.unit_index, bbox=bbox, confidence=raw_block.confidence,
                ),
            })
    return blocks


def _map_block_type(ir_block_type: str) -> str:
    """把 DocumentIR block_type 映射到 DocumentBlock.block_type。"""
    mapping = {
        "paragraph": "text",
        "heading": "title",
        "image": "figure_caption",
        "table": "table_cell",
        "unknown": "text",
    }
    return mapping.get(ir_block_type, "text")


def _infer_semantic_role(*, text: str, block_type: str, heading_level: Optional[int], style_hints: dict) -> str:
    """Infer a conservative role; this is a hint, never a publish decision."""
    compact = (text or "").strip().replace(" ", "")
    if any(token in compact for token in ("练习", "习题", "思考题", "试一试", "课后题")):
        return "practice_suggestion"
    if any(token in compact for token in ("示例", "例题", "案例", "例如", "代码演示")):
        return "example"
    if block_type == "title" or heading_level is not None or style_hints.get("is_heading"):
        return "section_title" if any(token in compact for token in ("章", "节", "单元", "模块", "部分")) else "knowledge_title"
    return "explanation"


class ParsePipelineError(Exception):
    """解析流水线错误，携带 error_code 与 message。"""

    def __init__(self, error_code: str, message: str) -> None:
        super().__init__(message)
        self.error_code = error_code
        self.message = message
