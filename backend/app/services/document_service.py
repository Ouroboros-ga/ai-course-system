"""
文档解析服务
负责文件解析、Markdown转换、智课脚本生成、RAG预处理
"""

import os
import json
import re
import logging
from pathlib import Path
from typing import Optional, Dict, Any, List, Tuple
from dataclasses import dataclass, field
from datetime import datetime

from app.common.llm_client import llm_client, Message
from app.common.RAG import rag_pipeline
from app.common.prompts.document_analysis import (
    KNOWLEDGE_EXTRACTION_PROMPT,
    build_knowledge_extraction_prompt
)
from app.common.prompts.knowledge_to_script import (
    KNOWLEDGE_TO_SCRIPT_PROMPT,
    build_knowledge_to_script_prompt
)

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


@dataclass
class ParseResult:
    """文件解析结果"""
    markdown_content: str
    filename: str
    file_path: str
    file_size: int
    parse_method: str
    error: Optional[str] = None
    ai_formatted: Optional[List[Dict[str, Any]]] = None
    doc_title: Optional[str] = None


@dataclass
class StructureResult:
    """结构化解析结果"""
    groups: List[Dict[str, Any]] = field(default_factory=list)
    texts: List[Dict[str, Any]] = field(default_factory=list)
    tables: List[Dict[str, Any]] = field(default_factory=list)
    pictures: List[Dict[str, Any]] = field(default_factory=list)
    raw_content: str = ""


@dataclass
class ScriptNode:
    """脚本节点"""
    chapter_id: str
    node_type: str
    title: str
    content: str
    page_start: int = 1
    page_end: int = 1
    duration: int = 60
    is_key_point: bool = False
    timestamp_start: float = 0.0
    timestamp_end: float = 0.0


@dataclass
class ScriptResult:
    """智课脚本结果"""
    title: str
    summary: str
    keywords: List[str]
    total_duration: int
    nodes: List[ScriptNode]
    script_content: Dict[str, Any]
    beautiful_markdown: str = ""


@dataclass
class RAGProcessResult:
    """RAG预处理结果"""
    formula_count: int = 0
    table_count: int = 0
    domain_term_count: int = 0
    tree_node_count: int = 0
    processed_text: str = ""
    knowledge_points: List[Dict[str, Any]] = field(default_factory=list)
    error: Optional[str] = None


@dataclass
class DocumentProcessResult:
    """完整文档处理结果"""
    parse_result: ParseResult
    structure_result: StructureResult
    script_result: ScriptResult
    rag_result: RAGProcessResult
    mind_map: Dict[str, Any]


class DocumentParser:
    """文档解析器"""
    
    @staticmethod
    async def parse_file(file_path: Path, filename: str) -> ParseResult:
        """
        解析文件，返回Markdown内容
        
        Args:
            file_path: 文件路径
            filename: 文件名
            
        Returns:
            ParseResult: 解析结果
        """
        file_size = file_path.stat().st_size if file_path.exists() else 0
        ai_formatted = None
        doc_title = None
        
        json_path = file_path.parent / f"{file_path.stem}_docling.json"
        if json_path.exists():
            try:
                with open(json_path, 'r', encoding='utf-8') as f:
                    docling_json = json.load(f)
                    ai_formatted = docling_json.get('ai_formatted', [])
                    doc_title = docling_json.get('title', Path(filename).stem)
                    logger.info(f"[DocumentParser] 找到docling JSON文件，加载ai_formatted数据: {len(ai_formatted)} 条")
            except Exception as e:
                logger.warning(f"[DocumentParser] 读取docling JSON失败: {e}")
        
        try:
            markdown_content = await DocumentParser._parse_with_docling(file_path, filename)
            return ParseResult(
                markdown_content=markdown_content,
                filename=filename,
                file_path=str(file_path),
                file_size=file_size,
                parse_method="docling",
                ai_formatted=ai_formatted,
                doc_title=doc_title
            )
        except Exception as e:
            logger.warning(f"Docling解析失败，使用备用方法: {e}")
            markdown_content = await DocumentParser._fallback_parse(file_path, filename)
            return ParseResult(
                markdown_content=markdown_content,
                filename=filename,
                file_path=str(file_path),
                file_size=file_size,
                parse_method="fallback"
            )
    
    @staticmethod
    async def _parse_with_docling(file_path: Path, filename: str) -> str:
        """
        使用Docling解析文件，返回Markdown内容（增强版，保留完整层级结构）
        对PPTX文件，使用iterate_items提取页码信息，生成带"第N页"标记的Markdown
        """
        try:
            from docling.document_converter import DocumentConverter
            from docling.backend.pypdfium2_backend import PyPdfiumDocumentBackend

            logger.info(f"[Docling] 开始解析: {filename}")

            # 使用增强配置以保留更好的层级结构
            converter = DocumentConverter(
                allowed_formats=["pdf", "docx", "pptx", "html", "image"],
                format_options={
                    "pdf": {
                        "do_ocr": True,
                        "ocr_lang": "zh_cn+en",
                    }
                }
            )

            result = converter.convert(str(file_path))

            is_pptx = filename.lower().endswith(('.pptx', '.ppt'))

            if is_pptx:
                # PPTX文件：使用iterate_items提取页码，生成带页码标记的Markdown
                markdown_content = DocumentParser._build_pptx_markdown_with_pages(result)
            else:
                # 其他格式：使用标准Markdown导出
                markdown_content = result.document.export_to_markdown()

            # 后处理：确保标题层级清晰
            markdown_content = DocumentParser._post_process_markdown(markdown_content, filename)

            logger.info(f"[Docling] 解析完成，生成 {len(markdown_content)} 字符Markdown")
            return markdown_content

        except ImportError:
            raise ImportError("Docling未安装")
        except Exception as e:
            logger.error(f"[Docling] 解析失败: {e}")
            raise Exception(f"Docling解析失败: {e}")

    @staticmethod
    def _build_pptx_markdown_with_pages(result) -> str:
        """
        从Docling的PPTX解析结果中构建带页码标记的Markdown
        使用iterate_items遍历文档项，根据prov中的page_no插入"第N页"标记
        """
        parts = []
        current_page = 0

        try:
            for item, level in result.document.iterate_items():
                # 从provenance中提取页码
                page_no = 0
                if hasattr(item, 'prov') and item.prov:
                    page_no = item.prov[0].page_no + 1  # Docling页码从0开始

                # 页码变化时插入标记
                if page_no > 0 and page_no != current_page:
                    current_page = page_no
                    parts.append(f"\n## 第{current_page}页\n")

                # 提取文本内容
                if hasattr(item, 'text') and item.text:
                    text = item.text.strip()
                    if text:
                        # 根据item类型决定格式
                        label = getattr(item, 'label', None)
                        label_val = label.value if hasattr(label, 'value') else str(label) if label else ''

                        if label_val in ('title', 'section_header', 'header'):
                            parts.append(f"\n### {text}\n")
                        elif label_val in ('list_item',):
                            parts.append(f"- {text}")
                        else:
                            parts.append(text)
        except Exception as e:
            logger.warning(f"[Docling] iterate_items失败，回退到标准导出: {e}")
            return result.document.export_to_markdown()

        content = "\n".join(parts)

        # 如果没有提取到任何页码标记，回退到python-pptx直接解析
        if current_page == 0:
            logger.warning("[Docling] PPTX未提取到页码信息，回退到python-pptx")
            try:
                from pptx import Presentation
                prs = Presentation(str(result.input.path) if hasattr(result, 'input') else '')
                slide_parts = []
                for slide_num, slide in enumerate(prs.slides, 1):
                    title = ""
                    if slide.shapes.title and slide.shapes.title.text:
                        title = slide.shapes.title.text.strip()
                    slide_parts.append(f"\n## 第{slide_num}页{': ' + title if title else ''}\n")
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text = shape.text.strip()
                            if text != title:
                                slide_parts.append(text)
                content = "\n".join(slide_parts)
            except Exception:
                pass

        return content

    @staticmethod
    def _post_process_markdown(content: str, filename: str) -> str:
        """
        后处理Markdown内容：
        1. 清理乱码字符
        2. 确保标题层级规范
        3. 移除空行过多的情况
        4. 标准化特殊符号
        """
        import re

        lines = content.split("\n")
        processed_lines = []
        prev_was_empty = False

        for line in lines:
            # 清理每行的控制字符和乱码
            cleaned = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', '', line)
            cleaned = re.sub(r'[□■◆●★☆♦♠♣♥▪▫◇]', '', cleaned)

            # 处理空行（避免连续多个空行）
            if not cleaned.strip():
                if not prev_was_empty:
                    processed_lines.append("")
                    prev_was_empty = True
                continue

            prev_was_empty = False

            # 确保标题格式规范
            if cleaned.startswith("#"):
                # 标准化标题：移除多余空格
                match = re.match(r'(#+)\s*(.+)', cleaned)
                if match:
                    level = len(match.group(1))
                    title_text = match.group(2).strip()
                    # 清理标题中的乱码
                    title_text = re.sub(r'[^\u4e00-\u9fa5a-zA-Z0-9\s\-\(\)（）【】《》、.。,，:：!！?？]', '', title_text)
                    if title_text and len(title_text) >= 2:
                        cleaned = f"{'#' * level} {title_text}"
                    else:
                        continue  # 跳过无效标题
                else:
                    continue

            processed_lines.append(cleaned)

        # 确保文档有标题
        if processed_lines and not processed_lines[0].startswith("#"):
            processed_lines.insert(0, f"# {filename}")

        result = "\n".join(processed_lines)

        # 最终清理：移除首尾空白
        result = result.strip()

        return result
    
    @staticmethod
    async def _fallback_parse(file_path: Path, filename: str) -> str:
        """
        备用解析方法（当Docling不可用时）
        """
        suffix = file_path.suffix.lower()
        
        if suffix == ".pdf":
            return await DocumentParser._parse_pdf(file_path, filename)
        elif suffix in [".docx", ".doc"]:
            return await DocumentParser._parse_docx(file_path, filename)
        elif suffix in [".pptx", ".ppt"]:
            return await DocumentParser._parse_pptx(file_path, filename)
        elif suffix in [".txt", ".md", ".json", ".py", ".js", ".html", ".css"]:
            return await DocumentParser._parse_text(file_path, filename)
        else:
            return f"# {filename}\n\n[不支持的文件格式: {suffix}]"
    
    @staticmethod
    async def _parse_pdf(file_path: Path, filename: str) -> str:
        """解析PDF文件"""
        try:
            import pdfplumber
            text_parts = [f"# {Path(filename).stem}\n"]
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    text = page.extract_text()
                    if text:
                        text_parts.append(f"\n## 第{page_num}页\n")
                        text_parts.append(f"{text}\n")
            return "\n".join(text_parts)
        except ImportError:
            return f"# {filename}\n\n[PDF文件需要安装pdfplumber库才能解析]"
    
    @staticmethod
    async def _parse_docx(file_path: Path, filename: str) -> str:
        """解析Word文件"""
        try:
            from docx import Document
            doc = Document(file_path)
            text_parts = [f"# {Path(filename).stem}\n"]
            
            for para in doc.paragraphs:
                if para.text.strip():
                    if para.style.name.startswith('Heading'):
                        level = int(para.style.name.replace('Heading ', '')) if para.style.name != 'Heading' else 1
                        text_parts.append(f"\n{'#' * level} {para.text.strip()}\n")
                    else:
                        text_parts.append(f"{para.text.strip()}\n")
            
            return "\n".join(text_parts)
        except ImportError:
            return f"# {filename}\n\n[Word文件需要安装python-docx库才能解析]"
    
    @staticmethod
    async def _parse_pptx(file_path: Path, filename: str) -> str:
        """解析PPT文件"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_parts = [f"# {Path(filename).stem}\n"]
            text_parts.append("\n> 本文档由PPT自动解析生成\n")
            
            for slide_num, slide in enumerate(prs.slides, 1):
                title = ""
                if slide.shapes.title and slide.shapes.title.text:
                    title = slide.shapes.title.text.strip()
                
                text_parts.append(f"\n## 第{slide_num}页{': ' + title if title else ''}\n")
                
                content_items = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        if text != title:
                            content_items.append(text)
                
                for item in content_items:
                    if len(item) < 50 and not item.endswith(('。', '，', '：', '.', ',', ':')):
                        if item.endswith('?') or item.endswith('？'):
                            text_parts.append(f"\n### ❓ {item}\n")
                        else:
                            text_parts.append(f"\n### {item}\n")
                    else:
                        text_parts.append(f"\n{item}\n")
            
            return "\n".join(text_parts)
        except ImportError:
            return f"# {filename}\n\n[PPT文件需要安装python-pptx库才能解析]"
    
    @staticmethod
    async def _parse_text(file_path: Path, filename: str) -> str:
        """解析纯文本文件"""
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()
        suffix = file_path.suffix.lower()
        if suffix == ".md":
            return content
        else:
            return f"# {filename}\n\n```\n{content}\n```"


class StructureParser:
    """结构化解析器 - 增强版"""

    @staticmethod
    def parse_markdown_to_structure(markdown_content: str, filename: str) -> StructureResult:
        """
        将Markdown内容解析为结构化数据
        增强版：正确提取docling的层级结构，生成丰富的教学内容
        """
        lines = markdown_content.split("\n")
        texts = []
        groups = []
        current_group = None
        current_group_idx = 0
        text_idx = 0
        current_page_no = 1  # 跟踪当前页码（从PPTX的"第N页"标题中提取）

        # 清理和预处理文本
        def clean_text(text: str) -> str:
            """清理文本，去除乱码和特殊字符"""
            if not text:
                return text
            # 移除控制字符但保留换行和常见标点
            text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', '', text)
            # 移除特殊方块符号
            text = re.sub(r'[□■◆●★☆♦♠♣♥▪▫◇○●□■]', '', text)
            # 合并多个空格为单个空格
            text = re.sub(r' +', ' ', text)
            # 移除首尾空白
            text = text.strip()
            return text

        for line in lines:
            original_line = line
            line = line.strip()

            if not line:
                continue

            # 清理文本
            cleaned_line = clean_text(line)

            # 检测标题层级（支持多级标题）
            if line.startswith("#"):
                # 计算标题级别
                level = 0
                temp_line = line
                while temp_line.startswith("#") and level < 6:
                    level += 1
                    temp_line = temp_line[1:]

                title_text = clean_text(temp_line.strip())

                # 过滤无效或过短的标题
                if not title_text or len(title_text) < 2:
                    continue

                # 从"第N页"格式的标题中提取页码（PPTX解析产物）
                page_match = re.match(r'第(\d+)页', title_text)
                if page_match:
                    current_page_no = int(page_match.group(1))

                # 确定节点类型标签
                if level == 1:
                    label = "chapter"
                    content_layer = "title"
                elif level == 2:
                    label = "section"
                    content_layer = "heading"
                elif level == 3:
                    label = "subsection"
                    content_layer = "subheading"
                else:
                    label = "paragraph"
                    content_layer = "text"

                group_data = {
                    "self_ref": f"#/groups/{current_group_idx}",
                    "name": title_text,
                    "label": label,
                    "content_layer": content_layer,
                    "level": level,
                    "page_no": current_page_no,
                }
                groups.append(group_data)
                current_group = group_data
                current_group_idx += 1

            else:
                # 普通文本内容 - 收集到当前分组下
                if len(cleaned_line) > 8:  # 过滤太短的内容（至少8个字符才有意义）
                    text_data = {
                        "self_ref": f"#/texts/{text_idx}",
                        "label": "paragraph",
                        "text": cleaned_line,
                        "page_no": current_page_no,
                        "group_id": current_group.get("self_ref") if current_group else None,
                        "group_level": current_group.get("level", 1) if current_group else 1,
                    }
                    texts.append(text_data)
                    text_idx += 1

        # 如果没有检测到任何分组，创建默认分组
        if not groups:
            groups.append({
                "self_ref": "#/groups/0",
                "name": filename,
                "label": "section",
                "content_layer": "body",
                "level": 1,
            })

        logger.info(f"[StructureParser] 解析完成: {len(groups)} 个分组, {len(texts)} 条文本")

        return StructureResult(
            groups=groups,
            texts=texts,
            tables=[],
            pictures=[],
            raw_content=markdown_content
        )


class KnowledgeExtractor:
    """知识点提取器 - 使用LLM从文档中提取结构化知识点"""

    @staticmethod
    async def extract_knowledge_points(
        document_content: str,
        filename: str = "",
        max_content_length: int = 8000
    ) -> str:
        """
        从文档内容中提取知识点，返回结构化Markdown

        Args:
            document_content: 文档内容
            filename: 文件名
            max_content_length: 最大内容长度

        Returns:
            str: 结构化的Markdown知识点文档
        """
        logger.info(f"[KnowledgeExtractor] 开始提取知识点: {filename}")

        # 截断内容以适应模型上下文限制
        if len(document_content) > max_content_length:
            truncated_content = document_content[:max_content_length]
            truncated_content += f"\n\n[内容已截断，原长度: {len(document_content)} 字符]"
        else:
            truncated_content = document_content

        # 使用提示词构建函数
        user_prompt = build_knowledge_extraction_prompt(truncated_content, filename)

        try:
            messages = [
                Message(role="system", content=KNOWLEDGE_EXTRACTION_PROMPT),
                Message(role="user", content=user_prompt)
            ]

            logger.info(f"[KnowledgeExtractor] 发送AI请求，内容长度: {len(user_prompt)} 字符")
            response = await llm_client.chat(messages, temperature=0.3)
            logger.info(f"[KnowledgeExtractor] 收到AI响应，长度: {len(response.content)} 字符")

            return response.content

        except Exception as e:
            logger.error(f"[KnowledgeExtractor] 知识点提取失败: {e}")
            # 返回一个基本的错误提示Markdown
            return f"""# {filename or '文档'} - 知识点提取失败

由于技术原因，无法自动提取知识点。请手动整理文档内容。

**错误信息**: {str(e)}
"""

    @staticmethod
    def parse_knowledge_markdown(markdown_content: str) -> List[Dict[str, Any]]:
        """
        解析知识点Markdown，提取结构化的知识点列表

        Args:
            markdown_content: 知识点Markdown内容

        Returns:
            List[Dict]: 知识点列表，每项包含 id, title, level, content
        """
        knowledge_points = []
        lines = markdown_content.split('\n')

        current_kp = None
        content_buffer = []

        for line in lines:
            line = line.strip()
            if not line:
                continue

            # 检测标题行
            if line.startswith('#'):
                # 保存上一个知识点
                if current_kp:
                    current_kp['content'] = '\n'.join(content_buffer).strip()
                    knowledge_points.append(current_kp)
                    content_buffer = []

                # 解析标题级别
                level = len(line) - len(line.lstrip('#'))
                title = line.lstrip('#').strip()

                current_kp = {
                    'id': f"知识点{len(knowledge_points) + 1}",
                    'title': title,
                    'level': level,
                    'content': ''
                }
            else:
                # 累积内容
                if current_kp:
                    content_buffer.append(line)

        # 保存最后一个知识点
        if current_kp:
            current_kp['content'] = '\n'.join(content_buffer).strip()
            knowledge_points.append(current_kp)

        logger.info(f"[KnowledgeExtractor] 解析出 {len(knowledge_points)} 个知识点")
        return knowledge_points

    @staticmethod
    async def generate_script_from_knowledge(
        knowledge_markdown: str,
        filename: str = "",
        max_content_length: int = 8000
    ) -> Dict[str, Any]:
        """
        从知识点 Markdown 生成结构化智课脚本

        Args:
            knowledge_markdown: 知识点 Markdown 内容
            filename: 文件名
            max_content_length: 最大内容长度

        Returns:
            Dict: 结构化脚本 JSON
        """
        logger.info(f"[KnowledgeExtractor] 开始从知识点生成脚本: {filename}")

        # 截断内容以适应模型上下文限制
        if len(knowledge_markdown) > max_content_length:
            truncated_content = knowledge_markdown[:max_content_length]
            truncated_content += f"\n\n[内容已截断，原长度: {len(knowledge_markdown)} 字符]"
        else:
            truncated_content = knowledge_markdown

        # 使用提示词构建函数
        user_prompt = build_knowledge_to_script_prompt(truncated_content, filename)

        try:
            messages = [
                Message(role="system", content=KNOWLEDGE_TO_SCRIPT_PROMPT),
                Message(role="user", content=user_prompt)
            ]

            logger.info(f"[KnowledgeExtractor] 发送脚本生成请求，内容长度: {len(user_prompt)} 字符")
            response = await llm_client.chat(messages, temperature=0.5)
            logger.info(f"[KnowledgeExtractor] 收到脚本响应，长度: {len(response.content)} 字符")

            # 提取 JSON
            json_match = re.search(r'\{[\s\S]*\}', response.content)
            if json_match:
                script_content = json.loads(json_match.group())
            else:
                # 如果无法解析 JSON，返回默认结构
                script_content = KnowledgeExtractor._create_default_script(filename, knowledge_markdown)

            return script_content

        except Exception as e:
            logger.error(f"[KnowledgeExtractor] 脚本生成失败: {e}")
            return KnowledgeExtractor._create_default_script(filename, knowledge_markdown)

    @staticmethod
    def _create_default_script(filename: str, knowledge_markdown: str) -> Dict[str, Any]:
        """
        创建默认脚本（当AI调用失败时）
        """
        lines = [l.strip() for l in knowledge_markdown.split("\n") if l.strip().startswith('#')]

        sections = []

        # 开场白
        sections.append({
            "type": "opening",
            "id": "sec_000",
            "title": "课程开场",
            "content": f"同学们好！欢迎学习《{Path(filename).stem if filename else '本课程'}》。在今天的课程中，我们将一起探索这个有趣的主题。希望通过今天的学习，大家能够掌握核心概念，并能够灵活运用到实际问题中。",
            "duration": 45,
            "tone": "enthusiastic",
            "transitions": {
                "next": "接下来，让我们进入今天的第一个知识点。"
            }
        })

        # 知识点
        for idx, line in enumerate(lines[:8]):
            if idx == 0:  # 跳过主标题
                continue

            title = line.lstrip('#').strip()
            prev_transition = "首先，" if idx == 1 else ("最后，" if idx == len(lines[:8]) - 1 else "接下来，")
            next_transition = "" if idx == len(lines[:8]) - 1 else "理解了概念后，我们继续往下看。"

            sections.append({
                "type": "knowledge_point",
                "id": f"sec_{idx:03d}",
                "title": title[:40] + ("..." if len(title) > 40 else ""),
                "definition": title,
                "explanation": f"这是关于{title[:20]}的详细解释。",
                "examples": [f"示例：相关应用场景"],
                "content": f"{prev_transition}我们来学习{title[:40]}。{title}这个概念非常重要。",
                "duration": 90,
                "difficulty": "medium",
                "is_key_point": idx % 2 == 0,
                "transitions": {
                    "prev": prev_transition.strip('，'),
                    "next": next_transition
                }
            })

        # 总结语
        sections.append({
            "type": "summary",
            "id": "sec_999",
            "title": "课程总结",
            "key_points": ["核心要点回顾"],
            "content": "今天我们学习了本课程的核心内容，希望大家课后能够复习巩固。",
            "duration": 60,
            "next_preview": "下节课我们将继续深入学习。"
        })

        return {
            "title": Path(filename).stem if filename else "课程",
            "summary": f"本课程《{Path(filename).stem if filename else '课程'}》共包含 {len(sections)} 个教学环节。",
            "keywords": ["知识点", "课程"],
            "total_duration": sum(s.get("duration", 60) for s in sections),
            "sections": sections
        }


class ScriptGenerator:
    """智课脚本生成器"""
    
    @staticmethod
    async def generate_script_from_ai_formatted(
        ai_formatted: List[Dict[str, Any]],
        filename: str,
        doc_title: str = None
    ) -> ScriptResult:
        """
        从ai_formatted数据生成智课脚本
        先提取结构，再使用AI将原始内容转为适合TTS语音讲解的文本
        
        Args:
            ai_formatted: docling生成的ai_formatted数组
            filename: 文件名
            doc_title: 文档标题
            
        Returns:
            ScriptResult: 脚本结果
        """
        logger.info(f"[ScriptGenerator] 从ai_formatted生成脚本: {filename}")
        
        if not ai_formatted:
            return ScriptGenerator._create_empty_script(filename)
        
        merged_sections = []
        pending_title = None
        
        for idx, section in enumerate(ai_formatted):
            title = section.get('title', f'知识点{idx+1}')
            content = section.get('content', '')
            label = section.get('label', 'text')
            
            if label in ['title_page', 'toc_page']:
                if content and len(content.strip()) >= 10:
                    merged_sections.append({
                        'title': title,
                        'content': content,
                        'label': label,
                        'slide': section.get('slide', 1),
                    })
                continue
            
            if not content or len(content.strip()) < 10:
                pending_title = title
                continue
            
            if pending_title:
                if title != pending_title:
                    title = f"{pending_title} - {title}"
                pending_title = None
            
            merged_sections.append({
                'title': title,
                'content': content,
                'label': label,
                'slide': section.get('slide', 1),
            })
        
        if pending_title:
            if merged_sections:
                last = merged_sections[-1]
                last['title'] = f"{last['title']} - {pending_title}"
            else:
                merged_sections.append({
                    'title': pending_title,
                    'content': '本节内容待补充。',
                    'label': 'section_title',
                    'slide': 1,
                })
        
        merged_sections = await ScriptGenerator._analyze_and_filter_sections(merged_sections, doc_title or filename)
        
        merged_sections = await ScriptGenerator._convert_to_tts_content(merged_sections, doc_title or filename)
        
        nodes = []
        total_duration = 0
        
        for idx, section in enumerate(merged_sections):
            node_type = ScriptGenerator._determine_node_type(section['label'], idx, len(merged_sections))
            duration = ScriptGenerator._estimate_duration(section['content'])
            is_key = section.get('is_key_point', idx % 3 == 0)
            
            nodes.append(ScriptNode(
                chapter_id=f"chap_{idx:03d}",
                node_type=node_type,
                title=section['title'],
                content=section['content'],
                page_start=section['slide'] + 1,
                page_end=section['slide'] + 1,
                duration=duration,
                is_key_point=is_key,
            ))
            total_duration += duration
        
        if not nodes:
            return ScriptGenerator._create_empty_script(filename)
        
        title = doc_title or Path(filename).stem
        summary = f"本课程《{title}》共包含 {len(nodes)} 个教学节点，系统讲解核心知识点。"
        keywords = ScriptGenerator._extract_keywords_from_nodes(nodes)
        
        logger.info(f"[ScriptGenerator] 从ai_formatted生成完成: {len(nodes)} 个节点")

        # 计算预估时间戳（基于duration字段按比例分配）
        nodes = ScriptGenerator._estimate_timestamps(nodes)

        return ScriptResult(
            title=title,
            summary=summary,
            keywords=keywords,
            total_duration=total_duration,
            nodes=nodes,
            script_content={
                "title": title,
                "summary": summary,
                "keywords": keywords,
                "total_duration": total_duration,
                "nodes": [
                    {
                        "chapter_id": n.chapter_id,
                        "node_type": n.node_type,
                        "title": n.title,
                        "content": n.content,
                        "duration": n.duration,
                        "is_key_point": n.is_key_point,
                    }
                    for n in nodes
                ]
            },
            beautiful_markdown=""
        )
    
    @staticmethod
    async def _analyze_and_filter_sections(
        sections: List[Dict[str, Any]],
        doc_title: str
    ) -> List[Dict[str, Any]]:
        """
        使用AI分析课件section列表，过滤非知识内容，合并相关section，
        识别知识点层级关系，标记重点节点
        """
        if not sections or len(sections) <= 2:
            return sections
        
        sections_summary = ""
        for idx, section in enumerate(sections):
            title = section.get('title', f'节点{idx+1}')
            content = section.get('content', '')[:100]
            sections_summary += f"\n[{idx+1}] 标题：{title}\n    内容：{content}...\n"
        
        system_prompt = """你是一位专业的课程结构分析专家。你的任务是分析PPT课件提取的节点列表，识别哪些是真正的知识点，哪些是非知识内容。

## 需要过滤的非知识内容（标记为 "skip"）
1. 教室/地点信息（如"工程综合训练中心巡天405A"）
2. 学时/课时信息（如"40学时（理论课36学时+实践课4学时）"）
3. 考核/评分标准（如"总成绩 = 平时成绩(30%)+期末考试"）
4. 文件名/课程代码（如"1 电路概念定律.pptx"）
5. 纯目录页（无实质内容的章节列表）
6. 重复的章节标题（与前一节点内容高度重复）

## 需要保留的知识内容
1. 具体的概念定义和原理解释
2. 公式推导和计算方法
3. 案例分析和应用示例
4. 章节概览（包含学习目标的）
5. 习题和问答内容

## 需要合并的情况
- 如果相邻的2-3个section内容都很短（<50字），且属于同一主题，合并为一个节点
- 章节标题 + 紧跟的概述内容 → 合并

## 重点标记（is_key_point）
- 核心概念/定义 → true
- 重要公式/定理 → true
- 一般性描述 → false
- 习题/问答 → false

## 输出格式
严格按JSON数组格式输出，每个元素对应一个处理后的节点：

```json
[
  {"action": "skip", "reason": "教室地点信息"},
  {"action": "skip", "reason": "学时信息"},
  {"action": "skip", "reason": "评分标准"},
  {"action": "keep", "is_key_point": true, "suggested_title": "电路的基本概念"},
  {"action": "merge", "merge_with": 4, "is_key_point": true, "suggested_title": "电路的作用与组成部分"},
  {"action": "keep", "is_key_point": false, "suggested_title": "电压和电流的参考方向"}
]
```

注意：
- 数组长度必须与输入节点数量完全一致
- merge_with 表示合并到第几个节点（从1开始计数）
- suggested_title 是优化后的节点标题
- 只输出JSON数组，不要其他内容"""

        user_prompt = f"""课程：{doc_title}

以下是{len(sections)}个课件节点的标题和内容摘要：

{sections_summary}

请分析每个节点，输出处理建议的JSON数组："""

        try:
            messages = [
                Message(role="system", content=system_prompt),
                Message(role="user", content=user_prompt)
            ]
            
            logger.info(f"[ScriptGenerator] 知识结构分析: 发送{len(sections)}个节点到AI")
            response = await llm_client.chat(messages, temperature=0.3, max_tokens=4096)
            logger.info(f"[ScriptGenerator] 知识结构分析: 收到AI响应")
            
            import json as json_module
            import re
            
            content = response.content.strip()
            json_match = re.search(r'\[.*\]', content, re.DOTALL)
            if not json_match:
                logger.warning("[ScriptGenerator] 知识结构分析: AI返回非JSON格式，跳过过滤")
                return sections
            
            actions = json_module.loads(json_match.group())
            
            if len(actions) != len(sections):
                logger.warning(
                    f"[ScriptGenerator] 知识结构分析: AI返回{len(actions)}项，期望{len(sections)}项，跳过过滤"
                )
                return sections
            
            result = []
            merge_targets = {}
            
            for idx, action in enumerate(actions):
                act = action.get('action', 'keep')
                
                if act == 'skip':
                    logger.info(f"[ScriptGenerator] 过滤节点[{idx+1}]: {sections[idx]['title']} - {action.get('reason', '')}")
                    continue
                
                section = sections[idx].copy()
                
                suggested_title = action.get('suggested_title')
                if suggested_title and len(suggested_title.strip()) > 0:
                    section['title'] = suggested_title
                
                section['is_key_point'] = action.get('is_key_point', False)
                
                if act == 'merge':
                    merge_to = action.get('merge_with', idx + 1) - 1
                    if 0 <= merge_to < len(sections):
                        if merge_to not in merge_targets:
                            merge_targets[merge_to] = []
                        merge_targets[merge_to].append(section)
                        continue
                
                result.append(section)
            
            for target_idx, merged_sections in merge_targets.items():
                for r_section in result:
                    target_title = sections[target_idx]['title']
                    if r_section['title'] == target_title or target_title in r_section['title']:
                        merged_content = "\n".join([ms['content'] for ms in merged_sections if ms.get('content')])
                        if merged_content:
                            r_section['content'] = r_section.get('content', '') + "\n" + merged_content
                        break
            
            if not result:
                logger.warning("[ScriptGenerator] 知识结构分析: 过滤后无节点，保留原始数据")
                return sections
            
            logger.info(f"[ScriptGenerator] 知识结构分析: {len(sections)} → {len(result)} 个节点（过滤{len(sections)-len(result)}个）")
            return result
            
        except Exception as e:
            logger.warning(f"[ScriptGenerator] 知识结构分析失败: {e}，跳过过滤")
            return sections

    @staticmethod
    async def _convert_to_tts_content(
        sections: List[Dict[str, Any]],
        doc_title: str
    ) -> List[Dict[str, Any]]:
        """
        使用AI批量将原始课件内容转为适合TTS语音讲解的文本
        将所有section的内容一次性发送给AI处理，减少API调用次数
        """
        if not sections:
            return sections
        
        # 构建待转换的内容摘要
        sections_text = ""
        for idx, section in enumerate(sections):
            sections_text += f"\n---\n[节点{idx+1}] 标题：{section['title']}\n内容：{section['content']}\n"
        
        system_prompt = """你是一位专业的课程讲解稿撰写专家。你的任务是将PPT课件中提取的原始文本改写为适合TTS语音播报的讲解稿。

## 原始文本的问题
课件提取的原始文本通常存在以下问题，你必须全部修正：
1. 包含文件名（如"金材-第二章 1.pptx"）—— 必须去除，替换为实际的知识点名称
2. 包含模板标记（如【概念定义】【原理解释】）—— 必须去除，用自然过渡语替代
3. 内容空洞泛化（如"它涉及多个要素的相互作用"）—— 必须根据标题和上下文补充具体的专业内容
4. 编号符号（如①②③、1. 2. 3.）—— 改为自然语言过渡
5. Markdown格式（如#、**、-）—— 去除所有格式标记
6. 图片占位符（如[图片]）—— 去除

## 讲解稿写作规范
1. **口语化表达**：像老师在课堂上讲课一样自然，使用"我们来看""接下来""需要注意的是""简单来说"等过渡语
2. **内容充实**：如果原文内容空洞，根据标题补充合理的专业知识和解释，不要写空话
3. **逻辑清晰**：按"引入→定义→解释→举例→小结"的结构组织
4. **长度适中**：每个节点300-800字，信息密度高但不冗长
5. **直接输出**：只输出讲解稿纯文本，不要任何额外标记或说明

## 输出格式
严格按以下格式输出每个节点的讲解稿，节点之间用 === 分隔：

[节点1的讲解稿文本]
===
[节点2的讲解稿文本]
===
[节点3的讲解稿文本]

注意：节点数量和顺序必须与输入完全一致，不要遗漏或合并任何节点。"""

        user_prompt = f"""课程：{doc_title}

以下是{len(sections)}个课件节点的原始内容，请改写为适合语音播报的讲解稿：

{sections_text}

请按格式输出{len(sections)}个节点的讲解稿，用 === 分隔："""

        try:
            messages = [
                Message(role="system", content=system_prompt),
                Message(role="user", content=user_prompt)
            ]
            
            logger.info(f"[ScriptGenerator] TTS转换: 发送{len(sections)}个节点到AI")
            response = await llm_client.chat(messages, temperature=0.7, max_tokens=8192)
            logger.info(f"[ScriptGenerator] TTS转换: 收到AI响应，长度{len(response.content)}字符")
            
            # 解析AI返回的内容
            tts_texts = response.content.strip().split('===')
            tts_texts = [t.strip() for t in tts_texts if t.strip()]
            
            # 如果返回的节点数量不匹配，逐个增强
            if len(tts_texts) != len(sections):
                logger.warning(
                    f"[ScriptGenerator] TTS转换: AI返回{len(tts_texts)}个节点，期望{len(sections)}个，回退到逐个增强"
                )
                for section in sections:
                    section['content'] = await ScriptGenerator._enhance_content_with_ai(
                        section['title'], section['content'], 0
                    )
                return sections
            
            # 将AI生成的讲解稿替换到sections中
            for idx, tts_text in enumerate(tts_texts):
                # 去除可能残留的[节点N]标记
                cleaned = re.sub(r'^\[节点\d+\]\s*', '', tts_text)
                if len(cleaned) >= 100:
                    sections[idx]['content'] = cleaned
                else:
                    logger.warning(f"[ScriptGenerator] TTS转换: 节点{idx+1}内容过短({len(cleaned)}字)，保留原始内容")
            
            logger.info(f"[ScriptGenerator] TTS转换: 成功转换{len(sections)}个节点")
            return sections
            
        except Exception as e:
            logger.warning(f"[ScriptGenerator] TTS转换失败: {e}，使用逐个增强")
            for section in sections:
                section['content'] = await ScriptGenerator._enhance_content_with_ai(
                    section['title'], section['content'], 0
                )
            return sections
    
    @staticmethod
    async def _enhance_content_with_ai(
        title: str,
        original_content: str,
        section_index: int
    ) -> str:
        """
        使用AI增强内容（只优化content，不改变标题）
        """
        if len(original_content) >= 300 and len(original_content) <= 800:
            return original_content
        
        if len(original_content) < 50:
            return ScriptGenerator._expand_content(original_content, section_index, "lecture")
        
        try:
            system_prompt = """你是一位专业的课程讲解稿撰写专家。请将以下课件内容改写为适合语音播报的讲解稿。

## 核心要求
1. 语言风格：口语化、自然流畅，像老师在课堂上讲课一样
2. 去除所有不适合朗读的内容：文件名（如xxx.pptx）、模板标记（如【概念定义】）、编号符号（如①②③）
3. 用自然的过渡语连接内容，如"首先""接下来""需要注意的是""简单来说"
4. 适当补充解释和举例，让抽象概念更易懂
5. 内容长度控制在300-800字之间
6. 直接输出讲解稿文本，不要添加任何额外说明或格式标记"""

            user_prompt = f"""标题：{title}

原始课件内容：
{original_content}

请改写为适合语音播报的讲解稿："""

            messages = [
                Message(role="system", content=system_prompt),
                Message(role="user", content=user_prompt)
            ]
            
            response = await llm_client.chat(messages)
            enhanced = response.content.strip()
            
            if len(enhanced) < 100:
                return original_content
            
            return enhanced
            
        except Exception as e:
            logger.warning(f"[ScriptGenerator] AI增强失败: {e}，使用原始内容")
            return original_content
    
    @staticmethod
    def _determine_node_type(label: str, index: int, total: int) -> str:
        """根据标签和位置确定节点类型"""
        if label == 'section_title':
            return 'lecture'
        elif label == 'heading':
            return 'lecture'
        elif index == total - 1:
            return 'summary'
        elif index % 4 == 3:
            return 'question'
        else:
            return 'lecture'
    
    @staticmethod
    def _estimate_duration(content: str) -> int:
        """根据内容长度估算时长（秒）"""
        char_count = len(content)
        if char_count < 100:
            return 60
        elif char_count < 200:
            return 90
        elif char_count < 300:
            return 120
        else:
            return min(180, char_count // 2)
    
    @staticmethod
    def _estimate_timestamps(nodes: List[ScriptNode]) -> List[ScriptNode]:
        """
        根据节点的duration字段预估时间戳
        按顺序累加计算每个节点的timestamp_start和timestamp_end
        
        Args:
            nodes: 脚本节点列表
            
        Returns:
            更新了时间戳的节点列表
        """
        current_time = 0.0
        for node in nodes:
            node.timestamp_start = round(current_time, 2)
            current_time += node.duration
            node.timestamp_end = round(current_time, 2)
        
        logger.info(f"[Timestamp] 已为 {len(nodes)} 个节点预估时间戳，总时长: {current_time:.1f}秒")
        return nodes
    
    @staticmethod
    def _extract_keywords_from_nodes(nodes: List[ScriptNode]) -> List[str]:
        """从节点中提取关键词"""
        keywords = set()
        for node in nodes[:5]:
            title_words = [w for w in node.title if len(w) > 1]
            keywords.update(title_words[:3])
        return list(keywords)[:5] if keywords else ["知识点", "课程"]
    
    @staticmethod
    def _create_empty_script(filename: str) -> ScriptResult:
        """创建空脚本"""
        return ScriptResult(
            title=Path(filename).stem,
            summary="课程内容解析失败",
            keywords=["知识点"],
            total_duration=0,
            nodes=[],
            script_content={},
            beautiful_markdown=""
        )
    
    @staticmethod
    async def generate_script(
        markdown_content: str,
        filename: str,
        max_content_length: int = 20000
    ) -> ScriptResult:
        """
        调用AI生成智课脚本
        
        Args:
            markdown_content: Markdown内容
            filename: 文件名
            max_content_length: 最大内容长度
            
        Returns:
            ScriptResult: 脚本结果
        """
        logger.info(f"[ScriptGenerator] 开始生成脚本: {filename}")
        
        if len(markdown_content) > max_content_length:
            truncated_content = markdown_content[:max_content_length]
            truncated_content += f"\n\n[内容已截断，原长度: {len(markdown_content)} 字符]"
        else:
            truncated_content = markdown_content
        
        system_prompt = """你是一位专业的课程设计师和讲解稿撰写专家。请根据用户提供的文档内容，生成一份**高质量的智课脚本**，内容将用于TTS语音合成和数字人视频播报。

## 核心要求

### 1. 内容长度要求（严格）
- **每个节点的content字段必须在300-800字之间（中文）**
- 不足300字的节点必须补充更多细节和说明
- 超过800字的节点需要适当精简但保留核心要点
- 统计字数时不包括标点符号和空格
- **生成的内容总文本量必须大于原文档的文本量**，通过补充解释、举例、类比等方式扩展内容

### 2. TTS语音播报要求（最重要）
content字段将被TTS语音合成引擎朗读，必须符合以下规范：
- **口语化表达**：像老师在课堂上讲课一样自然，不要写成书面论文风格
- **去除不适合朗读的内容**：文件名（如xxx.pptx）、模板标记（如【概念定义】【原理解释】）、编号符号（如①②③）
- **自然过渡语**：使用"我们来看""接下来""需要注意的是""简单来说""举个例子"等口语化过渡
- **避免空洞泛化**：不要写"它涉及多个要素的相互作用"这类空话，要写具体内容
- **纯文本输出**：content中不要包含Markdown格式（#、**、-）、特殊符号（□■◆●）、图片占位符

### 3. 结构化要求
- **严格按照文档的层级结构组织内容**
- 保持原文档的章节顺序和逻辑关系
- **文档中的每个章节、小节都必须对应至少一个教学节点，不得遗漏任何章节**
- 如果原文档有N个章节/小节标题，生成的节点数量不得少于N个
- 标题必须清晰、准确、无乱码，使用规范的中文表述

### 4. 内容质量标准
- **概念解释**：定义清晰，通俗易懂，避免过于学术化的表述，用类比帮助理解
- **原理讲解**：深入浅出，包含"是什么→为什么→怎么用"的完整逻辑链，详细阐述内在机制
- **实例说明**：至少提供2-3个贴近实际的具体案例，包含详细的分析过程
- **应用场景**：说明知识点的实际应用价值和适用范围，给出具体应用实例
- **过渡衔接**：自然流畅，使用口语化的过渡语
- **内容扩展**：在原文基础上补充背景知识、相关概念对比、常见误区、实际工程应用等

### 5. 标题规范
- 使用简洁明了的中文标题（4-15个字）
- 避免使用特殊符号、乱码字符
- 标题要能准确概括该节点的核心内容
- 禁止出现：□、■、◆、●等乱码符号

## 输出格式

请以JSON格式返回，结构如下：
{
    "title": "课程标题（使用清晰的中文，无乱码）",
    "summary": "课程整体摘要（100-200字，口语化）",
    "keywords": ["关键词1", "关键词2"],
    "total_duration": 总时长(秒),
    "nodes": [
        {
            "chapter_id": "chap_000",
            "node_type": "lecture",
            "title": "节点标题（中文，4-15字，无特殊符号）",
            "content": "口语化的讲解稿文本，300-800字，适合TTS朗读",
            "page_start": 1,
            "page_end": 1,
            "duration": 90,
            "is_key_point": false
        }
    ]
}

## 节点类型说明
- lecture: 知识点讲解（主要内容）
- question: 互动提问（穿插在知识点之间）
- summary: 章节总结（每个大章节结束时）
- interactive: 交互环节（引导思考）

## 重要提示
1. **content字段是唯一最重要的字段**，必须确保300-800字的口语化讲解稿
2. **标题必须使用纯中文**，禁止任何乱码或特殊符号
3. **保持文档原有的层级结构**，不要随意重组或打乱顺序
4. **内容要实用且丰富**，学生能够直接理解和应用
5. **不得遗漏原文档中的任何章节**，每个章节都必须有对应的节点
6. **通过补充解释、举例、类比等方式，使生成内容的文本量超过原文档**
7. 返回的必须是有效的JSON格式"""

        user_prompt = f"""请根据以下文档内容生成高质量的智课脚本：

文件名: {filename}

文档Markdown内容：
{truncated_content}

生成要求：
1. 分析文档的层级结构（# ## ### 标题）
2. 按照原有章节顺序拆分为教学节点，**每个章节/小节都必须有对应节点**
3. 每个节点的content严格控制在300-800字
4. 标题使用清晰中文，无乱码
5. 内容要丰富详实，包含足够的细节、例子和解释
6. **生成内容的总文本量必须大于原文档**

请立即生成JSON格式的智课脚本："""

        try:
            messages = [
                Message(role="system", content=system_prompt),
                Message(role="user", content=user_prompt)
            ]
            
            logger.info(f"[ScriptGenerator] 发送AI请求，内容长度: {len(user_prompt)} 字符")
            response = await llm_client.chat(messages)
            logger.info(f"[ScriptGenerator] 收到AI响应，长度: {len(response.content)} 字符")
            
            json_match = re.search(r'\{[\s\S]*\}', response.content)
            if json_match:
                script_content = json.loads(json_match.group())
            else:
                script_content = ScriptGenerator._create_default_script(filename, markdown_content)
            
        except Exception as e:
            logger.warning(f"[ScriptGenerator] AI调用失败: {e}，使用默认脚本")
            script_content = ScriptGenerator._create_default_script(filename, markdown_content)
        
        return ScriptGenerator._build_script_result(script_content, filename, markdown_content)
    
    @staticmethod
    def _create_default_script(filename: str, content: str) -> dict:
        """
        创建默认脚本（当AI调用失败时）
        确保每个节点内容在150-300字之间
        """
        lines = [l.strip() for l in content.split("\n") if l.strip() and len(l.strip()) > 10]

        nodes = []

        for idx, line in enumerate(lines[:12]):
            node_type = "lecture"
            if idx == len(lines[:12]) - 1:
                node_type = "summary"
            elif idx % 4 == 3:
                node_type = "question"

            # 生成150-300字的详细内容
            base_content = line
            expanded_content = ScriptGenerator._expand_content(base_content, idx, node_type)

            nodes.append({
                "chapter_id": f"chap_{idx:03d}",
                "node_type": node_type,
                "title": ScriptGenerator._clean_title(line[:30]),
                "content": expanded_content,
                "page_start": 1,
                "page_end": 1,
                "duration": 90,
                "is_key_point": idx % 3 == 0,
            })

        if not nodes:
            default_content = f"欢迎学习《{Path(filename).stem}》课程。本课程将带领大家系统地学习和掌握相关知识点。通过本课程的学习，你将能够理解核心概念、掌握基本方法，并能够将所学知识应用到实际问题中。让我们开始这段学习之旅吧！"
            nodes = [{
                "chapter_id": "chap_000",
                "node_type": "lecture",
                "title": "课程导入与学习目标",
                "content": default_content,
                "page_start": 1,
                "page_end": 1,
                "duration": 90,
                "is_key_point": True,
            }]

        return {
            "title": Path(filename).stem,
            "summary": f"本课程《{Path(filename).stem}》共包含 {len(nodes)} 个教学节点，系统讲解核心知识点。",
            "keywords": ["知识点", "课程", Path(filename).stem],
            "total_duration": sum(n["duration"] for n in nodes),
            "nodes": nodes,
        }

    @staticmethod
    def _expand_content(base_text: str, index: int, node_type: str) -> str:
        """
        将基础文本扩展为150-300字的详细内容
        """
        templates = {
            "lecture": f"""首先，我们来学习"{base_text}"这个重要的知识点。

【概念定义】
{base_text}是指在特定领域或情境下具有特定含义和作用的核心概念。它是理解和掌握相关知识体系的基础。

【原理解释】
从本质上讲，{base_text}的形成和发展遵循一定的规律和机制。它涉及多个要素的相互作用，包括理论基础、实践应用以及发展趋势等方面。深入理解这一概念，需要我们从多个角度进行分析。

【实际案例】
例如，在实际应用中，{base_text}常常体现在以下几个方面：第一，它可以用来解决具体问题；第二，它为后续学习奠定基础；第三，它与其他知识点存在密切联系。

【学习要点】
掌握{base_text}的关键在于：理解其核心内涵，熟悉其应用场景，并能够举一反三。建议结合实例进行练习，加深理解。""",

            "question": f"""【思考时间】关于"{base_text}"的互动问答

在学习了前面的内容后，请大家思考以下问题：

问题：如何理解{base_text}的核心要点？在实际中应该如何运用？

【提示方向】
1. 从定义出发，思考{base_text}的基本特征
2. 考虑它在不同场景下的表现形式
3. 思考它与其他知识点的关联

请结合自己的理解进行回答，这将有助于检验你对知识点的掌握程度。""",

            "summary": f"""【本节小结】{base_text}

通过本节的学习，我们系统地了解了{base_text}的相关内容。主要收获包括：

**核心要点回顾**
1. 掌握了{base_text}的基本概念和内涵
2. 理解了其原理机制和应用场景
3. 通过实例加深了对知识的理解

**知识框架梳理**
{base_text}作为重要知识点，在整个知识体系中起着承上启下的作用。它既是对前面内容的深化，也为后续学习打下基础。

**学习建议**
建议大家课后及时复习，并结合练习题巩固所学内容。如有疑问，可以随时提问讨论。"""
        }

        content = templates.get(node_type, templates["lecture"])

        # 确保字数在150-300字之间（中文）
        char_count = len(content.replace(" ", "").replace("\n", ""))
        if char_count < 150:
            content += "\n\n此外，还需要注意的是，这一知识点在实际应用中具有重要的价值和意义。通过不断的学习和实践，我们可以更好地掌握和运用它。"
        elif char_count > 300:
            sentences = content.split("。")
            content = "。".join(sentences[:len(sentences)//2]) + "。"

        return content

    @staticmethod
    def _clean_title(title: str) -> str:
        """清理标题，去除乱码和特殊字符"""
        import re
        title = re.sub(r'[□■◆●★☆♦♠♣♥▪▫◇○●□■]', '', title)
        title = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', '', title)
        title = title.strip()
        if not title or len(title) < 2:
            title = "知识点详解"
        elif len(title) > 20:
            title = title[:20]
        return title
    
    @staticmethod
    def _build_script_result(script_content: dict, filename: str, markdown_content: str) -> ScriptResult:
        """构建脚本结果对象"""
        summary_text = script_content.get("summary", f"本课程《{Path(filename).stem}》包含 {len(script_content.get('nodes', []))} 个知识点。")
        keywords = script_content.get("keywords", ["知识点", "课程", Path(filename).stem])
        
        nodes = []
        for node_data in script_content.get("nodes", []):
            nodes.append(ScriptNode(
                chapter_id=node_data.get("chapter_id", "chap_000"),
                node_type=node_data.get("node_type", "lecture"),
                title=node_data.get("title", "未命名节点"),
                content=node_data.get("content", ""),
                page_start=node_data.get("page_start", 1),
                page_end=node_data.get("page_end", 1),
                duration=node_data.get("duration", 60),
                is_key_point=node_data.get("is_key_point", False),
            ))
        
        logger.info(f"[ScriptGenerator] 生成完成: {len(nodes)} 个节点")
        
        return ScriptResult(
            title=script_content.get("title", Path(filename).stem),
            summary=summary_text,
            keywords=keywords,
            total_duration=script_content.get("total_duration", sum(n.duration for n in nodes)),
            nodes=nodes,
            script_content=script_content,
            beautiful_markdown=markdown_content
        )


class MindMapGenerator:
    """思维导图生成器"""
    
    @staticmethod
    def generate(script_result: ScriptResult) -> Dict[str, Any]:
        """
        根据脚本内容生成思维导图JSON结构
        """
        nodes = script_result.nodes
        title = script_result.title
        keywords = script_result.keywords
        
        children = []
        
        for node in nodes:
            child = {
                "text": node.title,
            }
            if node.is_key_point:
                child["highlight"] = True
            children.append(child)
        
        if not children:
            for kw in keywords[:5]:
                children.append({"text": kw})
        
        if not children:
            children = [
                {"text": "知识点1"},
                {"text": "知识点2"},
                {"text": "知识点3"},
            ]
        
        return {
            "text": title,
            "children": children
        }
    
    @staticmethod
    def generate_from_ai_formatted(
        ai_formatted: List[Dict[str, Any]],
        doc_title: str = None
    ) -> Dict[str, Any]:
        """
        从ai_formatted数据生成层级化思维导图JSON结构
        保持docling JSON的原始层级结构
        """
        if not ai_formatted:
            return {"text": doc_title or "文档", "children": []}
        
        root_title = doc_title or "文档"
        root = {"text": root_title, "children": []}
        
        toc_ended = False
        content_sections = []
        
        for section in ai_formatted:
            title = section.get('title', '').strip()
            content = section.get('content', '').strip()
            label = section.get('label', 'text')
            
            if not title or len(title) < 2:
                continue
            
            if label in ['title_page']:
                if len(title) >= 2:
                    root_title = title
                    root["text"] = root_title
                continue
            
            if label in ['toc_page']:
                continue
            
            has_content = bool(content) and len(content) > 10
            
            if not toc_ended:
                if label == 'section_title' and not has_content:
                    continue
                else:
                    toc_ended = True
            
            content_sections.append({
                "title": title,
                "has_content": has_content,
                "label": label,
            })
        
        MindMapGenerator._build_tree(root, content_sections)
        MindMapGenerator._sort_children(root)
        MindMapGenerator._clean_empty_children(root)
        return root
    
    @staticmethod
    def _detect_level(title: str, label: str) -> int:
        if re.match(r'^\d+\.\d+\.\d+', title):
            return 3
        if re.match(r'^\d+\.\d+', title):
            return 2
        if re.match(r'^第[一二三四五六七八九十\d]+[章节]', title):
            return 1
        if label == 'heading':
            return 3
        if label == 'section_title':
            return 2
        return 3
    
    @staticmethod
    def _extract_section_number(title: str):
        m = re.match(r'^(\d+\.\d+(?:\.\d+)?)', title)
        return m.group(1) if m else None
    
    @staticmethod
    def _build_tree(root: Dict, sections: List[Dict]) -> None:
        stack = [root]
        stack_levels = [0]
        stack_numbers = [None]
        
        for sec in sections:
            title = sec["title"]
            level = MindMapGenerator._detect_level(title, sec["label"])
            sec_num = MindMapGenerator._extract_section_number(title)
            
            while len(stack) > 1 and stack_levels[-1] >= level:
                stack.pop()
                stack_levels.pop()
                stack_numbers.pop()
            
            if sec_num and level >= 2:
                parent_num = sec_num.rsplit('.', 1)[0] if '.' in sec_num else None
                if parent_num:
                    for i in range(len(stack) - 1, 0, -1):
                        if stack_numbers[i] == parent_num:
                            while len(stack) > i + 1:
                                stack.pop()
                                stack_levels.pop()
                                stack_numbers.pop()
                            break
            
            parent = stack[-1]
            if "children" not in parent:
                parent["children"] = []
            
            existing = None
            if sec_num:
                for child in parent["children"]:
                    child_num = MindMapGenerator._extract_section_number(child.get("text", ""))
                    if child_num and child_num == sec_num:
                        existing = child
                        break
            
            if existing:
                if sec["has_content"] and not existing.get("has_content"):
                    existing["has_content"] = True
                stack.append(existing)
                stack_levels.append(level)
                stack_numbers.append(sec_num)
            else:
                node = {"text": title}
                if sec["has_content"]:
                    node["has_content"] = True
                parent["children"].append(node)
                stack.append(node)
                stack_levels.append(level)
                stack_numbers.append(sec_num)
    
    @staticmethod
    def _section_sort_key(title: str) -> tuple:
        num = MindMapGenerator._extract_section_number(title)
        if num:
            parts = num.split('.')
            return tuple(int(p) for p in parts)
        return (999,)
    
    @staticmethod
    def _sort_children(node: Dict[str, Any]):
        if "children" not in node:
            return
        node["children"].sort(key=lambda c: MindMapGenerator._section_sort_key(c.get("text", "")))
        for child in node["children"]:
            MindMapGenerator._sort_children(child)
    
    @staticmethod
    def _clean_empty_children(node: Dict[str, Any]):
        """递归清理空的children数组"""
        if "children" in node and not node["children"]:
            del node["children"]
        elif "children" in node:
            for child in node["children"]:
                MindMapGenerator._clean_empty_children(child)


class RAGProcessor:
    """RAG预处理器"""
    
    @staticmethod
    def process(
        markdown_content: str,
        doc_name: str,
        doc_id: str
    ) -> RAGProcessResult:
        """
        执行RAG预处理流水线
        
        Args:
            markdown_content: Markdown内容
            doc_name: 文档名称
            doc_id: 文档ID
            
        Returns:
            RAGProcessResult: RAG处理结果
        """
        logger.info(f"[RAGProcessor] 开始RAG预处理: {doc_name}")
        
        try:
            rag_result = rag_pipeline.process_document(
                markdown_text=markdown_content,
                doc_name=doc_name,
                doc_id=doc_id,
            )
            
            knowledge_points = RAGProcessor._extract_knowledge_points(rag_result)
            
            result = RAGProcessResult(
                formula_count=rag_result.formula_result.formula_count,
                table_count=len(rag_result.table_results),
                domain_term_count=rag_result.tokenize_result.domain_term_count,
                tree_node_count=rag_result.tree_result.total_nodes,
                processed_text=rag_result.processed_text,
                knowledge_points=knowledge_points,
            )
            
            logger.info(f"[RAGProcessor] 预处理完成:")
            logger.info(f"  - 公式: {result.formula_count}个")
            logger.info(f"  - 表格: {result.table_count}个")
            logger.info(f"  - 领域术语: {result.domain_term_count}个")
            logger.info(f"  - 知识树节点: {result.tree_node_count}个")
            logger.info(f"  - 知识点: {len(knowledge_points)}个")
            
            return result
            
        except Exception as e:
            logger.error(f"[RAGProcessor] 预处理失败: {e}")
            return RAGProcessResult(error=str(e))
    
    @staticmethod
    def _extract_knowledge_points(rag_result) -> List[Dict[str, Any]]:
        """
        从RAG结果中提取知识点列表
        """
        knowledge_points = []
        
        if hasattr(rag_result, 'tree_result') and rag_result.tree_result:
            tree = rag_result.tree_result.tree
            if tree:
                RAGProcessor._traverse_tree(tree, knowledge_points)
        
        return knowledge_points
    
    @staticmethod
    def _traverse_tree(node, knowledge_points: List[Dict[str, Any]], path: str = ""):
        """
        遍历知识树，提取知识点
        """
        if node is None:
            return
        
        current_path = f"{path}/{node.title}" if path else node.title
        
        if node.content and len(node.content.strip()) > 20:
            knowledge_points.append({
                "id": f"知识点{len(knowledge_points) + 1}",
                "title": node.title,
                "content": node.content[:500],
                "path": current_path,
                "level": node.level,
            })
        
        for child in node.children:
            RAGProcessor._traverse_tree(child, knowledge_points, current_path)


class DocumentService:
    """文档处理服务"""

    def __init__(self):
        self.parser = DocumentParser()
        self.structure_parser = StructureParser()
        self.knowledge_extractor = KnowledgeExtractor()
        self.script_generator = ScriptGenerator()
        self.mind_map_generator = MindMapGenerator()
        self.rag_processor = RAGProcessor()
    
    async def process_document(
        self,
        file_path: Path,
        filename: str,
        enable_rag: bool = True,
        enable_script: bool = True
    ) -> DocumentProcessResult:
        """
        完整的文档处理流程
        
        Args:
            file_path: 文件路径
            filename: 文件名
            enable_rag: 是否启用RAG预处理
            enable_script: 是否生成智课脚本
            
        Returns:
            DocumentProcessResult: 完整处理结果
        """
        logger.info(f"[DocumentService] 开始处理文档: {filename}")
        
        parse_result = await self.parser.parse_file(file_path, filename)
        
        structure_result = self.structure_parser.parse_markdown_to_structure(
            parse_result.markdown_content, filename
        )
        
        if enable_script:
            if parse_result.ai_formatted:
                logger.info(f"[DocumentService] 使用ai_formatted数据生成脚本")
                script_result = await self.script_generator.generate_script_from_ai_formatted(
                    parse_result.ai_formatted,
                    filename,
                    parse_result.doc_title
                )
            else:
                logger.info(f"[DocumentService] 使用markdown数据生成脚本")
                script_result = await self.script_generator.generate_script(
                    parse_result.markdown_content, filename
                )
        else:
            script_result = ScriptResult(
                title=Path(filename).stem,
                summary="",
                keywords=[],
                total_duration=0,
                nodes=[],
                script_content={},
                beautiful_markdown=parse_result.markdown_content
            )
        
        if enable_rag:
            rag_result = self.rag_processor.process(
                parse_result.markdown_content,
                filename,
                str(hash(file_path))
            )
        else:
            rag_result = RAGProcessResult()
        
        if parse_result.ai_formatted:
            logger.info(f"[DocumentService] 使用ai_formatted数据生成思维导图")
            mind_map = self.mind_map_generator.generate_from_ai_formatted(
                parse_result.ai_formatted,
                parse_result.doc_title
            )
        else:
            mind_map = self.mind_map_generator.generate(script_result)
        
        logger.info(f"[DocumentService] 文档处理完成: {filename}")
        
        return DocumentProcessResult(
            parse_result=parse_result,
            structure_result=structure_result,
            script_result=script_result,
            rag_result=rag_result,
            mind_map=mind_map
        )
    
    async def parse_only(self, file_path: Path, filename: str) -> ParseResult:
        """
        仅解析文件，不进行其他处理
        """
        return await self.parser.parse_file(file_path, filename)

    async def extract_knowledge_only(
        self,
        markdown_content: str,
        filename: str = ""
    ) -> Dict[str, Any]:
        """
        仅提取知识点

        Args:
            markdown_content: Markdown内容
            filename: 文件名

        Returns:
            Dict: 包含 knowledge_markdown 和 knowledge_points 的结果
        """
        knowledge_markdown = await self.knowledge_extractor.extract_knowledge_points(
            markdown_content, filename
        )
        knowledge_points = self.knowledge_extractor.parse_knowledge_markdown(
            knowledge_markdown
        )
        return {
            "knowledge_markdown": knowledge_markdown,
            "knowledge_points": knowledge_points
        }

    async def generate_script_from_knowledge_only(
        self,
        knowledge_markdown: str,
        filename: str = ""
    ) -> Dict[str, Any]:
        """
        仅从知识点生成智课脚本

        Args:
            knowledge_markdown: 知识点 Markdown 内容
            filename: 文件名

        Returns:
            Dict: 结构化脚本 JSON
        """
        return await self.knowledge_extractor.generate_script_from_knowledge(
            knowledge_markdown, filename
        )

    async def extract_and_generate_script(
        self,
        markdown_content: str,
        filename: str = ""
    ) -> Dict[str, Any]:
        """
        完整流程：提取知识点并生成智课脚本

        Args:
            markdown_content: 文档 Markdown 内容
            filename: 文件名

        Returns:
            Dict: 包含 knowledge_markdown, knowledge_points, script 的完整结果
        """
        logger.info(f"[DocumentService] 开始完整流程：提取知识点并生成脚本 - {filename}")

        # 步骤1：提取知识点
        knowledge_result = await self.extract_knowledge_only(markdown_content, filename)
        knowledge_markdown = knowledge_result["knowledge_markdown"]
        knowledge_points = knowledge_result["knowledge_points"]

        # 步骤2：从知识点生成脚本
        script = await self.knowledge_extractor.generate_script_from_knowledge(
            knowledge_markdown, filename
        )

        logger.info(f"[DocumentService] 完整流程完成：提取了 {len(knowledge_points)} 个知识点，生成了 {len(script.get('sections', []))} 个教学环节")

        return {
            "knowledge_markdown": knowledge_markdown,
            "knowledge_points": knowledge_points,
            "script": script
        }

    async def generate_script_only(
        self,
        markdown_content: str,
        filename: str
    ) -> ScriptResult:
        """
        仅生成智课脚本（基于原始文档内容）
        """
        return await self.script_generator.generate_script(markdown_content, filename)

    def process_rag_only(
        self,
        markdown_content: str,
        doc_name: str,
        doc_id: str
    ) -> RAGProcessResult:
        """
        仅执行RAG预处理
        """
        return self.rag_processor.process(markdown_content, doc_name, doc_id)


document_service = DocumentService()
