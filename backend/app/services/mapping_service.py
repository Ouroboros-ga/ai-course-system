"""
知识点↔PPT页面映射引擎服务
提供自动映射（基于已有page_start/page_end + AI语义匹配）和手动调整功能
"""

import json
import logging
from typing import List, Dict, Any, Optional

from sqlmodel import Session, select
from app.core.time_utils import utcnow_naive
from app.models.course_model import (
    Course,
    CourseScript,
    ScriptNode,
    DoclingText,
    DoclingDocument,
)
from app.models.mapping_model import KnowledgePageMap
from app.common.llm_client import llm_client, Message

logger = logging.getLogger(__name__)


class MappingService:
    """知识点↔PPT页面映射引擎"""

    @staticmethod
    def get_active_script(session: Session, course_id: int) -> Optional[CourseScript]:
        """获取课程的激活脚本"""
        return session.exec(
            select(CourseScript).where(
                CourseScript.course_id == course_id,
                CourseScript.is_active == True,
            )
        ).first()

    @staticmethod
    def get_page_texts(session: Session, course_id: int) -> List[Dict[str, Any]]:
        """获取PPT逐页文本内容（用于前端展示和AI匹配）"""
        docling_doc = session.exec(
            select(DoclingDocument).where(
                DoclingDocument.course_id == course_id,
            ).order_by(DoclingDocument.created_at.desc())
        ).first()

        if not docling_doc:
            return []

        # 从 raw_json 中提取逐页内容
        if docling_doc.raw_json and "texts" in docling_doc.raw_json:
            page_map: Dict[int, List[str]] = {}
            for text_data in docling_doc.raw_json["texts"]:
                page_no = text_data.get("page_no", 1)
                text = text_data.get("text", "")
                if text.strip():
                    page_map.setdefault(page_no, []).append(text.strip())

            # 也从 groups 中提取标题信息，补充到对应页码
            if "groups" in docling_doc.raw_json:
                for group_data in docling_doc.raw_json["groups"]:
                    name = group_data.get("name", "")
                    group_page = group_data.get("page_no", 0)
                    if name and len(name) >= 2 and group_page > 0:
                        page_map.setdefault(group_page, []).insert(0, f"[{name}]")

            pages = []
            for page_no in sorted(page_map.keys()):
                pages.append({
                    "page_no": page_no,
                    "text": "\n".join(page_map[page_no]),
                })

            # 检测是否所有文本都在第1页（旧数据问题），对PPTX文件尝试从源文件重新提取
            if len(pages) <= 1 and docling_doc.origin_filename:
                is_pptx = docling_doc.origin_filename.lower().endswith(('.pptx', '.ppt'))
                if is_pptx:
                    re_extracted = MappingService._re_extract_pptx_pages(docling_doc)
                    if re_extracted:
                        return re_extracted

            return pages

        # fallback: 从 docling_texts 表查询
        texts = session.exec(
            select(DoclingText).where(
                DoclingText.doc_id == docling_doc.id,
            ).order_by(DoclingText.page_no, DoclingText.sort_order)
        ).all()

        page_map: Dict[int, List[str]] = {}
        for t in texts:
            page_map.setdefault(t.page_no, []).append(t.text.strip())

        pages = []
        for page_no in sorted(page_map.keys()):
            pages.append({
                "page_no": page_no,
                "text": "\n".join(page_map[page_no]),
            })

        # 同样检测旧数据问题
        if len(pages) <= 1 and docling_doc.origin_filename:
            is_pptx = docling_doc.origin_filename.lower().endswith(('.pptx', '.ppt'))
            if is_pptx:
                re_extracted = MappingService._re_extract_pptx_pages(docling_doc)
                if re_extracted:
                    return re_extracted

        return pages

    @staticmethod
    def _re_extract_pptx_pages(docling_doc) -> List[Dict[str, Any]]:
        """
        当数据库中的page_no全部为1时，从原始PPTX文件重新提取逐页文本
        """
        file_path = docling_doc.source_file_path
        if not file_path:
            return []

        import os
        if not os.path.exists(file_path):
            return []

        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            pages = []
            for slide_num, slide in enumerate(prs.slides, 1):
                texts = []
                title = ""
                if slide.shapes.title and slide.shapes.title.text:
                    title = slide.shapes.title.text.strip()
                    texts.append(title)
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        if text != title:
                            texts.append(text)
                if texts:
                    pages.append({
                        "page_no": slide_num,
                        "text": "\n".join(texts),
                    })
            logger.info(f"[MappingService] 从PPTX重新提取 {len(pages)} 页内容")
            return pages
        except Exception as e:
            logger.warning(f"[MappingService] PPTX重新提取失败: {e}")
            return []

    @staticmethod
    def auto_map_from_nodes(session: Session, course_id: int) -> List[KnowledgePageMap]:
        """
        自动生成映射：基于 ScriptNode 已有的 page_start/page_end 字段
        当 page_start==page_end==1（未映射）时，尝试从 DoclingText/DoclingGroup 推断页码
        """
        script = MappingService.get_active_script(session, course_id)
        if not script:
            return []

        # 删除旧的自动映射（保留手动调整的）
        old_maps = session.exec(
            select(KnowledgePageMap).where(
                KnowledgePageMap.course_id == course_id,
                KnowledgePageMap.is_manual == False,
            )
        ).all()
        for old in old_maps:
            session.delete(old)
        session.commit()

        # 获取脚本节点
        nodes = session.exec(
            select(ScriptNode).where(
                ScriptNode.script_id == script.id,
            ).order_by(ScriptNode.node_index)
        ).all()

        # 获取页面文本，用于推断未映射节点的页码
        page_texts = MappingService.get_page_texts(session, course_id)
        total_pages = len(page_texts) if page_texts else 1

        # 构建标题→页码的映射表（从页面文本中提取）
        title_to_page = {}
        for page in page_texts:
            text = page.get("text", "")
            page_no = page.get("page_no", 1)
            # 取第一行作为标题线索
            first_line = text.split("\n")[0].strip() if text else ""
            if first_line and len(first_line) >= 2:
                title_to_page[first_line] = page_no

        mappings = []
        for idx, node in enumerate(nodes):
            # 跳过已有手动映射的节点
            existing_manual = session.exec(
                select(KnowledgePageMap).where(
                    KnowledgePageMap.node_id == node.id,
                    KnowledgePageMap.is_manual == True,
                )
            ).first()

            if existing_manual:
                mappings.append(existing_manual)
                continue

            page_start = node.page_start
            page_end = node.page_end

            # 当页码为默认值(1,1)时，尝试推断
            if page_start == 1 and page_end == 1 and total_pages > 1:
                # 策略1：根据节点标题匹配页面文本
                matched_page = None
                for title, pno in title_to_page.items():
                    if node.title and node.title in title:
                        matched_page = pno
                        break

                if matched_page:
                    page_start = matched_page
                    page_end = matched_page
                else:
                    # 策略2：按节点索引均匀分配页面
                    # 多个节点可以映射到同一页，确保 page_start <= page_end <= total_pages
                    nodes_count = len(nodes)
                    if nodes_count > 0 and total_pages > 0:
                        # 每页对应多少个节点（向上取整，确保最后一页也能覆盖）
                        nodes_per_page = max(1, (nodes_count + total_pages - 1) // total_pages)
                        page_start = min(idx // nodes_per_page + 1, total_pages)
                        page_end = page_start

            mapping = KnowledgePageMap(
                course_id=course_id,
                script_id=script.id,
                node_id=node.id,
                page_start=page_start,
                page_end=page_end,
                confidence=0.7 if (page_start != 1 or page_end != 1 or total_pages <= 1) else 0.3,
                is_manual=False,
            )
            session.add(mapping)
            mappings.append(mapping)

        session.commit()
        for m in mappings:
            session.refresh(m)

        return mappings

    @staticmethod
    async def ai_match_mapping(
        session: Session,
        course_id: int,
    ) -> List[KnowledgePageMap]:
        """
        AI语义匹配：使用LLM分析知识点标题/内容与PPT页面文本的语义相似度，
        重新计算映射关系
        """
        script = MappingService.get_active_script(session, course_id)
        if not script:
            return []

        nodes = session.exec(
            select(ScriptNode).where(
                ScriptNode.script_id == script.id,
            ).order_by(ScriptNode.node_index)
        ).all()

        page_texts = MappingService.get_page_texts(session, course_id)

        if not nodes or not page_texts:
            return []

        # 构建LLM匹配请求
        node_descriptions = []
        for node in nodes:
            node_descriptions.append({
                "node_id": node.id,
                "title": node.title or "",
                "content_preview": (node.content or "")[:200],
            })

        page_descriptions = []
        for page in page_texts:
            page_descriptions.append({
                "page_no": page["page_no"],
                "text_preview": page["text"][:300],
            })

        system_prompt = """你是一个教学内容分析专家。你需要将知识点与PPT页面的内容进行语义匹配。

给定一组知识点和PPT页面内容，请为每个知识点找到最匹配的PPT页面范围。

输出JSON格式：
{
  "mappings": [
    {
      "node_id": 知识点ID,
      "page_start": 起始页码,
      "page_end": 结束页码,
      "confidence": 置信度(0-1)
    }
  ]
}

匹配原则：
1. 一个知识点可以对应1-3页PPT
2. 页码范围应该连续
3. 不同知识点可以对应相同的PPT页面（一对多关系）
4. 置信度反映匹配的确定程度，0.5以下表示不确定
5. 只输出JSON，不要其他文字"""

        user_prompt = f"""请为以下知识点匹配PPT页面：

知识点列表：
{json.dumps(node_descriptions, ensure_ascii=False, indent=2)}

PPT页面内容：
{json.dumps(page_descriptions, ensure_ascii=False, indent=2)}"""

        try:
            messages = [
                Message(role="system", content=system_prompt),
                Message(role="user", content=user_prompt),
            ]
            response = await llm_client.chat(messages, temperature=0.1)

            # 解析LLM返回的JSON
            content = response.content.strip()
            # 提取JSON部分
            if "```json" in content:
                content = content.split("```json")[1].split("```")[0].strip()
            elif "```" in content:
                content = content.split("```")[1].split("```")[0].strip()

            result = json.loads(content)
            ai_mappings = result.get("mappings", [])

            # 删除旧的自动映射
            old_maps = session.exec(
                select(KnowledgePageMap).where(
                    KnowledgePageMap.course_id == course_id,
                    KnowledgePageMap.is_manual == False,
                )
            ).all()
            for old in old_maps:
                session.delete(old)
            session.commit()

            # 创建新的AI映射
            node_map = {n.id: n for n in nodes}
            mappings = []

            # 保留手动调整的映射
            manual_maps = session.exec(
                select(KnowledgePageMap).where(
                    KnowledgePageMap.course_id == course_id,
                    KnowledgePageMap.is_manual == True,
                )
            ).all()
            manual_node_ids = {m.node_id for m in manual_maps}
            mappings.extend(manual_maps)

            for ai_map in ai_mappings:
                node_id = ai_map.get("node_id")
                if node_id in manual_node_ids or node_id not in node_map:
                    continue

                mapping = KnowledgePageMap(
                    course_id=course_id,
                    script_id=script.id,
                    node_id=node_id,
                    page_start=ai_map.get("page_start", 1),
                    page_end=ai_map.get("page_end", 1),
                    confidence=ai_map.get("confidence", 0.5),
                    is_manual=False,
                )
                session.add(mapping)
                mappings.append(mapping)

            session.commit()
            for m in mappings:
                session.refresh(m)

            return mappings

        except Exception as e:
            logger.error(f"[MappingService] AI匹配失败: {e}")
            # fallback: 使用基于节点的自动映射
            return MappingService.auto_map_from_nodes(session, course_id)

    @staticmethod
    def update_node_mapping(
        session: Session,
        course_id: int,
        node_id: int,
        page_start: int,
        page_end: int,
    ) -> KnowledgePageMap:
        """手动调整单个知识点的映射关系"""
        mapping = session.exec(
            select(KnowledgePageMap).where(
                KnowledgePageMap.course_id == course_id,
                KnowledgePageMap.node_id == node_id,
            )
        ).first()

        if mapping:
            mapping.page_start = page_start
            mapping.page_end = page_end
            mapping.is_manual = True
            mapping.confidence = 1.0
            mapping.updated_at = utcnow_naive()
            session.add(mapping)
        else:
            script = MappingService.get_active_script(session, course_id)
            mapping = KnowledgePageMap(
                course_id=course_id,
                script_id=script.id if script else 0,
                node_id=node_id,
                page_start=page_start,
                page_end=page_end,
                confidence=1.0,
                is_manual=True,
            )
            session.add(mapping)

        session.commit()
        session.refresh(mapping)
        return mapping

    @staticmethod
    def batch_update_mappings(
        session: Session,
        course_id: int,
        updates: List[Dict[str, Any]],
    ) -> List[KnowledgePageMap]:
        """批量更新映射关系"""
        mappings = []
        for update in updates:
            node_id = update.get("node_id")
            page_start = update.get("page_start")
            page_end = update.get("page_end")
            if node_id and page_start and page_end:
                m = MappingService.update_node_mapping(
                    session, course_id, node_id, page_start, page_end
                )
                mappings.append(m)
        return mappings

    @staticmethod
    def apply_mappings_to_script(session: Session, course_id: int) -> bool:
        """
        将映射关系应用到脚本节点
        更新 ScriptNode 的 page_start/page_end，触发后续视频生成管线重新解析
        """
        mappings = session.exec(
            select(KnowledgePageMap).where(
                KnowledgePageMap.course_id == course_id,
            )
        ).all()

        if not mappings:
            return False

        mapping_by_node = {m.node_id: m for m in mappings}

        for node_id, mapping in mapping_by_node.items():
            node = session.get(ScriptNode, node_id)
            if node:
                node.page_start = mapping.page_start
                node.page_end = mapping.page_end
                session.add(node)

        session.commit()
        return True

    @staticmethod
    def calculate_timestamps_from_audio(
        session: Session,
        script_id: int,
        audio_duration: float,
    ) -> bool:
        """
        根据TTS合成的音频总时长，按比例重新计算各节点的精确时间戳
        
        在视频生成管线中，TTS合成完成后应调用此方法更新精确时间戳
        
        Args:
            session: 数据库会话
            script_id: 脚本ID
            audio_duration: 音频实际总时长(秒)
            
        Returns:
            bool: 是否成功更新
        """
        nodes = session.exec(
            select(ScriptNode)
            .where(ScriptNode.script_id == script_id)
            .order_by(ScriptNode.node_index)
        ).all()

        if not nodes:
            logger.warning(f"[Timestamp] 脚本 {script_id} 无节点，无法计算时间戳")
            return False

        total_duration = sum(node.duration for node in nodes)

        if total_duration == 0:
            logger.warning(f"[Timestamp] 脚本 {script_id} 节点总时长为0，使用均分策略")
            per_node_duration = audio_duration / len(nodes)
            current_time = 0.0
            for node in nodes:
                node.timestamp_start = round(current_time, 2)
                current_time += per_node_duration
                node.timestamp_end = round(current_time, 2)
                session.add(node)
        else:
            # 按duration占比计算实际时间戳（更精确）
            current_time = 0.0
            for node in nodes:
                actual_duration = (node.duration / total_duration) * audio_duration
                node.timestamp_start = round(current_time, 2)
                current_time += actual_duration
                node.timestamp_end = round(current_time, 2)
                session.add(node)

        session.commit()
        logger.info(
            f"[Timestamp] 已为脚本 {script_id} 的 {len(nodes)} 个节点更新精确时间戳"
            f"（基于音频时长 {audio_duration:.1f}秒）"
        )
        return True

    @staticmethod
    def get_mapping_detail(session: Session, course_id: int) -> Dict[str, Any]:
        """获取课程完整的映射详情（含节点信息和页面信息）"""
        script = MappingService.get_active_script(session, course_id)

        nodes = []
        if script:
            script_nodes = session.exec(
                select(ScriptNode).where(
                    ScriptNode.script_id == script.id,
                ).order_by(ScriptNode.node_index)
            ).all()

            mappings = session.exec(
                select(KnowledgePageMap).where(
                    KnowledgePageMap.course_id == course_id,
                )
            ).all()
            mapping_by_node = {m.node_id: m for m in mappings}

            for node in script_nodes:
                mapping = mapping_by_node.get(node.id)
                nodes.append({
                    "node_id": node.id,
                    "node_index": node.node_index,
                    "title": node.title,
                    "content_preview": (node.content or "")[:100],
                    "node_type": node.node_type.value if hasattr(node.node_type, "value") else node.node_type,
                    "is_key_point": node.is_key_point,
                    "page_start": mapping.page_start if mapping else node.page_start,
                    "page_end": mapping.page_end if mapping else node.page_end,
                    "confidence": mapping.confidence if mapping else 0.0,
                    "is_manual": mapping.is_manual if mapping else False,
                    "mapping_id": mapping.id if mapping else None,
                })

        page_texts = MappingService.get_page_texts(session, course_id)

        return {
            "course_id": course_id,
            "script_id": script.id if script else None,
            "nodes": nodes,
            "pages": page_texts,
            "total_nodes": len(nodes),
            "total_pages": len(page_texts),
        }


mapping_service = MappingService()
