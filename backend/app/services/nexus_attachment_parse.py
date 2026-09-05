"""NX-A1 附件解析层：八格式薄提取 + 引用定位。

设计取舍（见 v1.3 C2）：
- 与课程 ParserProvider 用**同一底层引擎**（fitz/python-docx/openpyxl/
  python-pptx/Pillow/LibreOffice），但不经过课程入库流水线
  （DocumentParseService/DocumentIR/Evidence/Graph），不伪造 course_id，
  产出只进 Nexus 会话上下文；
- 输出统一 block 模型，每个 block 自带 locator（页/slide/段落/表/单元格），
  模型引用时必须使用 locator，不得编造页码；
- 图片走"直传优先"：原图字节保留供未来视觉模型，OCR 只是可选文本补充，
  无 OCR/无视觉一律如实标注，不冒充看图；
- DOC/PPT（OLE）必须经 LibreOffice 转换；本机与服务器均无 soffice 时如实
  返回 CONVERT_UNAVAILABLE，不删目标格式（no-go 换方案不删目标）。

纯函数模块：无 DB、无网络（OCR 除外，由调用方注入结果），可完整单测。
"""

from __future__ import annotations

import io
import shutil
import subprocess
import tempfile
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

# 允许的附件格式（扩展名小写）。contentKind 决定解析器与引用维度。
ALLOWED_EXTENSIONS = ("pdf", "docx", "jpg", "jpeg", "png", "xlsx", "pptx", "ppt", "doc")
EXT_TO_KIND = {
    "pdf": "pdf",
    "docx": "docx",
    "jpg": "image",
    "jpeg": "image",
    "png": "image",
    "xlsx": "xlsx",
    "pptx": "pptx",
    "ppt": "ole",
    "doc": "ole",
}

# v1.3 C2 首版预算（解析层硬上限，超限走 partial/拒绝并如实标注）。
MAX_PDF_PAGES = 200
MAX_IMAGE_MEGAPIXELS = 20
MAX_WORKBOOK_CELLS = 100_000
MAX_BLOCK_TEXT_CHARS = 4000  # 单 block 截断上限（定位不断）
MAX_PARSED_BYTES = 256 * 1024  # 落盘 parsed.json 上限


class AttachmentParseError(Exception):
    """解析失败：携带机器可读 code（fail-closed，不伪造内容）。"""

    def __init__(self, code: str, detail: str = "") -> None:
        super().__init__(detail or code)
        self.code = code
        self.detail = detail


@dataclass
class ParsedAttachment:
    kind: str
    blocks: list[dict[str, Any]] = field(default_factory=list)
    truncated: bool = False  # 是否因预算截断（partial 依据之一）
    warnings: list[str] = field(default_factory=list)
    stats: dict[str, Any] = field(default_factory=dict)


def sniff_kind(data: bytes, filename: str) -> str:
    """魔数嗅探 + 扩展名交叉验证：不信任客户端声明的类型。"""
    ext = (filename.rsplit(".", 1)[-1] if "." in filename else "").lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise AttachmentParseError("ATTACHMENT_TYPE_UNSUPPORTED", f"不支持的格式：{ext or '（无扩展名）'}")
    kind = ""
    if data[:4] == b"%PDF":
        kind = "pdf"
    elif data[:8] == b"\x89PNG\r\n\x1a\n":
        kind = "image-png"
    elif data[:3] == b"\xff\xd8\xff":
        kind = "image-jpg"
    elif data[:2] == b"PK":
        kind = "zip"
    elif data[:8] == b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1":
        kind = "ole"
    else:
        raise AttachmentParseError("ATTACHMENT_UNRECOGNIZED", "无法识别的文件头")
    expected = EXT_TO_KIND[ext]
    if expected == "pdf" and kind != "pdf":
        raise AttachmentParseError("ATTACHMENT_TYPE_MISMATCH", "扩展名与文件头不一致（pdf）")
    if expected == "image" and kind not in ("image-png", "image-jpg"):
        raise AttachmentParseError("ATTACHMENT_TYPE_MISMATCH", "扩展名与文件头不一致（图片）")
    if expected in ("docx", "xlsx", "pptx") and kind != "zip":
        raise AttachmentParseError("ATTACHMENT_TYPE_MISMATCH", "扩展名与文件头不一致（OOXML）")
    if expected == "ole" and kind != "ole":
        raise AttachmentParseError("ATTACHMENT_TYPE_MISMATCH", "扩展名与文件头不一致（OLE）")
    if kind in ("image-png", "image-jpg"):
        return "image"
    if kind == "zip":
        # 轻量校验 OOXML 容器（不解压全部内容，防 ZIP 炸弹第一步）。
        try:
            with zipfile.ZipFile(io.BytesIO(data)) as zf:
                names = zf.namelist()
            if "[Content_Types].xml" not in names:
                raise AttachmentParseError("ATTACHMENT_MALFORMED", "OOXML 容器缺少 [Content_Types].xml")
        except zipfile.BadZipFile as error:
            raise AttachmentParseError("ATTACHMENT_MALFORMED", "ZIP 容器损坏") from error
        return expected
    return expected


def _clip(text: str) -> str:
    text = (text or "").strip()
    return text[:MAX_BLOCK_TEXT_CHARS]


def parse_pdf(data: bytes) -> ParsedAttachment:
    """PDF：逐页文本（fitz），无文本页标记 image-only（OCR 由服务层按需补）。"""
    try:
        import pymupdf
    except ImportError:
        import fitz as pymupdf  # type: ignore[no-redef]
    try:
        doc = pymupdf.open(stream=data, filetype="pdf")
    except Exception as error:
        raise AttachmentParseError("ATTACHMENT_MALFORMED", f"PDF 无法打开：{error}") from error
    out = ParsedAttachment(kind="pdf", stats={"pages": len(doc)})
    try:
        for i, page in enumerate(doc):
            if i >= MAX_PDF_PAGES:
                out.truncated = True
                out.warnings.append(f"仅解析前 {MAX_PDF_PAGES} 页（共 {len(doc)} 页）")
                break
            text = _clip(page.get_text("text"))
            if text:
                out.blocks.append({"kind": "text", "locator": f"p{i + 1}", "text": text})
            else:
                out.blocks.append({
                    "kind": "image-page", "locator": f"p{i + 1}",
                    "text": "", "needs_ocr": True,
                })
    finally:
        doc.close()
    return out


def parse_docx(data: bytes) -> ParsedAttachment:
    """DOCX：段落（段号）+ 表格（表号/行列），标题样式保留级别。"""
    import docx

    try:
        document = docx.Document(io.BytesIO(data))
    except Exception as error:
        raise AttachmentParseError("ATTACHMENT_MALFORMED", f"DOCX 无法打开：{error}") from error
    out = ParsedAttachment(kind="docx")
    for idx, para in enumerate(document.paragraphs, start=1):
        text = _clip(para.text)
        if not text:
            continue
        block: dict[str, Any] = {"kind": "text", "locator": f"para{idx}", "text": text}
        style = (para.style.name if para.style is not None else "") or ""
        if style.startswith("Heading"):
            block["heading"] = style
        out.blocks.append(block)
    for tidx, table in enumerate(document.tables, start=1):
        rows = [[_clip(cell.text) for cell in row.cells] for row in table.rows]
        rows = [r for r in rows if any(r)]
        if rows:
            out.blocks.append({
                "kind": "table", "locator": f"table{tidx}",
                "rows": len(rows), "cols": max(len(r) for r in rows),
                "text": "\n".join(" | ".join(r) for r in rows)[:MAX_BLOCK_TEXT_CHARS],
            })
    out.stats = {"paragraphs": len(document.paragraphs), "tables": len(document.tables)}
    return out


def parse_xlsx(data: bytes) -> ParsedAttachment:
    """XLSX：按 sheet/范围输出（openpyxl 只读），只取缓存值不执行公式。

    超过 10 万非空单元格即停（partial），先到先得按 sheet 顺序。
    """
    import openpyxl

    try:
        workbook = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    except Exception as error:
        raise AttachmentParseError("ATTACHMENT_MALFORMED", f"XLSX 无法打开：{error}") from error
    out = ParsedAttachment(kind="xlsx")
    cells = 0
    try:
        for sheet in workbook.worksheets:
            if cells >= MAX_WORKBOOK_CELLS:
                out.truncated = True
                out.warnings.append(f"超过 {MAX_WORKBOOK_CELLS} 非空单元格上限，后续 sheet 未解析")
                break
            lines: list[str] = []
            max_row = min(sheet.max_row or 0, 2000)
            max_col = min(sheet.max_column or 0, 100)
            for row in sheet.iter_rows(min_row=1, max_row=max_row, max_col=max_col, values_only=True):
                values = [(str(v) if v is not None else "") for v in row]
                while values and not values[-1]:
                    values.pop()
                if not any(values):
                    continue
                cells += sum(1 for v in values if v)
                if cells > MAX_WORKBOOK_CELLS:
                    out.truncated = True
                    out.warnings.append(f"超过 {MAX_WORKBOOK_CELLS} 非空单元格上限，已截断")
                    break
                lines.append(" | ".join(values))
            if lines:
                out.blocks.append({
                    "kind": "table", "locator": f"sheet:{sheet.title}",
                    "rows": len(lines), "text": "\n".join(lines)[:MAX_BLOCK_TEXT_CHARS],
                })
    finally:
        workbook.close()
    out.stats = {"sheets": len(workbook.sheetnames), "cells": cells}
    return out


def parse_pptx(data: bytes) -> ParsedAttachment:
    """PPTX：幻灯片号/文字/备注/表格（python-pptx 原生，不调课程管线）。"""
    import pptx

    try:
        presentation = pptx.Presentation(io.BytesIO(data))
    except Exception as error:
        raise AttachmentParseError("ATTACHMENT_MALFORMED", f"PPTX 无法打开：{error}") from error
    out = ParsedAttachment(kind="pptx", stats={"slides": len(presentation.slides)})
    for idx, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_table:
                rows = []
                for row in shape.table.rows:
                    values = [_clip(cell.text) for cell in row.cells]
                    if any(values):
                        rows.append(" | ".join(values))
                if rows:
                    texts.append("[表格]\n" + "\n".join(rows))
            elif shape.has_text_frame:
                text = _clip(shape.text)
                if text:
                    texts.append(text)
        notes = ""
        if slide.has_notes_slide:
            notes = _clip(slide.notes_slide.notes_text_frame.text)
        body = "\n".join(texts).strip()
        if body or notes:
            block: dict[str, Any] = {"kind": "text", "locator": f"slide{idx}"}
            if body:
                block["text"] = body[:MAX_BLOCK_TEXT_CHARS]
            if notes:
                block["notes"] = notes[:MAX_BLOCK_TEXT_CHARS]
            out.blocks.append(block)
    return out


def parse_image(data: bytes, ext: str) -> ParsedAttachment:
    """图片：校验解码 + 像素预算 + 尺寸元数据；原图直传（vision 预留），
    OCR 文本由服务层按需注入（本层不强制 OCR）。"""
    from PIL import Image

    try:
        image = Image.open(io.BytesIO(data))
        image.load()
    except Exception as error:
        raise AttachmentParseError("ATTACHMENT_MALFORMED", f"图片无法解码：{error}") from error
    width, height = image.size
    megapixels = width * height / 1_000_000
    if megapixels > MAX_IMAGE_MEGAPIXELS:
        raise AttachmentParseError(
            "ATTACHMENT_TOO_LARGE", f"图片 {megapixels:.1f}MP 超过 {MAX_IMAGE_MEGAPIXELS}MP 上限")
    out = ParsedAttachment(kind="image", stats={
        "width": width, "height": height, "format": image.format or ext,
        "vision": "unavailable",  # 无视觉模型配置，直传保留字节，消费侧如实降级
        "ocr": "pending",  # 服务层按 PaddleOCR 可用性回填 text/unavailable
    })
    out.blocks.append({
        "kind": "image", "locator": "img1",
        "text": f"[图片 {width}x{height}]",
        "width": width, "height": height,
    })
    return out


def convert_ole_to_pdf(data: bytes, source_ext: str, timeout_s: int = 120) -> bytes:
    """DOC/PPT 经 LibreOffice 转 PDF（--headless，禁宏/禁外链更新）。

    无 soffice 二进制即抛 CONVERT_UNAVAILABLE（两地均无，不删目标格式）。
    """
    binary = shutil.which("soffice") or shutil.which("libreoffice")
    if binary is None:
        raise AttachmentParseError(
            "ATTACHMENT_CONVERT_UNAVAILABLE",
            "DOC/PPT 需 LibreOffice 转换，当前环境未安装（已记录为待补能力）")
    with tempfile.TemporaryDirectory(prefix="nexus-ole-") as tmpdir:
        src = Path(tmpdir) / f"source.{source_ext}"
        src.write_bytes(data)
        try:
            subprocess.run(
                [binary, "--headless", "--nolockcheck", "--nodefault",
                 "--nofirststartwizard", "--norestore",
                 "--convert-to", "pdf", "--outdir", tmpdir, str(src)],
                capture_output=True, timeout=timeout_s, check=False,
            )
        except subprocess.TimeoutExpired as error:
            raise AttachmentParseError("ATTACHMENT_PARSE_TIMEOUT", "Office 转换超时") from error
        pdf_path = Path(tmpdir) / f"source.pdf"
        if not pdf_path.exists():
            raise AttachmentParseError("ATTACHMENT_MALFORMED", "Office 转换未产出 PDF")
        return pdf_path.read_bytes()


def parse_bytes(data: bytes, filename: str) -> ParsedAttachment:
    """入口：嗅探→分发→结构化产出。OLE 先转 PDF 再按 PDF 解析（页码为转换版页码）。"""
    kind = sniff_kind(data, filename)
    if kind == "pdf":
        return parse_pdf(data)
    if kind == "docx":
        return parse_docx(data)
    if kind == "xlsx":
        return parse_xlsx(data)
    if kind == "pptx":
        return parse_pptx(data)
    if kind == "image":
        ext = filename.rsplit(".", 1)[-1].lower()
        return parse_image(data, ext)
    if kind == "ole":
        ext = filename.rsplit(".", 1)[-1].lower()
        pdf_bytes = convert_ole_to_pdf(data, ext)
        parsed = parse_pdf(pdf_bytes)
        parsed.kind = "ole-pdf"
        parsed.warnings.insert(0, "DOC/PPT 经转换解析，页码为转换版页码，排版/备注可能损失")
        return parsed
    raise AttachmentParseError("ATTACHMENT_TYPE_UNSUPPORTED", f"未知 kind：{kind}")
