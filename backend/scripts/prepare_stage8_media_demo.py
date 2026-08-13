"""Prepare and publish a synthetic Stage 8 media demo course.

This helper exists only for local browser acceptance.  It never calls a paid
Provider: media audio must still be created through the registered Stage 8
batch endpoints while ``MEDIA_DEMO_MODE=true``.  The script only creates the
course/outline/script/PPT source needed by that flow and, after an active media
release exists, freezes a matching synthetic ``CourseRelease`` for learner
playback.

Safety rules:

* the requested course id must be absent or already carry this script's marker;
* existing non-demo courses are never modified;
* ``publish`` refuses incomplete or mismatched media releases;
* no API key, Provider speaker id, or real student data is read.
"""
from __future__ import annotations

import argparse
import hashlib
import json
import os
import sys
import tempfile
from pathlib import Path

os.environ.setdefault("AI_COURSE_SKIP_STARTUP_SIDE_EFFECTS", "1")
os.environ.setdefault("MEDIA_DEMO_MODE", "true")
sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402
from sqlmodel import Session, func, select  # noqa: E402

from app.core.time_utils import utcnow_aware  # noqa: E402
from app.models.access_control_model import CourseCapability  # noqa: E402
from app.models.course_build_model import (  # noqa: E402
    CourseRelease,
    ReleaseStatus,
    SourceMaterial,
    SourceMaterialVersion,
    MaterialStatus,
)
from app.models.course_model import Course, CourseStatus, StudentEnrollment  # noqa: E402
from app.models.course_outline_model import (  # noqa: E402
    CourseOutlineNode,
    CourseOutlineVersion,
    CoursePptMapping,
    OutlineLifecycleStatus,
    OutlineNodeType,
    TeachingScriptNode,
    TeachingScriptVersion,
)
from app.models.database import engine  # noqa: E402
from app.models.media_release_model import (  # noqa: E402
    MediaRelease,
    MediaReleaseItem,
    MediaReleaseStatus,
)
from app.models.user_model import User  # noqa: E402
from app.services.course_access_service import (  # noqa: E402
    activate_student_membership,
    establish_course_access_baseline,
)
from app.services.course_build_service import course_release_service  # noqa: E402
from app.services.object_storage import get_object_storage  # noqa: E402
from app.services.platform_media_preset_service import ensure_platform_presets  # noqa: E402
from app.services.unified_learning_service import ordered_outline_nodes  # noqa: E402


DEMO_MARKER_PREFIX = "local-stage8-media-demo-"
COURSE_TITLE = "课程 87 · 数字人媒体联调验收"
MATERIAL_ID_TEMPLATE = "sm_stage8_media_demo_{course_id}"
MATERIAL_VERSION_TEMPLATE = "smv_stage8_media_demo_{course_id}_v1"
OUTLINE_VERSION_TEMPLATE = "ov_stage8_media_demo_{course_id}_v1"
SCRIPT_VERSION_TEMPLATE = "tsv_stage8_media_demo_{course_id}_v1"

LESSON_ITEMS = (
    (
        "二分查找的前提与边界",
        "二分查找适用于已经有序的数据。每一步比较中间元素，并根据大小关系缩小搜索区间。"
        "实现时要明确左右边界是闭区间还是半开区间，整个循环必须保持同一个不变量。"
        "当目标不存在时，区间会稳定收缩为空，算法应返回明确的未找到结果。",
    ),
    (
        "循环不变量与中点计算",
        "循环不变量描述每次迭代开始时仍然成立的事实。对于闭区间写法，目标若存在就始终位于左右边界之间。"
        "中点可以写成左边界加上区间长度的一半，这种方式也能避免某些语言中的整数溢出。"
        "更新边界时必须排除已经比较过的中点，否则可能出现死循环。",
    ),
    (
        "复杂度与调试方法",
        "二分查找每次把候选区间缩小一半，因此时间复杂度是对数级，额外空间通常是常数级。"
        "调试时可以记录左边界、中点和右边界，检查区间是否严格缩小。"
        "还要覆盖空数组、单元素、首尾命中和目标不存在等边界用例。",
    ),
)


def _marker(course_id: int) -> str:
    return f"{DEMO_MARKER_PREFIX}{course_id}"


def _require_user(session: Session, username: str) -> User:
    user = session.exec(select(User).where(User.username == username)).first()
    if user is None:
        raise RuntimeError(f"local demo user {username!r} does not exist")
    return user


def _build_demo_pptx() -> bytes:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    palette = (
        (RGBColor(22, 47, 77), RGBColor(228, 239, 249)),
        (RGBColor(34, 82, 76), RGBColor(227, 242, 238)),
        (RGBColor(91, 62, 39), RGBColor(247, 238, 226)),
    )
    for index, ((title, body), (ink, paper)) in enumerate(zip(LESSON_ITEMS, palette), start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = paper

        heading = slide.shapes.add_textbox(Inches(0.9), Inches(0.75), Inches(11.5), Inches(1.0))
        heading_frame = heading.text_frame
        heading_frame.clear()
        heading_run = heading_frame.paragraphs[0].add_run()
        heading_run.text = f"{index:02d}  {title}"
        heading_run.font.name = "Microsoft YaHei"
        heading_run.font.size = Pt(30)
        heading_run.font.bold = True
        heading_run.font.color.rgb = ink

        content = slide.shapes.add_textbox(Inches(1.05), Inches(2.0), Inches(11.0), Inches(3.7))
        content_frame = content.text_frame
        content_frame.word_wrap = True
        content_frame.clear()
        for paragraph_index, sentence in enumerate(body.split("。")):
            sentence = sentence.strip()
            if not sentence:
                continue
            paragraph = content_frame.paragraphs[0] if paragraph_index == 0 else content_frame.add_paragraph()
            paragraph.text = f"• {sentence}。"
            paragraph.font.name = "Microsoft YaHei"
            paragraph.font.size = Pt(21)
            paragraph.font.color.rgb = ink
            paragraph.space_after = Pt(14)
            paragraph.alignment = PP_ALIGN.LEFT

        footer = slide.shapes.add_textbox(Inches(0.95), Inches(6.75), Inches(11.4), Inches(0.35))
        footer_paragraph = footer.text_frame.paragraphs[0]
        footer_paragraph.text = "Stage 8 本地合成媒体验收 · 非正式课程内容"
        footer_paragraph.font.name = "Microsoft YaHei"
        footer_paragraph.font.size = Pt(10)
        footer_paragraph.font.color.rgb = ink

    descriptor, filename = tempfile.mkstemp(prefix="stage8-media-demo-", suffix=".pptx")
    os.close(descriptor)
    path = Path(filename).resolve()
    try:
        presentation.save(path)
        return path.read_bytes()
    finally:
        path.unlink(missing_ok=True)


def _ensure_course(session: Session, *, course_id: int, teacher: User) -> Course:
    marker = _marker(course_id)
    course = session.get(Course, course_id)
    if course is not None and course.fanya_course_id != marker:
        raise RuntimeError(
            f"course {course_id} already exists and is not the synthetic Stage 8 demo; refusing to modify it"
        )
    if course is None:
        course = Course(
            id=course_id,
            fanya_course_id=marker,
            fanya_course_name=COURSE_TITLE,
            title=COURSE_TITLE,
            description="仅用于本地 Fake WAV、PPT、字幕与 PixiJS 数字人联调。",
            teacher_id=int(teacher.id),
            status=CourseStatus.DRAFT,
            is_ai_generated=False,
            total_nodes=len(LESSON_ITEMS),
            total_pages=len(LESSON_ITEMS),
            source_file_name="stage8-media-demo.pptx",
            source_mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        session.add(course)
        session.flush()
    establish_course_access_baseline(session, course_id, int(teacher.id))
    capabilities = session.exec(
        select(CourseCapability).where(CourseCapability.course_id == course_id)
    ).first()
    if capabilities is None:
        capabilities = CourseCapability(course_id=course_id)
    capabilities.learning = True
    capabilities.course_building = True
    capabilities.knowledge_graph = True
    capabilities.evidence = True
    capabilities.cognitive_analysis = True
    session.add(capabilities)
    return course


def _ensure_material(session: Session, *, course_id: int, teacher: User, course: Course) -> SourceMaterialVersion:
    material_id = MATERIAL_ID_TEMPLATE.format(course_id=course_id)
    version_id = MATERIAL_VERSION_TEMPLATE.format(course_id=course_id)
    material = session.exec(select(SourceMaterial).where(SourceMaterial.material_id == material_id)).first()
    version = session.exec(select(SourceMaterialVersion).where(SourceMaterialVersion.version_id == version_id)).first()
    if material is not None and material.course_id != course_id:
        raise RuntimeError("synthetic material id belongs to a different course")
    if version is not None and version.course_id != course_id:
        raise RuntimeError("synthetic material version belongs to a different course")

    pptx_bytes = _build_demo_pptx()
    content_hash = hashlib.sha256(pptx_bytes).hexdigest()
    object_key = f"course-materials/course{course_id}/stage8-demo-{content_hash[:16]}.pptx"
    storage = get_object_storage()
    if not storage.exists(object_key):
        storage.put(
            object_key,
            pptx_bytes,
            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    if material is None:
        material = SourceMaterial(
            material_id=material_id,
            course_id=course_id,
            name="Stage 8 本地媒体验收课件",
            material_type="slide",
            material_role="primary_courseware",
            include_in_course_corpus=False,
            source_kind="local_demo_fixture",
            status=MaterialStatus.PARSED,
            created_by=int(teacher.id),
        )
        session.add(material)
        session.flush()
    if version is None:
        version = SourceMaterialVersion(
            version_id=version_id,
            material_id=material.material_id,
            course_id=course_id,
            version=1,
            file_path=object_key,
            file_hash=content_hash,
            file_size=len(pptx_bytes),
            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            parse_status=MaterialStatus.PARSED,
            is_current=True,
            created_by=int(teacher.id),
        )
        session.add(version)
        session.flush()
    material.current_version_id = version.version_id
    material.status = MaterialStatus.PARSED
    session.add(material)
    course.source_file_path = object_key
    course.source_file_name = "stage8-media-demo.pptx"
    session.add(course)
    return version


def _ensure_outline_and_scripts(
    session: Session,
    *,
    course_id: int,
    teacher: User,
    material_version: SourceMaterialVersion,
) -> tuple[CourseOutlineVersion, TeachingScriptVersion]:
    outline_id = OUTLINE_VERSION_TEMPLATE.format(course_id=course_id)
    script_id = SCRIPT_VERSION_TEMPLATE.format(course_id=course_id)
    outline = session.exec(
        select(CourseOutlineVersion).where(CourseOutlineVersion.outline_version_id == outline_id)
    ).first()
    if outline is None:
        outline = CourseOutlineVersion(
            outline_version_id=outline_id,
            course_id=course_id,
            version=1,
            lifecycle_status=OutlineLifecycleStatus.DRAFT,
            generation_source="local_demo_fixture",
            review_status="teacher_edited",
            created_by=int(teacher.id),
        )
        session.add(outline)
        session.flush()
    if outline.course_id != course_id:
        raise RuntimeError("synthetic outline id belongs to a different course")

    chapter_id = f"on_stage8_media_demo_{course_id}_chapter"
    chapter = session.exec(
        select(CourseOutlineNode).where(CourseOutlineNode.outline_node_id == chapter_id)
    ).first()
    if chapter is None:
        chapter = CourseOutlineNode(
            outline_node_id=chapter_id,
            outline_version_id=outline.outline_version_id,
            course_id=course_id,
            node_type=OutlineNodeType.CHAPTER,
            title="算法基础",
            order_index=0,
            content_hash=hashlib.sha256("算法基础".encode("utf-8")).hexdigest(),
        )
        session.add(chapter)
        session.flush()

    script = session.exec(
        select(TeachingScriptVersion).where(TeachingScriptVersion.script_version_id == script_id)
    ).first()
    if script is None:
        script = TeachingScriptVersion(
            script_version_id=script_id,
            course_id=course_id,
            outline_version_id=outline.outline_version_id,
            version=1,
            lifecycle_status=OutlineLifecycleStatus.DRAFT,
            generation_source="local_demo_fixture",
            review_status="teacher_edited",
            created_by=int(teacher.id),
        )
        session.add(script)
        session.flush()

    for index, (title, body) in enumerate(LESSON_ITEMS, start=1):
        outline_node_id = f"on_stage8_media_demo_{course_id}_kp{index}"
        concept_key = f"kg_stage8_media_demo_{course_id}_kp{index}"
        outline_node = session.exec(
            select(CourseOutlineNode).where(CourseOutlineNode.outline_node_id == outline_node_id)
        ).first()
        if outline_node is None:
            outline_node = CourseOutlineNode(
                outline_node_id=outline_node_id,
                outline_version_id=outline.outline_version_id,
                course_id=course_id,
                parent_node_id=chapter.outline_node_id,
                node_type=OutlineNodeType.KNOWLEDGE_POINT,
                title=title,
                order_index=index - 1,
                page_range=str(index),
                # The local fixture deliberately uses deterministic synthetic
                # graph keys so a release-pinned review can be exercised.
                knowledge_graph_node_id=concept_key,
                content_hash=hashlib.sha256(title.encode("utf-8")).hexdigest(),
            )
            session.add(outline_node)
            session.flush()
        elif not outline_node.knowledge_graph_node_id:
            # Older copies of this local-only fixture predate learning
            # adjustments.  Backfill only this script-owned synthetic node.
            outline_node.knowledge_graph_node_id = concept_key
            session.add(outline_node)

        script_node_id = f"tsn_stage8_media_demo_{course_id}_kp{index}"
        script_node = session.exec(
            select(TeachingScriptNode).where(TeachingScriptNode.script_node_id == script_node_id)
        ).first()
        if script_node is None:
            session.add(TeachingScriptNode(
                script_node_id=script_node_id,
                script_version_id=script.script_version_id,
                course_id=course_id,
                outline_node_id=outline_node.outline_node_id,
                content=body,
                style="beginner",
                evidence_refs=[],
                source_block_refs=[],
                content_hash=hashlib.sha256(body.encode("utf-8")).hexdigest(),
            ))

        mapping = session.exec(select(CoursePptMapping).where(
            CoursePptMapping.course_id == course_id,
            CoursePptMapping.outline_node_id == outline_node.outline_node_id,
            CoursePptMapping.material_version_id == material_version.version_id,
        )).first()
        if mapping is None:
            session.add(CoursePptMapping(
                course_id=course_id,
                outline_node_id=outline_node.outline_node_id,
                material_version_id=material_version.version_id,
                page_start=index,
                page_end=index,
                page_refs=[index],
                confidence=1.0,
                status="draft",
                teacher_locked=True,
                created_by=int(teacher.id),
                updated_by=int(teacher.id),
            ))
    session.flush()
    return outline, script


def prepare(course_id: int, teacher_username: str, student_username: str) -> dict[str, object]:
    with Session(engine) as session:
        teacher = _require_user(session, teacher_username)
        student = _require_user(session, student_username)
        course = _ensure_course(session, course_id=course_id, teacher=teacher)
        material_version = _ensure_material(
            session,
            course_id=course_id,
            teacher=teacher,
            course=course,
        )
        outline, script = _ensure_outline_and_scripts(
            session,
            course_id=course_id,
            teacher=teacher,
            material_version=material_version,
        )
        activate_student_membership(session, course_id, int(student.id))
        enrollment = session.exec(select(StudentEnrollment).where(
            StudentEnrollment.student_id == int(student.id),
            StudentEnrollment.course_id == course_id,
        )).first()
        if enrollment is None:
            session.add(StudentEnrollment(
                student_id=int(student.id),
                course_id=course_id,
                total_nodes_count=len(LESSON_ITEMS),
                is_active=True,
            ))
        ensure_platform_presets(session)
        session.commit()
        return {
            "mode": "prepared",
            "course_id": course_id,
            "course_status": course.status.value,
            "teacher": teacher.username,
            "student": student.username,
            "outline_version_id": outline.outline_version_id,
            "script_version_id": script.script_version_id,
            "material_version_id": material_version.version_id,
            "knowledge_point_count": len(LESSON_ITEMS),
            "paid_provider_called": False,
        }


def publish(course_id: int) -> dict[str, object]:
    with Session(engine) as session:
        course = session.get(Course, course_id)
        if course is None or course.fanya_course_id != _marker(course_id):
            raise RuntimeError("publish is limited to this script's synthetic demo course")
        media = session.exec(select(MediaRelease).where(
            MediaRelease.course_id == course_id,
            MediaRelease.status == MediaReleaseStatus.ACTIVE,
        ).order_by(MediaRelease.version_number.desc())).first()
        if media is None:
            raise RuntimeError("course has no active MediaRelease; finish the real media build flow first")
        if not media.audio_playlist_object_key or not media.audio_playlist_sha256:
            raise RuntimeError("active MediaRelease has no frozen audio-playlist/v1")
        storage = get_object_storage()
        if not storage.exists(media.audio_playlist_object_key):
            raise RuntimeError("active MediaRelease playlist object is unavailable")

        items = list(session.exec(select(MediaReleaseItem).where(
            MediaReleaseItem.course_id == course_id,
            MediaReleaseItem.release_id == media.release_id,
        )).all())
        if not items or any(item.status != "ready" for item in items):
            raise RuntimeError("active MediaRelease does not contain an all-ready item set")
        script_node_ids = [item.node_id for item in items]
        script_nodes = list(session.exec(select(TeachingScriptNode).where(
            TeachingScriptNode.id.in_(script_node_ids),
        )).all())
        if len(script_nodes) != len(items):
            raise RuntimeError("media item script nodes are incomplete")
        script_version_ids = {node.script_version_id for node in script_nodes}
        if len(script_version_ids) != 1:
            raise RuntimeError("media items span multiple script versions")
        script = session.exec(select(TeachingScriptVersion).where(
            TeachingScriptVersion.script_version_id == next(iter(script_version_ids)),
            TeachingScriptVersion.course_id == course_id,
        )).first()
        if script is None:
            raise RuntimeError("media script version is unavailable")
        outline = session.exec(select(CourseOutlineVersion).where(
            CourseOutlineVersion.outline_version_id == script.outline_version_id,
            CourseOutlineVersion.course_id == course_id,
        )).first()
        if outline is None:
            raise RuntimeError("media outline version is unavailable")

        expected_ids = [node.outline_node_id for node in ordered_outline_nodes(
            session,
            outline_version_id=outline.outline_version_id,
            knowledge_points_only=True,
        )]
        actual_ids = [item.outline_node_id for item in sorted(items, key=lambda row: row.order_index)]
        if expected_ids != actual_ids:
            raise RuntimeError(
                f"media playlist does not match outline pre-order: expected={expected_ids}, actual={actual_ids}"
            )

        existing = session.exec(select(CourseRelease).where(
            CourseRelease.course_id == course_id,
            CourseRelease.is_active == True,  # noqa: E712
            CourseRelease.status == ReleaseStatus.PUBLISHED,
        )).first()
        if existing is not None and (
            (existing.media_snapshot or {}).get("media_release_id") == media.release_id
            and (existing.media_snapshot or {}).get("playlist_content_hash") == media.audio_playlist_sha256
        ):
            return {
                "mode": "already_published",
                "course_id": course_id,
                "course_release_id": existing.release_id,
                "media_release_id": media.release_id,
                "playlist_content_hash": media.audio_playlist_sha256,
                "paid_provider_called": False,
            }

        if existing is not None:
            existing.is_active = False
            existing.status = ReleaseStatus.SUPERSEDED
            session.add(existing)

        outline.lifecycle_status = OutlineLifecycleStatus.PUBLISHED
        script.lifecycle_status = OutlineLifecycleStatus.PUBLISHED
        session.add(outline)
        session.add(script)
        mappings = list(session.exec(select(CoursePptMapping).where(
            CoursePptMapping.course_id == course_id,
            CoursePptMapping.outline_node_id.in_(expected_ids),
        )).all())
        for mapping in mappings:
            mapping.status = "published"
            session.add(mapping)

        latest_version = session.exec(select(func.max(CourseRelease.version)).where(
            CourseRelease.course_id == course_id,
        )).one() or 0
        page_mapping_items = [
            {
                "outline_node_id": mapping.outline_node_id,
                "material_version_id": mapping.material_version_id,
                "page_start": mapping.page_start,
                "page_end": mapping.page_end,
                "page_refs": mapping.page_refs or [],
            }
            for mapping in mappings
        ]
        release = CourseRelease(
            course_id=course_id,
            version=int(latest_version) + 1,
            prev_release_id=existing.release_id if existing else None,
            status=ReleaseStatus.PUBLISHED,
            is_active=True,
            structure_snapshot={
                "outline_version_id": outline.outline_version_id,
                "knowledge_point_ids": expected_ids,
            },
            scripts_snapshot={
                "script_version_id": script.script_version_id,
                "script_node_ids": [node.script_node_id for node in script_nodes],
            },
            page_mappings_snapshot={"items": page_mapping_items},
            media_snapshot={
                "media_release_id": media.release_id,
                "playlist_content_hash": media.audio_playlist_sha256,
                "audio_playlist_object_key": media.audio_playlist_object_key,
                "ppt_manifest_object_key": media.ppt_manifest_object_key,
                "avatar_preset_id": media.avatar_preset_id,
                "avatar_preset_version": media.avatar_preset_version,
            },
            outline_version_id=outline.outline_version_id,
            script_version_id=script.script_version_id,
            quality_gate_passed=True,
            publication_check_snapshot={
                "fixture": "stage8-local-browser-acceptance/v1",
                "note": "Synthetic local fixture; not a production quality-gate result.",
            },
            label="Stage 8 本地媒体播放验收",
            release_notes="仅用于本地 Fake WAV、PPT、字幕与 PixiJS 数字人回归。",
            published_by=course.teacher_id,
            published_at=utcnow_aware(),
            created_by=course.teacher_id,
        )
        release.content_hash = course_release_service._release_content_hash(release)
        session.add(release)
        session.flush()
        course_release_service._record_frozen_artifacts(session, release=release)
        course.status = CourseStatus.PUBLISHED
        course.updated_at = utcnow_aware()
        session.add(course)
        session.commit()
        return {
            "mode": "published",
            "course_id": course_id,
            "course_release_id": release.release_id,
            "media_release_id": media.release_id,
            "playlist_content_hash": media.audio_playlist_sha256,
            "outline_version_id": outline.outline_version_id,
            "script_version_id": script.script_version_id,
            "knowledge_point_count": len(expected_ids),
            "paid_provider_called": False,
        }


def status(course_id: int) -> dict[str, object]:
    with Session(engine) as session:
        course = session.get(Course, course_id)
        if course is None:
            return {"course_id": course_id, "exists": False}
        media = session.exec(select(MediaRelease).where(
            MediaRelease.course_id == course_id,
            MediaRelease.status == MediaReleaseStatus.ACTIVE,
        ).order_by(MediaRelease.version_number.desc())).first()
        release = course_release_service.get_active_release(session, course_id=course_id)
        return {
            "course_id": course_id,
            "exists": True,
            "synthetic_demo": course.fanya_course_id == _marker(course_id),
            "course_status": course.status.value,
            "active_media_release_id": media.release_id if media else None,
            "playlist_content_hash": media.audio_playlist_sha256 if media else None,
            "active_course_release_id": release.release_id if release else None,
            "paid_provider_called": False,
        }


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("command", choices=("prepare", "publish", "status"))
    parser.add_argument("--course-id", type=int, default=87)
    parser.add_argument("--teacher", default="TTT")
    parser.add_argument("--student", default="SSS")
    args = parser.parse_args()

    if args.command == "prepare":
        result = prepare(args.course_id, args.teacher, args.student)
    elif args.command == "publish":
        result = publish(args.course_id)
    else:
        result = status(args.course_id)
    print(json.dumps(result, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
