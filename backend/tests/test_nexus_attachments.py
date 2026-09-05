"""NX-A1/NX-E1：附件八格式解析 + 生命周期 + run 恢复查询。

存储隔离：对象存储换成 tmp LocalStorageProvider；元数据表为可移植 DDL，
SQLite/PG 双跑（恢复语义必须本地可测）。
合成样例：pdf/docx/xlsx/pptx/png/jpg 六格式真实生成并解析断言；
doc/ppt 本地与服务器均无 LibreOffice，如实 failed（CONVERT_UNAVAILABLE），
目标保留（no-go 换方案不删目标）。
"""

import io
import json
import uuid

import pytest

from app.api.v1.endpoints import nexus_internal, nexus_proxy
from app.core.security import create_access_token, get_password_hash
from app.models.access_control_model import PlatformPermission, PlatformPermissionAssignment
from app.models.user_model import User, UserRole


@pytest.fixture
def isolated_storage(tmp_path, monkeypatch):
    from app.services import object_storage as os_mod
    from app.services.object_storage import LocalStorageProvider

    provider = LocalStorageProvider(str(tmp_path / "obj"))
    monkeypatch.setattr(os_mod, "_object_storage", provider)
    return provider


@pytest.fixture
def internal_configured(monkeypatch):
    monkeypatch.setattr(nexus_internal.settings, "NEXUS_INTERNAL_TOKEN", "internal-token-1")


@pytest.fixture
def worker_configured(monkeypatch):
    monkeypatch.setattr(nexus_proxy.settings, "REPRO_WORKER_URL", "http://127.0.0.1:8400")
    monkeypatch.setattr(nexus_proxy.settings, "REPRO_WORKER_TOKEN", "worker-token-1")


@pytest.fixture
def nexus_student_token(session, student_user):
    session.add(PlatformPermissionAssignment(
        user_id=student_user.id,
        permission=PlatformPermission.NEXUS_USE,
    ))
    session.commit()
    return create_access_token({
        "sub": str(student_user.id),
        "username": student_user.username,
        "role": student_user.role.value,
        "school_id": student_user.school_id or "test-school",
    })


def _auth(token: str) -> dict[str, str]:
    return {"Authorization": f"Bearer {token}"}


def _internal_headers(user_id) -> dict[str, str]:
    return {"Authorization": "Bearer internal-token-1", "X-Nexus-User-Id": str(user_id)}


# ---------------------------------------------------------------------------
# 合成样例（六格式真实生成）
# ---------------------------------------------------------------------------


def _pdf_bytes() -> bytes:
    import pymupdf

    doc = pymupdf.open()
    p1 = doc.new_page()
    p1.insert_text((72, 72), "Transformer attention QKV")
    p2 = doc.new_page()
    p2.insert_text((72, 72), "second page FFN")
    raw = doc.tobytes()
    doc.close()
    return raw


def _docx_bytes() -> bytes:
    import docx

    document = docx.Document()
    document.add_heading("综述标题", level=1)
    document.add_paragraph("注意力机制正文")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "模型"
    table.cell(0, 1).text = "分数"
    table.cell(1, 0).text = "A"
    table.cell(1, 1).text = "0.9"
    buf = io.BytesIO()
    document.save(buf)
    return buf.getvalue()


def _xlsx_bytes() -> bytes:
    import openpyxl

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "指标"
    sheet["A1"] = "模型"
    sheet["B1"] = "val_loss"
    sheet["A2"] = "nanoGPT"
    sheet["B2"] = 1.8857
    buf = io.BytesIO()
    workbook.save(buf)
    return buf.getvalue()


def _pptx_bytes() -> bytes:
    import pptx

    presentation = pptx.Presentation()
    layout = presentation.slide_layouts[6]  # blank
    slide = presentation.slides.add_slide(layout)
    box = slide.shapes.add_textbox(0, 0, 1000000, 1000000)
    box.text_frame.text = "幻灯片正文注意力"
    slide.notes_slide.notes_text_frame.text = "备注讲解稿"
    buf = io.BytesIO()
    presentation.save(buf)
    return buf.getvalue()


def _png_bytes() -> bytes:
    from PIL import Image

    buf = io.BytesIO()
    Image.new("RGB", (64, 32), color="red").save(buf, format="PNG")
    return buf.getvalue()


def _jpg_bytes() -> bytes:
    from PIL import Image

    buf = io.BytesIO()
    Image.new("RGB", (32, 32), color="blue").save(buf, format="JPEG")
    return buf.getvalue()


def _upload(client, token, filename, data, session_id=""):
    return client.post(
        "/api/v1/nexus/attachments",
        files={"file": (filename, data, "application/octet-stream")},
        data={"session_id": session_id},
        headers=_auth(token),
    )


# ---------------------------------------------------------------------------
# NX-A1：六格式真实解析
# ---------------------------------------------------------------------------


def test_upload_pdf_ready_with_page_locators(client, nexus_student_token, isolated_storage):
    response = _upload(client, nexus_student_token, "paper.pdf", _pdf_bytes())
    assert response.status_code == 200
    body = response.json()
    assert body["status"] == "ready"
    assert body["stats"]["pages"] == 2
    detail = client.get(
        f"/api/v1/nexus/attachments/{body['attachment_id']}?include_blocks=true",
        headers=_auth(nexus_student_token),
    ).json()
    locators = [b["locator"] for b in detail["content"]["blocks"]]
    assert locators == ["p1", "p2"]
    assert "QKV" in detail["content"]["blocks"][0]["text"]


def test_upload_docx_para_and_table(client, nexus_student_token, isolated_storage):
    body = _upload(client, nexus_student_token, "note.docx", _docx_bytes()).json()
    assert body["status"] == "ready"
    detail = client.get(
        f"/api/v1/nexus/attachments/{body['attachment_id']}?include_blocks=true",
        headers=_auth(nexus_student_token),
    ).json()
    by_locator = {b["locator"]: b for b in detail["content"]["blocks"]}
    assert "注意力机制正文" in by_locator["para2"]["text"]
    assert by_locator["para1"]["heading"] == "Heading 1"
    assert by_locator["table1"]["rows"] == 2


def test_upload_xlsx_sheet_range(client, nexus_student_token, isolated_storage):
    body = _upload(client, nexus_student_token, "m.xlsx", _xlsx_bytes()).json()
    assert body["status"] == "ready"
    detail = client.get(
        f"/api/v1/nexus/attachments/{body['attachment_id']}?include_blocks=true",
        headers=_auth(nexus_student_token),
    ).json()
    block = detail["content"]["blocks"][0]
    assert block["locator"] == "sheet:指标"
    assert "1.8857" in block["text"]


def test_upload_pptx_slide_and_notes(client, nexus_student_token, isolated_storage):
    body = _upload(client, nexus_student_token, "d.pptx", _pptx_bytes()).json()
    assert body["status"] == "ready"
    detail = client.get(
        f"/api/v1/nexus/attachments/{body['attachment_id']}?include_blocks=true",
        headers=_auth(nexus_student_token),
    ).json()
    block = detail["content"]["blocks"][0]
    assert block["locator"] == "slide1"
    assert "幻灯片正文注意力" in block["text"]
    assert "备注讲解稿" in block["notes"]


def test_upload_images_direct_honest_degrade(client, nexus_student_token, isolated_storage):
    for name, data in (("a.png", _png_bytes()), ("b.jpg", _jpg_bytes())):
        body = _upload(client, nexus_student_token, name, data).json()
        # 直传成功（ready），视觉明确 unavailable，OCR 本地无服务 unavailable。
        assert body["status"] == "ready", body
        assert body["stats"]["vision"] == "unavailable"
        assert body["stats"]["ocr"] == "unavailable"


def test_upload_doc_without_libreoffice_honest_failed(
    client, nexus_student_token, isolated_storage
):
    fake_ole = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 512
    body = _upload(client, nexus_student_token, "old.doc", fake_ole).json()
    assert body["status"] == "failed"
    assert body["error_code"] == "ATTACHMENT_CONVERT_UNAVAILABLE"


def test_upload_rejects_mismatch_and_unsupported(client, nexus_student_token, isolated_storage):
    # 扩展名 pdf 但内容是 PNG。
    response = _upload(client, nexus_student_token, "fake.pdf", _png_bytes())
    assert response.status_code == 422
    # 不支持的格式。
    response = _upload(client, nexus_student_token, "a.zip", b"PK\x03\x04" + b"\x00" * 64)
    assert response.status_code == 422
    # 空文件。
    response = _upload(client, nexus_student_token, "e.pdf", b"")
    assert response.status_code == 422


# ---------------------------------------------------------------------------
# NX-A1：生命周期（绑定/隔离/删除/过期）
# ---------------------------------------------------------------------------


def test_bind_session_and_mismatch(client, nexus_student_token, isolated_storage):
    aid = _upload(client, nexus_student_token, "p.pdf", _pdf_bytes()).json()["attachment_id"]
    ok = client.post(
        f"/api/v1/nexus/attachments/{aid}/bind",
        json={"session_id": "s1"},
        headers=_auth(nexus_student_token),
    )
    assert ok.status_code == 200
    assert ok.json()["session_id"] == "s1"
    # 同会话幂等。
    again = client.post(
        f"/api/v1/nexus/attachments/{aid}/bind",
        json={"session_id": "s1"},
        headers=_auth(nexus_student_token),
    )
    assert again.status_code == 200
    # 他会话 403。
    other = client.post(
        f"/api/v1/nexus/attachments/{aid}/bind",
        json={"session_id": "s2"},
        headers=_auth(nexus_student_token),
    )
    assert other.status_code == 403


def _make_other_token(session) -> str:
    """另一个有 Nexus 使用权但数据隔离的用户（跨用户是数据不可见，不是无权限）。"""
    other = User(
        username=f"nexus_other_{uuid.uuid4().hex[:8]}",
        real_name="Nexus Other",
        hashed_password=get_password_hash("test-password"),
        role=UserRole.STUDENT,
        is_active=True,
    )
    session.add(other)
    session.commit()
    session.refresh(other)
    session.add(PlatformPermissionAssignment(
        user_id=other.id,
        permission=PlatformPermission.NEXUS_USE,
    ))
    session.commit()
    return create_access_token({
        "sub": str(other.id), "username": other.username,
        "role": other.role.value, "school_id": "test-school",
    })


def test_owner_isolation(client, session, nexus_student_token, isolated_storage):
    other_token = _make_other_token(session)
    aid = _upload(client, nexus_student_token, "p.pdf", _pdf_bytes()).json()["attachment_id"]
    assert client.get(f"/api/v1/nexus/attachments/{aid}",
                      headers=_auth(other_token)).status_code == 404
    assert client.get(f"/api/v1/nexus/attachments/{aid}/download",
                      headers=_auth(other_token)).status_code == 404
    assert client.delete(f"/api/v1/nexus/attachments/{aid}",
                         headers=_auth(other_token)).status_code == 404


def test_delete_revokes_immediately(client, nexus_student_token, isolated_storage):
    aid = _upload(client, nexus_student_token, "p.pdf", _pdf_bytes()).json()["attachment_id"]
    assert client.delete(f"/api/v1/nexus/attachments/{aid}",
                         headers=_auth(nexus_student_token)).status_code == 200
    assert client.get(f"/api/v1/nexus/attachments/{aid}",
                      headers=_auth(nexus_student_token)).status_code == 404
    assert client.get(f"/api/v1/nexus/attachments/{aid}/download",
                      headers=_auth(nexus_student_token)).status_code == 404


def test_lazy_expiry_marks_expired(client, nexus_student_token, isolated_storage, session):
    from sqlalchemy import text as sql_text

    aid = _upload(client, nexus_student_token, "p.pdf", _pdf_bytes()).json()["attachment_id"]
    session.connection().execute(
        sql_text("UPDATE nexus_attachments SET expires_at=1.0 WHERE attachment_id=:aid"),
        {"aid": aid},
    )
    session.commit()
    body = client.get(f"/api/v1/nexus/attachments/{aid}",
                      headers=_auth(nexus_student_token)).json()
    assert body["status"] == "expired"


# ---------------------------------------------------------------------------
# NX-A1：内部 content 端点（Runtime 工具消费）
# ---------------------------------------------------------------------------


def test_internal_content_requires_session_binding(
    client, nexus_student_token, student_user, internal_configured, isolated_storage
):
    aid = _upload(client, nexus_student_token, "p.pdf", _pdf_bytes()).json()["attachment_id"]
    base = f"/api/v1/nexus-internal/attachments/{aid}/content"
    headers = _internal_headers(student_user.id)
    # 未绑定会话 → 403。
    assert client.get(base, headers=headers).status_code == 403
    # 错会话 → 403。
    wrong = dict(headers, **{"X-Nexus-Session-Id": "other"})
    assert client.get(base, headers=wrong).status_code == 403
    # 绑定后正确会话 → blocks + locator 精读。
    client.post(f"/api/v1/nexus/attachments/{aid}/bind", json={"session_id": "s1"},
                headers=_auth(nexus_student_token))
    right = dict(headers, **{"X-Nexus-Session-Id": "s1"})
    ok = client.get(base, headers=right)
    assert ok.status_code == 200
    assert [b["locator"] for b in ok.json()["data"]["blocks"]] == ["p1", "p2"]
    one = client.get(base, headers=right, params={"locator": "p2"})
    assert [b["locator"] for b in one.json()["data"]["blocks"]] == ["p2"]
    missing = client.get(base, headers=right, params={"locator": "p9"})
    assert missing.status_code == 422


# ---------------------------------------------------------------------------
# NX-E1：run 注册与恢复查询
# ---------------------------------------------------------------------------


def _record_run(client, student_user, run_id="run-1", session_id="s1", job_id="job-1"):
    return client.post(
        "/api/v1/nexus-internal/repro-runs",
        json={"run_id": run_id, "session_id": session_id, "tool": "run_reproduction",
              "preset_id": "nanogpt", "plan_hash": "ph", "approval_id": "apv-1",
              "job_id": job_id, "status": "submitted"},
        headers=_internal_headers(student_user.id),
    )


def test_run_record_and_session_recovery(
    client, session, nexus_student_token, student_user, internal_configured,
    worker_configured, isolated_storage, monkeypatch,
):
    import httpx

    assert _record_run(client, student_user).status_code == 200
    # 跨用户：他人列表为空、详情 404（不区分不存在与非归属）。
    other_token = _make_other_token(session)
    assert client.get("/api/v1/nexus/runs?session_id=s1",
                      headers=_auth(other_token)).json()["items"] == []
    assert client.get("/api/v1/nexus/runs/run-1",
                      headers=_auth(other_token)).status_code == 404

    # run_id 冲突属他人：拒绝覆盖（409），原行归属不变。
    conflict = client.post(
        "/api/v1/nexus-internal/repro-runs",
        json={"run_id": "run-1", "session_id": "sx", "tool": "run_reproduction",
              "preset_id": "nanogpt", "plan_hash": "ph", "approval_id": "apv-x",
              "job_id": "job-x", "status": "submitted"},
        headers={"Authorization": "Bearer internal-token-1",
                 "X-Nexus-User-Id": "999999"},
    )
    assert conflict.status_code == 409

    async def _live(request: httpx.Request) -> httpx.Response:
        return httpx.Response(200, json={
            "job_id": "job-1", "status": "running", "preset_id": "nanogpt",
            "steps_result": [], "artifacts": [],
        })

    real_client = httpx.AsyncClient

    def _factory(**kwargs):
        kwargs["transport"] = httpx.MockTransport(_live)
        return real_client(**kwargs)

    monkeypatch.setattr(nexus_proxy.httpx, "AsyncClient", _factory)
    items = client.get("/api/v1/nexus/runs?session_id=s1",
                       headers=_auth(nexus_student_token)).json()["items"]
    assert len(items) == 1
    assert items[0]["run_id"] == "run-1"
    assert items[0]["live"]["status"] == "running"
    detail = client.get("/api/v1/nexus/runs/run-1",
                        headers=_auth(nexus_student_token)).json()
    assert detail["job_id"] == "job-1"


def test_run_recovery_worker_missing_and_down(
    client, nexus_student_token, student_user, internal_configured, worker_configured,
    isolated_storage, monkeypatch,
):
    import httpx

    assert _record_run(client, student_user, run_id="run-9", job_id="job-9").status_code == 200

    async def _gone(request: httpx.Request) -> httpx.Response:
        return httpx.Response(404, json={"detail": "no such job"})

    real_client = httpx.AsyncClient

    def _factory(**kwargs):
        kwargs["transport"] = httpx.MockTransport(_gone)
        return real_client(**kwargs)

    monkeypatch.setattr(nexus_proxy.httpx, "AsyncClient", _factory)
    items = client.get("/api/v1/nexus/runs?session_id=s1",
                       headers=_auth(nexus_student_token)).json()["items"]
    # Worker 重启丢内存 → unknown + 不可恢复说明，不伪造终态。
    assert items[0]["live"]["status"] == "unknown"
    assert "不可恢复" in items[0]["live"]["note"]

    def _boom(request: httpx.Request) -> httpx.Response:
        raise httpx.ConnectError("down")

    def _factory2(**kwargs):
        kwargs["transport"] = httpx.MockTransport(_boom)
        return real_client(**kwargs)

    monkeypatch.setattr(nexus_proxy.httpx, "AsyncClient", _factory2)
    items = client.get("/api/v1/nexus/runs?session_id=s1",
                       headers=_auth(nexus_student_token)).json()["items"]
    # Worker 不可达 → 回落登记快照并标 stale。
    assert items[0]["live"]["status"] == "stale"
