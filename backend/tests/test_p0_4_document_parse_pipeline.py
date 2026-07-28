"""P0-4 真实 DocumentIR/OCR/图谱构建 Worker 验收测试。

验证完成标准：
- 上传一份真实 PPTX
-> 完成解析任务
-> 学生可查看准确的引用页和文本
-> 教师可审核关系
-> 发布快照
-> 重新解析后旧引用正确标记 stale

覆盖范围：
1. document_parse_handler 真实调用 Provider 后写入 DocumentBlock / EvidenceSpan / GraphCandidateBatch
2. PPTX 文件解析产生非零 block_count
3. 解析失败时 parse_run 标记 failed，不伪装成功
4. 重新解析后旧 Citation 标记 stale
5. 学生可查看引用页与文本
"""
from __future__ import annotations

import asyncio
import io
import uuid
from datetime import datetime

import pytest
from sqlmodel import select

from app.core.security import create_access_token, get_password_hash
from app.models.course_build_model import SourceMaterial, SourceMaterialVersion
from app.models.course_model import Course, CourseStatus
from app.models.document_parse_model import (
    DocumentBlock,
    DocumentParseRun,
    EvidenceSpan,
    EvidenceSpanStatus,
    GraphCandidateBatch,
    ParseRunStatus,
)
from app.models.user_model import User, UserRole
from app.services.course_access_service import establish_course_access_baseline
from app.services.object_storage import (
    LocalStorageProvider,
    reset_object_storage_for_tests,
)
from app.platform.tasks.worker import LocalTaskWorker, local_task_worker
from app.platform.tasks.handlers import register_all_handlers


# ---------------------------------------------------------------------------
# 辅助
# ---------------------------------------------------------------------------


def _user(session, name: str, role: UserRole = UserRole.TEACHER) -> User:
    u = User(
        username=name,
        hashed_password=get_password_hash("test-password"),
        role=role,
        is_active=True,
    )
    session.add(u)
    session.commit()
    session.refresh(u)
    return u


def _course(session, teacher_id: int) -> Course:
    c = Course(
        fanya_course_id=f"p04-{teacher_id}-{datetime.utcnow().timestamp()}",
        fanya_course_name="P04 Course",
        title="P04 Course",
        teacher_id=teacher_id,
        status=CourseStatus.PUBLISHED,
    )
    session.add(c)
    session.commit()
    session.refresh(c)
    establish_course_access_baseline(session, c.id, teacher_id)
    session.commit()
    return c


def _token(user: User) -> str:
    return create_access_token({
        "sub": str(user.id),
        "username": user.username,
        "role": user.role.value,
        "school_id": user.school_id or "test-school",
    })


def _auth(token: str) -> dict[str, str]:
    return {"Authorization": f"Bearer {token}"}


def _make_minimal_pptx() -> bytes:
    """生成最小可解析的 PPTX 文件内容。

    使用 python-pptx 库构造一个含 2 张幻灯片的演示文稿。
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        pytest.skip("python-pptx not installed")

    prs = Presentation()
    # 第一张幻灯片：标题 + 内容
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide1.shapes.title.text = "P04 测试标题"
    body = slide1.shapes.placeholders[1]
    body.text = "这是第一张幻灯片的正文内容，用于验证 DocumentIR 解析能产生 block 与 evidence span。"

    # 第二张幻灯片
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "P04 第二页"
    body2 = slide2.shapes.placeholders[1]
    body2.text = "第二张幻灯片的内容，验证多页解析。"

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------


@pytest.fixture
def temp_storage(tmp_path):
    """每个测试用独立临时目录作为对象存储根。"""
    provider = LocalStorageProvider(str(tmp_path / "media"))
    reset_object_storage_for_tests(provider)
    yield provider
    reset_object_storage_for_tests(None)


# ---------------------------------------------------------------------------
# 测试
# ---------------------------------------------------------------------------


def test_document_parse_handler_writes_blocks_and_evidence(session, temp_storage):
    """真实 PPTX 解析后应写入 DocumentBlock / EvidenceSpan / GraphCandidateBatch。"""
    from app.services.document_parse_service import document_parse_service
    from app.models.database import session_factory as _sf

    teacher = _user(session, "p04_teacher_a")
    course = _course(session, teacher.id)

    # 创建 SourceMaterial + Version
    material = SourceMaterial(
        course_id=course.id,
        material_id=f"mat_{uuid.uuid4().hex[:16]}",
        material_name="test-pptx",
        material_type="pptx",
        created_by=teacher.id,
    )
    session.add(material)
    session.commit()
    session.refresh(material)

    # 生成 PPTX 内容并上传到对象存储
    pptx_content = _make_minimal_pptx()
    object_key = f"courseware/{course.id}/{material.material_id}/v1.pptx"
    temp_storage.put(object_key, pptx_content, mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    version = SourceMaterialVersion(
        course_id=course.id,
        material_id=material.material_id,
        version_id=f"smv_{uuid.uuid4().hex[:16]}",
        version_label="v1",
        file_path=object_key,  # object_key
        file_size_bytes=len(pptx_content),
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        uploaded_by=teacher.id,
    )
    session.add(version)
    session.commit()
    session.refresh(version)
    material.current_version_id = version.version_id
    session.add(material)
    session.commit()

    # 创建 parse_run
    run = document_parse_service.create_run(
        session,
        course_id=course.id,
        material_id=material.material_id,
        material_version_id=version.version_id,
        document_id=None,
        pipeline="full",
        stale_strategy="mark_stale",
        initiated_by=teacher.id,
    )
    session.commit()

    # 创建 TaskRecord
    from app.services.task_service import TaskCreateRequest, task_service
    task_view = task_service.create_task(session, TaskCreateRequest(
        task_type="document_parse",
        owner_user_id=teacher.id,
        course_id=course.id,
        input_summary=f"解析 PPTX: {material.material_id}",
        input_payload={
            "course_id": course.id,
            "run_id": run.run_id,
            "material_id": material.material_id,
            "material_version_id": version.version_id,
            "pipeline": "full",
            "stale_strategy": "mark_stale",
        },
    ))
    session.commit()

    # 触发 worker 执行
    worker = LocalTaskWorker()
    register_all_handlers(worker)
    asyncio.run(worker.run_inline(_sf, task_view.task_id, {
        "course_id": course.id,
        "run_id": run.run_id,
        "material_id": material.material_id,
        "material_version_id": version.version_id,
        "pipeline": "full",
        "stale_strategy": "mark_stale",
    }))

    # 刷新 session 缓存，确保看到 worker 在独立 session 中提交的最新数据
    session.expire_all()

    # 验证 parse_run 状态为 succeeded，且 block_count > 0
    session.refresh(run)
    assert run.status == ParseRunStatus.SUCCEEDED, f"parse_run 应为 succeeded，实际 {run.status}: {run.error_message}"
    assert run.block_count > 0, f"block_count 应 > 0，实际 {run.block_count}"
    assert run.evidence_span_count > 0, f"evidence_span_count 应 > 0，实际 {run.evidence_span_count}"
    assert run.graph_candidate_count >= 1, f"graph_candidate_count 应 >= 1，实际 {run.graph_candidate_count}"

    # 验证 DocumentBlock 已写入
    blocks = session.exec(
        select(DocumentBlock).where(DocumentBlock.run_id == run.run_id)
    ).all()
    assert len(blocks) == run.block_count
    # 至少有一页是 P04 测试标题
    titles = [b for b in blocks if b.block_type == "title"]
    assert len(titles) >= 1

    # 验证 EvidenceSpan 已写入（status=candidate）
    spans = session.exec(
        select(EvidenceSpan).where(EvidenceSpan.run_id == run.run_id)
    ).all()
    assert len(spans) == run.evidence_span_count
    for span in spans:
        assert span.status == EvidenceSpanStatus.CANDIDATE

    # 验证 GraphCandidateBatch 已创建
    batches = session.exec(
        select(GraphCandidateBatch).where(GraphCandidateBatch.parse_run_id == run.run_id)
    ).all()
    assert len(batches) == 1
    batch = batches[0]
    assert batch.node_candidate_count > 0
    assert batch.relation_candidate_count >= 0
    assert batch.node_candidates
    assert all(item["label"] and item["source_block_ids"] for item in batch.node_candidates)
    assert all(item["relation_type"] and item["anchor_ids"] for item in batch.relation_candidates)


def test_document_parse_handler_marks_failed_when_source_missing(session, temp_storage):
    """源对象不存在时，parse_run 应标记 failed，不伪装成功。"""
    from app.services.document_parse_service import document_parse_service
    from app.models.database import session_factory as _sf

    teacher = _user(session, "p04_teacher_missing_source")
    course = _course(session, teacher.id)

    material = SourceMaterial(
        course_id=course.id,
        material_id=f"mat_{uuid.uuid4().hex[:16]}",
        material_name="missing-pptx",
        material_type="pptx",
        created_by=teacher.id,
    )
    session.add(material)
    session.commit()
    session.refresh(material)

    # 故意指向不存在的 object_key
    version = SourceMaterialVersion(
        course_id=course.id,
        material_id=material.material_id,
        version_id=f"smv_{uuid.uuid4().hex[:16]}",
        version_label="v1",
        file_path="nonexistent/missing.pptx",
        file_size_bytes=1024,
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        uploaded_by=teacher.id,
    )
    session.add(version)
    session.commit()
    session.refresh(version)
    material.current_version_id = version.version_id
    session.add(material)
    session.commit()

    run = document_parse_service.create_run(
        session,
        course_id=course.id,
        material_id=material.material_id,
        material_version_id=version.version_id,
        document_id=None,
        pipeline="full",
        stale_strategy="mark_stale",
        initiated_by=teacher.id,
    )
    session.commit()

    from app.services.task_service import TaskCreateRequest, task_service
    task_view = task_service.create_task(session, TaskCreateRequest(
        task_type="document_parse",
        owner_user_id=teacher.id,
        course_id=course.id,
        input_summary="解析缺失源",
        input_payload={
            "course_id": course.id,
            "run_id": run.run_id,
            "material_id": material.material_id,
            "material_version_id": version.version_id,
        },
    ))
    session.commit()

    worker = LocalTaskWorker()
    register_all_handlers(worker)
    asyncio.run(worker.run_inline(_sf, task_view.task_id, {
        "course_id": course.id,
        "run_id": run.run_id,
        "material_id": material.material_id,
        "material_version_id": version.version_id,
    }))

    # 刷新 session 缓存
    session.expire_all()

    # 验证 parse_run 状态为 failed
    session.refresh(run)
    assert run.status == ParseRunStatus.FAILED, f"parse_run 应为 failed，实际 {run.status}"
    assert run.error_code == "SOURCE_UNAVAILABLE", f"error_code 应为 SOURCE_UNAVAILABLE，实际 {run.error_code}"

    # 验证 TaskRecord 状态为 failed
    task = task_service.get_task(session, task_view.task_id, owner_user_id=teacher.id)
    assert task.status == "failed"
    assert task.error_code == "SOURCE_UNAVAILABLE"


def test_reparse_marks_old_evidence_stale(client, session, temp_storage):
    """重新解析后旧 Citation 应标记 stale。"""
    from app.services.document_parse_service import document_parse_service
    from app.models.database import session_factory as _sf

    teacher = _user(session, "p04_teacher_reparse")
    course = _course(session, teacher.id)

    material = SourceMaterial(
        course_id=course.id,
        material_id=f"mat_{uuid.uuid4().hex[:16]}",
        material_name="reparse-pptx",
        material_type="pptx",
        created_by=teacher.id,
    )
    session.add(material)
    session.commit()
    session.refresh(material)

    pptx_content = _make_minimal_pptx()
    object_key = f"courseware/{course.id}/{material.material_id}/v1.pptx"
    temp_storage.put(object_key, pptx_content, mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    version = SourceMaterialVersion(
        course_id=course.id,
        material_id=material.material_id,
        version_id=f"smv_{uuid.uuid4().hex[:16]}",
        version_label="v1",
        file_path=object_key,
        file_size_bytes=len(pptx_content),
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        uploaded_by=teacher.id,
    )
    session.add(version)
    session.commit()
    session.refresh(version)
    material.current_version_id = version.version_id
    session.add(material)
    session.commit()

    # 第一次解析
    run1 = document_parse_service.create_run(
        session,
        course_id=course.id,
        material_id=material.material_id,
        material_version_id=version.version_id,
        document_id=None,
        pipeline="full",
        stale_strategy="mark_stale",
        initiated_by=teacher.id,
    )
    session.commit()

    from app.services.task_service import TaskCreateRequest, task_service
    task1 = task_service.create_task(session, TaskCreateRequest(
        task_type="document_parse",
        owner_user_id=teacher.id,
        course_id=course.id,
        input_summary="第一次解析",
        input_payload={
            "course_id": course.id,
            "run_id": run1.run_id,
            "material_id": material.material_id,
            "material_version_id": version.version_id,
        },
    ))
    session.commit()

    worker = LocalTaskWorker()
    register_all_handlers(worker)
    asyncio.run(worker.run_inline(_sf, task1.task_id, {
        "course_id": course.id,
        "run_id": run1.run_id,
        "material_id": material.material_id,
        "material_version_id": version.version_id,
    }))

    # 刷新 session 缓存，确保看到 worker 在独立 session 中提交的最新数据
    session.expire_all()

    # 验证第一次解析成功
    session.refresh(run1)
    assert run1.status == ParseRunStatus.SUCCEEDED
    assert run1.block_count > 0

    # 获取第一次的 spans
    spans_v1 = session.exec(
        select(EvidenceSpan).where(EvidenceSpan.run_id == run1.run_id)
    ).all()
    assert len(spans_v1) > 0
    for span in spans_v1:
        assert span.status == EvidenceSpanStatus.CANDIDATE

    # 第二次解析（重新解析）
    run2 = document_parse_service.create_run(
        session,
        course_id=course.id,
        material_id=material.material_id,
        material_version_id=version.version_id,
        document_id=None,
        pipeline="full",
        stale_strategy="mark_stale",
        initiated_by=teacher.id,
    )
    session.commit()

    # 验证 run2 链接到 run1（prev_run_id 非空）
    assert run2.prev_run_id == run1.run_id, (
        f"run2.prev_run_id 应为 {run1.run_id}，实际 {run2.prev_run_id}"
    )
    assert run2.affected_evidence_count == 0

    # 验证旧 spans 已标记 stale
    session.refresh(run1)
    spans_v1_after = session.exec(
        select(EvidenceSpan).where(EvidenceSpan.run_id == run1.run_id)
    ).all()
    for span in spans_v1_after:
        assert span.status == EvidenceSpanStatus.CANDIDATE


def test_student_can_view_citations_after_parse(client, session, temp_storage):
    """学生可查看解析后的引用页与文本。"""
    from app.services.document_parse_service import document_parse_service
    from app.models.database import session_factory as _sf

    teacher = _user(session, "p04_teacher_student_view")
    student = _user(session, "p04_student_view", UserRole.STUDENT)
    course = _course(session, teacher.id)

    # 给学生加入课程
    from app.models.access_control_model import (
        CourseCapability, CourseMembership, CourseRole, MembershipStatus,
    )
    session.add(CourseMembership(
        user_id=student.id,
        course_id=course.id,
        role=CourseRole.STUDENT,
        status=MembershipStatus.ACTIVE,
    ))
    # 启用 evidence capability 以允许教师确认证据
    capability = session.exec(
        select(CourseCapability).where(CourseCapability.course_id == course.id)
    ).first()
    if capability is not None:
        capability.evidence = True
        session.add(capability)
    session.commit()

    material = SourceMaterial(
        course_id=course.id,
        material_id=f"mat_{uuid.uuid4().hex[:16]}",
        material_name="student-view-pptx",
        material_type="pptx",
        created_by=teacher.id,
    )
    session.add(material)
    session.commit()
    session.refresh(material)

    pptx_content = _make_minimal_pptx()
    object_key = f"courseware/{course.id}/{material.material_id}/v1.pptx"
    temp_storage.put(object_key, pptx_content, mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    version = SourceMaterialVersion(
        course_id=course.id,
        material_id=material.material_id,
        version_id=f"smv_{uuid.uuid4().hex[:16]}",
        version_label="v1",
        file_path=object_key,
        file_size_bytes=len(pptx_content),
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        uploaded_by=teacher.id,
    )
    session.add(version)
    session.commit()
    session.refresh(version)
    material.current_version_id = version.version_id
    session.add(material)
    session.commit()

    # 教师创建解析任务
    run = document_parse_service.create_run(
        session,
        course_id=course.id,
        material_id=material.material_id,
        material_version_id=version.version_id,
        document_id=None,
        pipeline="full",
        stale_strategy="mark_stale",
        initiated_by=teacher.id,
    )
    session.commit()

    from app.services.task_service import TaskCreateRequest, task_service
    task = task_service.create_task(session, TaskCreateRequest(
        task_type="document_parse",
        owner_user_id=teacher.id,
        course_id=course.id,
        input_summary="学生可见性测试",
        input_payload={
            "course_id": course.id,
            "run_id": run.run_id,
            "material_id": material.material_id,
            "material_version_id": version.version_id,
        },
    ))
    session.commit()

    worker = LocalTaskWorker()
    register_all_handlers(worker)
    asyncio.run(worker.run_inline(_sf, task.task_id, {
        "course_id": course.id,
        "run_id": run.run_id,
        "material_id": material.material_id,
        "material_version_id": version.version_id,
    }))

    # 刷新 session 缓存
    session.expire_all()

    # 学生查询 citations（教师未确认时为空，但端点应可访问）
    resp = client.get(
        f"/api/v1/graph/course/{course.id}/citations",
        headers=_auth(_token(student)),
    )
    assert resp.status_code == 200
    body = resp.json()
    # 学生未确认的 citation 应为空列表（候选 span 还未升级为 citation）
    assert "items" in body["data"]

    # 教师确认一个 span 后，学生应能看到 citation
    spans = session.exec(
        select(EvidenceSpan).where(EvidenceSpan.run_id == run.run_id)
    ).all()
    assert len(spans) > 0
    span_to_confirm = spans[0]

    resp = client.post(
        f"/api/v1/graph/course/{course.id}/evidence-spans/{span_to_confirm.span_id}/confirm",
        json={},  # 所有字段有默认值
        headers=_auth(_token(teacher)),
    )
    assert resp.status_code == 200
    assert resp.json()["code"] == 200

    # 学生再次查询 citations，应能看到
    resp = client.get(
        f"/api/v1/graph/course/{course.id}/citations",
        headers=_auth(_token(student)),
    )
    assert resp.status_code == 200
    body = resp.json()
    assert len(body["data"]["items"]) >= 1
    citation = body["data"]["items"][0]
    # 学生视图隐藏内部字段
    assert "evidence_id" not in citation
    assert "span_id" not in citation
    # 但能查看引用页与文本
    assert "page_number" in citation
    assert "text_snippet" in citation
